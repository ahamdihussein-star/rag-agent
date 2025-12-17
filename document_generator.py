from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
from langchain_openai import ChatOpenAI
from langchain_openai import OpenAIEmbeddings
from pinecone import Pinecone
from dotenv import load_dotenv
from datetime import datetime
import os
import json

# Load environment variables
load_dotenv()

# Initialize
pc = Pinecone(api_key=os.getenv("PINECONE_API_KEY"))
index = pc.Index(os.getenv("PINECONE_INDEX_NAME"))
embeddings = OpenAIEmbeddings(model="text-embedding-3-large")
llm = ChatOpenAI(model="gpt-4o", temperature=0.7)

def get_context(query: str, top_k: int = 5):
    """Get relevant context from Pinecone"""
    query_vector = embeddings.embed_query(query)
    results = index.query(
        vector=query_vector,
        top_k=top_k,
        include_metadata=True
    )
    
    context = ""
    for match in results['matches']:
        context += match['metadata']['text'] + "\n\n"
    
    return context

def generate_content(prompt: str, context: str):
    """Use LLM to generate content based on context"""
    
    full_prompt = f"""Based on the following context, {prompt}

Context:
{context}

Respond in JSON format as requested."""
    
    response = llm.invoke(full_prompt)
    return response.content

# ==================== PowerPoint ====================

def create_ppt_from_template(template_path: str, output_path: str, topic: str):
    """Create PowerPoint from template"""
    
    print(f"📊 Creating PowerPoint about: {topic}")
    
    # Get context from knowledge base
    context = get_context(topic)
    
    # Generate slide content
    prompt = f"""create a presentation about "{topic}".
    
Return JSON with this structure:
{{
    "title": "Presentation Title",
    "subtitle": "Subtitle",
    "slides": [
        {{
            "title": "Slide Title",
            "bullet_points": ["Point 1", "Point 2", "Point 3"]
        }}
    ]
}}

Generate 5-7 slides. JSON only:"""
    
    content = generate_content(prompt, context)
    
    # Parse JSON
    try:
        json_str = content.strip()
        if json_str.startswith("```"):
            json_str = json_str.split("```")[1]
            if json_str.startswith("json"):
                json_str = json_str[4:]
        json_str = json_str.strip()
        data = json.loads(json_str)
    except:
        print("❌ Failed to parse content")
        return None
    
    # Load template or create new
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    
    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = data.get("title", topic)
    subtitle.text = data.get("subtitle", "")
    
    # Content slides
    for slide_data in data.get("slides", []):
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = slide_data.get("title", "")
        
        body = slide.placeholders[1]
        tf = body.text_frame
        
        for i, point in enumerate(slide_data.get("bullet_points", [])):
            if i == 0:
                tf.text = point
            else:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
    
    # Save
    prs.save(output_path)
    print(f"✅ PowerPoint saved: {output_path}")
    return output_path

def create_ppt_simple(output_path: str, topic: str):
    """Create PowerPoint without template"""
    return create_ppt_from_template(None, output_path, topic)

# ==================== Word Document ====================

def create_word_from_template(template_path: str, output_path: str, topic: str, doc_type: str = "proposal"):
    """Create Word document from template"""
    
    print(f"📝 Creating {doc_type} about: {topic}")
    
    # Get context from knowledge base
    context = get_context(topic)
    
    # Generate document content
    prompt = f"""create a {doc_type} about "{topic}".
    
Return JSON with this structure:
{{
    "title": "Document Title",
    "sections": [
        {{
            "heading": "Section Heading",
            "content": "Section content paragraph..."
        }}
    ]
}}

Generate 4-6 sections. JSON only:"""
    
    content = generate_content(prompt, context)
    
    # Parse JSON
    try:
        json_str = content.strip()
        if json_str.startswith("```"):
            json_str = json_str.split("```")[1]
            if json_str.startswith("json"):
                json_str = json_str[4:]
        json_str = json_str.strip()
        data = json.loads(json_str)
    except:
        print("❌ Failed to parse content")
        return None
    
    # Load template or create new
    if template_path and os.path.exists(template_path):
        doc = Document(template_path)
    else:
        doc = Document()
    
    # Title
    doc.add_heading(data.get("title", topic), 0)
    
    # Date
    doc.add_paragraph(f"Date: {datetime.now().strftime('%B %d, %Y')}")
    doc.add_paragraph("")
    
    # Sections
    for section in data.get("sections", []):
        doc.add_heading(section.get("heading", ""), level=1)
        doc.add_paragraph(section.get("content", ""))
    
    # Save
    doc.save(output_path)
    print(f"✅ Word document saved: {output_path}")
    return output_path

def create_word_simple(output_path: str, topic: str, doc_type: str = "proposal"):
    """Create Word document without template"""
    return create_word_from_template(None, output_path, topic, doc_type)

# ==================== Excel ====================

def create_excel_from_template(template_path: str, output_path: str, topic: str, sheet_type: str = "bom"):
    """Create Excel from template"""
    
    print(f"📈 Creating {sheet_type} about: {topic}")
    
    # Get context from knowledge base
    context = get_context(topic)
    
    # Generate content based on sheet type
    if sheet_type == "bom":
        prompt = f"""create a Bill of Materials (BOM) for "{topic}".
        
Return JSON with this structure:
{{
    "title": "BOM Title",
    "items": [
        {{
            "item_no": "1",
            "description": "Item description",
            "quantity": 1,
            "unit": "EA",
            "unit_price": 100.00,
            "total": 100.00
        }}
    ],
    "grand_total": 500.00
}}

Generate 5-10 items. JSON only:"""
    else:
        prompt = f"""create a pricing sheet for "{topic}".
        
Return JSON with this structure:
{{
    "title": "Pricing Sheet Title",
    "items": [
        {{
            "service": "Service name",
            "description": "Description",
            "price": 1000.00
        }}
    ],
    "total": 5000.00
}}

Generate 5-8 items. JSON only:"""
    
    content = generate_content(prompt, context)
    
    # Parse JSON
    try:
        json_str = content.strip()
        if json_str.startswith("```"):
            json_str = json_str.split("```")[1]
            if json_str.startswith("json"):
                json_str = json_str[4:]
        json_str = json_str.strip()
        data = json.loads(json_str)
    except Exception as e:
        print(f"❌ Failed to parse content: {e}")
        return None
    
    # Load template or create new
    if template_path and os.path.exists(template_path):
        wb = load_workbook(template_path)
        ws = wb.active
    else:
        wb = Workbook()
        ws = wb.active
    
    # Styles
    header_font = Font(bold=True, size=12)
    header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
    header_font_white = Font(bold=True, size=12, color="FFFFFF")
    border = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin')
    )
    
    # Title
    ws['A1'] = data.get("title", topic)
    ws['A1'].font = Font(bold=True, size=16)
    ws.merge_cells('A1:F1')
    
    # Date
    ws['A2'] = f"Date: {datetime.now().strftime('%B %d, %Y')}"
    ws['A3'] = ""
    
    if sheet_type == "bom":
        # Headers
        headers = ["Item No", "Description", "Quantity", "Unit", "Unit Price", "Total"]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=4, column=col, value=header)
            cell.font = header_font_white
            cell.fill = header_fill
            cell.border = border
            cell.alignment = Alignment(horizontal='center')
        
        # Data
        for row, item in enumerate(data.get("items", []), 5):
            ws.cell(row=row, column=1, value=item.get("item_no", "")).border = border
            ws.cell(row=row, column=2, value=item.get("description", "")).border = border
            ws.cell(row=row, column=3, value=item.get("quantity", 0)).border = border
            ws.cell(row=row, column=4, value=item.get("unit", "")).border = border
            ws.cell(row=row, column=5, value=item.get("unit_price", 0)).border = border
            ws.cell(row=row, column=6, value=item.get("total", 0)).border = border
        
        # Grand Total
        last_row = 5 + len(data.get("items", []))
        ws.cell(row=last_row, column=5, value="Grand Total:").font = Font(bold=True)
        ws.cell(row=last_row, column=6, value=data.get("grand_total", 0)).font = Font(bold=True)
    
    else:  # pricing
        # Headers
        headers = ["Service", "Description", "Price"]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=4, column=col, value=header)
            cell.font = header_font_white
            cell.fill = header_fill
            cell.border = border
            cell.alignment = Alignment(horizontal='center')
        
        # Data
        for row, item in enumerate(data.get("items", []), 5):
            ws.cell(row=row, column=1, value=item.get("service", "")).border = border
            ws.cell(row=row, column=2, value=item.get("description", "")).border = border
            ws.cell(row=row, column=3, value=item.get("price", 0)).border = border
        
        # Total
        last_row = 5 + len(data.get("items", []))
        ws.cell(row=last_row, column=2, value="Total:").font = Font(bold=True)
        ws.cell(row=last_row, column=3, value=data.get("total", 0)).font = Font(bold=True)
    
    # Adjust column widths
    ws.column_dimensions['A'].width = 12
    ws.column_dimensions['B'].width = 40
    ws.column_dimensions['C'].width = 12
    ws.column_dimensions['D'].width = 10
    ws.column_dimensions['E'].width = 15
    ws.column_dimensions['F'].width = 15
    
    # Save
    wb.save(output_path)
    print(f"✅ Excel saved: {output_path}")
    return output_path

def create_excel_simple(output_path: str, topic: str, sheet_type: str = "bom"):
    """Create Excel without template"""
    return create_excel_from_template(None, output_path, topic, sheet_type)

# ==================== Main ====================

if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 4:
        print("Usage:")
        print("  python document_generator.py ppt <output.pptx> <topic>")
        print("  python document_generator.py word <output.docx> <topic>")
        print("  python document_generator.py excel <output.xlsx> <topic>")
        print("")
        print("Examples:")
        print("  python document_generator.py ppt presentation.pptx 'Boomi Integration'")
        print("  python document_generator.py word proposal.docx 'API Management Solution'")
        print("  python document_generator.py excel bom.xlsx 'Integration Project BOM'")
    else:
        doc_type = sys.argv[1].lower()
        output_path = sys.argv[2]
        topic = " ".join(sys.argv[3:])
        
        if doc_type == "ppt":
            create_ppt_simple(output_path, topic)
        elif doc_type == "word":
            create_word_simple(output_path, topic)
        elif doc_type == "excel":
            create_excel_simple(output_path, topic)
        else:
            print(f"❌ Unknown document type: {doc_type}")