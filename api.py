from fastapi import FastAPI, UploadFile, File, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_experimental.text_splitter import SemanticChunker
from pinecone import Pinecone
from dotenv import load_dotenv
from datetime import datetime, timedelta
from youtube_transcript_api import YouTubeTranscriptApi
from bs4 import BeautifulSoup
from rank_bm25 import BM25Okapi
import requests
import os
import uuid
import tempfile
import re
import json
import cohere
import numpy as np
from typing import List, Optional
import asyncio
from openai import OpenAI  # For Agentic RAG with function calling
# Import configuration
from config import (
    AgentConfig, PRESET_CONFIGS, get_config, get_preset_info,
    AGENT_SYSTEM_PROMPT, QUERY_REWRITE_PROMPT,
    STOP_WORDS, EXPAND_BUTTON_SELECTORS, load_dynamic_domains,
    API_PRICING, validate_temperature, validate_iterations, validate_top_k
)


# Unstructured.io imports for smart document parsing
try:
    from unstructured.partition.html import partition_html
    from unstructured.partition.pdf import partition_pdf
    from unstructured.partition.docx import partition_docx
    from unstructured.partition.xlsx import partition_xlsx
    from unstructured.partition.pptx import partition_pptx
    from unstructured.partition.text import partition_text
    from unstructured.chunking.title import chunk_by_title
    UNSTRUCTURED_AVAILABLE = True
    print("Ã¢Å“â€¦ Unstructured.io loaded successfully")
except ImportError as e:
    UNSTRUCTURED_AVAILABLE = False
    print(f"âš ï¸ Unstructured.io not available: {e}. Using fallback parsing.")

# Playwright for JavaScript-rendered pages
try:
    from playwright.async_api import async_playwright
    PLAYWRIGHT_AVAILABLE = True
    print("Ã¢Å“â€¦ Playwright loaded successfully")
except ImportError as e:
    PLAYWRIGHT_AVAILABLE = False
    print(f"âš ï¸ Playwright not available: {e}. Dynamic pages may not render fully.")

# Load environment variables
load_dotenv()

# Initialize
pc = Pinecone(api_key=os.getenv("PINECONE_API_KEY"))
index = pc.Index(os.getenv("PINECONE_INDEX_NAME"))
embeddings = OpenAIEmbeddings(model="text-embedding-3-large")
llm = ChatOpenAI(model="gpt-4o", temperature=0.7)

# OpenAI client for Agentic RAG (function calling)
openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# Initialize Cohere for Reranking
co = cohere.Client(os.getenv("COHERE_API_KEY"))

# Semantic Chunker - splits by meaning
# Get config for chunking settings
_config = AgentConfig()

semantic_chunker = SemanticChunker(
    embeddings=embeddings,
    breakpoint_threshold_type="percentile",
    breakpoint_threshold_amount=_config.semantic_threshold  # 75 for better context
)

# Fallback character-based splitter
fallback_splitter = RecursiveCharacterTextSplitter(
    chunk_size=_config.chunk_size,  # 3500 - larger chunks for better context
    chunk_overlap=_config.chunk_overlap  # 350 - overlap to prevent context loss
)

# Folders - Use Railway Volume if available, otherwise local
DATA_DIR = "/app/data" if os.path.exists("/app/data") else "."
UPLOADS_FOLDER = os.path.join(DATA_DIR, "uploads")
CONVERSATIONS_FOLDER = os.path.join(DATA_DIR, "conversations")
DOCUMENTS_FOLDER = os.path.join(DATA_DIR, "documents")
BM25_INDEX_FILE = os.path.join(DATA_DIR, "bm25_index.json")

os.makedirs(UPLOADS_FOLDER, exist_ok=True)
os.makedirs(CONVERSATIONS_FOLDER, exist_ok=True)
os.makedirs(DOCUMENTS_FOLDER, exist_ok=True)

# FastAPI app
app = FastAPI(title="RAG Agent API", version="3.0.0")  # Major update with configurable settings

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Request/Response models
class ChatRequest(BaseModel):
    message: str
    user_id: str = "default_user"
    conversation_id: Optional[str] = None

class Source(BaseModel):
    type: str
    title: Optional[str] = None
    source: str
    score: float
    download_url: Optional[str] = None

class ImageInfo(BaseModel):
    url: str
    alt: Optional[str] = None
    title: Optional[str] = None
    context: Optional[str] = None
    doc_title: Optional[str] = None

class ChatResponse(BaseModel):
    answer: str
    memories_found: int
    sources: List[Source] = []
    images: List[ImageInfo] = []
    conversation_id: str

class IngestResponse(BaseModel):
    success: bool
    message: str
    chunks: int = 0
    chunker_type: str = "semantic"
    element_types: Optional[dict] = None  # NEW: {"Table": 2, "NarrativeText": 5, ...}

class UrlRequest(BaseModel):
    url: str

# ==================== Pipeline Monitor ====================

import time
from dataclasses import dataclass, field, asdict
from typing import Dict, Any
from collections import deque

# Pipeline data storage
PIPELINE_LOG_FILE = os.path.join(DATA_DIR, "pipeline_logs.json")
MAX_PIPELINE_LOGS = 100  # Keep last 100 pipeline runs

@dataclass
class PipelineStep:
    """Represents a single step in the pipeline"""
    name: str
    status: str = "pending"  # pending, running, completed, error
    start_time: float = 0
    end_time: float = 0
    duration_ms: float = 0
    details: Dict[str, Any] = field(default_factory=dict)
    
    def start(self):
        self.status = "running"
        self.start_time = time.time()
    
    def complete(self, **details):
        self.status = "completed"
        self.end_time = time.time()
        self.duration_ms = round((self.end_time - self.start_time) * 1000, 2)
        self.details.update(details)
    
    def error(self, message: str):
        self.status = "error"
        self.end_time = time.time()
        self.duration_ms = round((self.end_time - self.start_time) * 1000, 2)
        self.details["error"] = message
    
    def to_dict(self):
        return {
            "name": self.name,
            "status": self.status,
            "duration_ms": self.duration_ms,
            "details": self.details
        }

class PipelineTracker:
    """Tracks pipeline execution for monitoring"""
    
    def __init__(self, pipeline_type: str, source: str = ""):
        self.id = f"pipe_{uuid.uuid4().hex[:12]}"
        self.pipeline_type = pipeline_type  # "ingestion" or "retrieval"
        self.source = source
        self.steps: Dict[str, PipelineStep] = {}
        self.start_time = time.time()
        self.end_time = 0
        self.total_duration_ms = 0
        self.status = "running"
        self.metadata: Dict[str, Any] = {}
        
        # Define steps based on pipeline type
        if pipeline_type == "ingestion":
            step_names = ["extract", "chunk", "embed", "store_vectors", "store_bm25"]
        else:  # retrieval
            step_names = ["embed_query", "pinecone_search", "bm25_search", "rrf_merge", "rerank", "generate"]
        
        for name in step_names:
            self.steps[name] = PipelineStep(name=name)
    
    def start_step(self, step_name: str):
        if step_name in self.steps:
            self.steps[step_name].start()
    
    def complete_step(self, step_name: str, **details):
        if step_name in self.steps:
            self.steps[step_name].complete(**details)
    
    def error_step(self, step_name: str, message: str):
        if step_name in self.steps:
            self.steps[step_name].error(message)
    
    def finish(self, **metadata):
        self.end_time = time.time()
        self.total_duration_ms = round((self.end_time - self.start_time) * 1000, 2)
        self.status = "completed"
        self.metadata.update(metadata)
        
        # Save to log
        save_pipeline_log(self.to_dict())
    
    def to_dict(self):
        return {
            "id": self.id,
            "pipeline_type": self.pipeline_type,
            "source": self.source,
            "status": self.status,
            "start_time": datetime.fromtimestamp(self.start_time).isoformat(),
            "total_duration_ms": self.total_duration_ms,
            "steps": {name: step.to_dict() for name, step in self.steps.items()},
            "metadata": self.metadata
        }

def load_pipeline_logs() -> List[dict]:
    """Load pipeline logs from file"""
    if os.path.exists(PIPELINE_LOG_FILE):
        try:
            with open(PIPELINE_LOG_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except:
            return []
    return []

def save_pipeline_log(log: dict):
    """Save a pipeline log entry"""
    logs = load_pipeline_logs()
    logs.append(log)
    # Keep only last MAX_PIPELINE_LOGS entries
    if len(logs) > MAX_PIPELINE_LOGS:
        logs = logs[-MAX_PIPELINE_LOGS:]
    with open(PIPELINE_LOG_FILE, 'w', encoding='utf-8') as f:
        json.dump(logs, f, indent=2)

# Current active pipelines (for real-time monitoring)
active_pipelines: Dict[str, PipelineTracker] = {}

# ==================== Cost Monitor ====================

# API Pricing (as of December 2024)
API_PRICING = {
    "openai": {
        "text-embedding-3-large": {
            "price_per_million_tokens": 0.13,
            "unit": "tokens"
        },
        "text-embedding-3-small": {
            "price_per_million_tokens": 0.02,
            "unit": "tokens"
        },
        "gpt-4o": {
            "input_price_per_million_tokens": 2.50,
            "output_price_per_million_tokens": 10.00,
            "unit": "tokens"
        },
        "gpt-4o-mini": {
            "input_price_per_million_tokens": 0.15,
            "output_price_per_million_tokens": 0.60,
            "unit": "tokens"
        },
        "whisper-1": {
            "price_per_minute": 0.006,
            "unit": "minutes"
        }
    },
    "cohere": {
        "rerank-v3.5": {
            "price_per_search": 0.0005,
            "unit": "searches"
        }
    },
    "pinecone": {
        "serverless": {
            "price_per_query": 0.0,  # Free tier: 100K queries/month
            "unit": "queries"
        }
    }
}

# Cost data storage
COST_LOG_FILE = os.path.join(DATA_DIR, "cost_logs.json")
COST_SUMMARY_FILE = os.path.join(DATA_DIR, "cost_summary.json")

@dataclass
class CostEntry:
    """Represents a single cost entry"""
    id: str
    timestamp: str
    service: str  # openai, cohere, pinecone
    model: str
    operation: str  # embedding, generation, transcription, reranking
    quantity: float  # tokens, minutes, searches
    unit: str
    unit_price: float
    total_cost: float
    source: str = ""  # document name or query
    pipeline_id: str = ""
    details: Dict[str, Any] = field(default_factory=dict)
    
    def to_dict(self):
        return {
            "id": self.id,
            "timestamp": self.timestamp,
            "service": self.service,
            "model": self.model,
            "operation": self.operation,
            "quantity": self.quantity,
            "unit": self.unit,
            "unit_price": self.unit_price,
            "total_cost": self.total_cost,
            "source": self.source,
            "pipeline_id": self.pipeline_id,
            "details": self.details
        }

class CostTracker:
    """Track API costs for monitoring"""
    
    @staticmethod
    def estimate_tokens(text: str) -> int:
        """Rough estimate: ~4 characters per token for English"""
        return max(1, len(text) // 4)
    
    @staticmethod
    def calculate_embedding_cost(text: str, model: str = "text-embedding-3-large") -> dict:
        """Calculate cost for embedding text"""
        tokens = CostTracker.estimate_tokens(text)
        pricing = API_PRICING["openai"].get(model, API_PRICING["openai"]["text-embedding-3-large"])
        cost = (tokens / 1_000_000) * pricing["price_per_million_tokens"]
        
        return {
            "tokens": tokens,
            "cost": round(cost, 8),
            "model": model,
            "formula": f"({tokens} ÃƒÂ· 1,000,000) Ãƒâ€” ${pricing['price_per_million_tokens']}",
            "text_preview": text[:200] + "..." if len(text) > 200 else text,
            "text_length": len(text)
        }
    
    @staticmethod
    def calculate_llm_cost(input_text: str, output_text: str, model: str = "gpt-4o", breakdown: dict = None) -> dict:
        """Calculate cost for LLM generation with detailed breakdown"""
        input_tokens = CostTracker.estimate_tokens(input_text)
        output_tokens = CostTracker.estimate_tokens(output_text)
        pricing = API_PRICING["openai"].get(model, API_PRICING["openai"]["gpt-4o"])
        
        input_cost = (input_tokens / 1_000_000) * pricing["input_price_per_million_tokens"]
        output_cost = (output_tokens / 1_000_000) * pricing["output_price_per_million_tokens"]
        total_cost = input_cost + output_cost
        
        result = {
            "input_tokens": input_tokens,
            "output_tokens": output_tokens,
            "total_tokens": input_tokens + output_tokens,
            "input_cost": round(input_cost, 8),
            "output_cost": round(output_cost, 8),
            "total_cost": round(total_cost, 8),
            "model": model,
            "formula": f"Input: ({input_tokens} ÃƒÂ· 1M) Ãƒâ€” ${pricing['input_price_per_million_tokens']} + Output: ({output_tokens} ÃƒÂ· 1M) Ãƒâ€” ${pricing['output_price_per_million_tokens']}",
            "input_preview": input_text[:300] + "..." if len(input_text) > 300 else input_text,
            "output_preview": output_text[:300] + "..." if len(output_text) > 300 else output_text
        }
        
        # Add detailed breakdown if provided
        if breakdown:
            result["breakdown"] = breakdown
        
        return result
    
    @staticmethod
    def calculate_prompt_breakdown(system_prompt: str, conversation_history: str, 
                                   document_context: str, user_question: str) -> dict:
        """Calculate detailed breakdown of prompt components"""
        components = {
            "system_prompt": {
                "text": system_prompt,
                "tokens": CostTracker.estimate_tokens(system_prompt),
                "percentage": 0
            },
            "conversation_history": {
                "text": conversation_history,
                "tokens": CostTracker.estimate_tokens(conversation_history) if conversation_history else 0,
                "percentage": 0
            },
            "document_context": {
                "text": document_context[:500] + "..." if len(document_context) > 500 else document_context,
                "tokens": CostTracker.estimate_tokens(document_context) if document_context else 0,
                "percentage": 0
            },
            "user_question": {
                "text": user_question,
                "tokens": CostTracker.estimate_tokens(user_question),
                "percentage": 0
            }
        }
        
        total_tokens = sum(c["tokens"] for c in components.values())
        
        # Calculate percentages
        for key in components:
            if total_tokens > 0:
                components[key]["percentage"] = round((components[key]["tokens"] / total_tokens) * 100, 1)
        
        return {
            "components": components,
            "total_input_tokens": total_tokens
        }
    
    @staticmethod
    def generate_optimization_suggestions(question: str, context_tokens: int, 
                                          total_tokens: int, operation: str = "retrieval") -> list:
        """Generate smart suggestions to optimize token usage"""
        suggestions = []
        
        if operation == "retrieval":
            # Question optimization
            question_tokens = CostTracker.estimate_tokens(question)
            
            # Check for unnecessary words
            filler_words = ["please", "can you", "could you", "i want to", "i need to", 
                          "help me", "tell me", "explain to me", "i would like"]
            question_lower = question.lower()
            found_fillers = [w for w in filler_words if w in question_lower]
            
            if found_fillers:
                suggestions.append({
                    "type": "question",
                    "priority": "medium",
                    "icon": "Ã¢Å“â€šï¸",
                    "title": "Remove filler words",
                    "description": f"Words like '{', '.join(found_fillers[:2])}' can be removed",
                    "potential_savings": f"~{len(found_fillers) * 3} tokens",
                    "example": {
                        "before": question[:100],
                        "after": question[:100].replace("please ", "").replace("can you ", "").replace("I want to ", "")
                    }
                })
            
            # Check question length
            if question_tokens > 50:
                suggestions.append({
                    "type": "question",
                    "priority": "high",
                    "icon": "ðŸ“Â",
                    "title": "Shorten your question",
                    "description": f"Your question uses {question_tokens} tokens. Try to be more concise.",
                    "potential_savings": f"~{question_tokens - 30} tokens"
                })
            
            # Context optimization
            if context_tokens > 5000:
                suggestions.append({
                    "type": "context",
                    "priority": "high",
                    "icon": "ðŸ“Å¡",
                    "title": "Large context detected",
                    "description": f"Document context is {context_tokens} tokens ({round(context_tokens/total_tokens*100)}% of total)",
                    "potential_savings": "Consider using more specific questions to retrieve smaller context"
                })
            
            # Cost warning
            if total_tokens > 8000:
                suggestions.append({
                    "type": "cost",
                    "priority": "warning",
                    "icon": "ðŸ’°",
                    "title": "High token usage",
                    "description": f"This query uses ~{total_tokens} tokens (Ã¢â€°Ë†${round(total_tokens/1000000*12.5, 4)})",
                    "potential_savings": "Break complex questions into smaller ones"
                })
        
        elif operation == "ingestion":
            if total_tokens > 10000:
                suggestions.append({
                    "type": "document",
                    "priority": "info",
                    "icon": "ðŸ“„",
                    "title": "Large document",
                    "description": f"Document has ~{total_tokens} tokens",
                    "potential_savings": "Consider splitting into smaller sections if not all content is needed"
                })
        
        return suggestions
    
    @staticmethod
    def calculate_whisper_cost(duration_seconds: float) -> dict:
        """Calculate cost for Whisper transcription"""
        duration_minutes = duration_seconds / 60
        pricing = API_PRICING["openai"]["whisper-1"]
        cost = duration_minutes * pricing["price_per_minute"]
        
        return {
            "duration_seconds": duration_seconds,
            "duration_minutes": round(duration_minutes, 2),
            "cost": round(cost, 8),
            "model": "whisper-1",
            "formula": f"{round(duration_minutes, 2)} minutes Ãƒâ€” ${pricing['price_per_minute']}/min"
        }
    
    @staticmethod
    def calculate_rerank_cost(num_searches: int = 1) -> dict:
        """Calculate cost for Cohere reranking"""
        pricing = API_PRICING["cohere"]["rerank-v3.5"]
        cost = num_searches * pricing["price_per_search"]
        
        return {
            "searches": num_searches,
            "cost": round(cost, 8),
            "model": "rerank-v3.5",
            "formula": f"{num_searches} Ãƒâ€” ${pricing['price_per_search']}/search"
        }
    
    @staticmethod
    def log_cost(service: str, model: str, operation: str, quantity: float, 
                 unit: str, unit_price: float, total_cost: float,
                 source: str = "", pipeline_id: str = "", **details):
        """Log a cost entry"""
        entry = CostEntry(
            id=f"cost_{uuid.uuid4().hex[:12]}",
            timestamp=datetime.now().isoformat(),
            service=service,
            model=model,
            operation=operation,
            quantity=quantity,
            unit=unit,
            unit_price=unit_price,
            total_cost=total_cost,
            source=source,
            pipeline_id=pipeline_id,
            details=details
        )
        
        # Save to log
        save_cost_log(entry.to_dict())
        
        # Update summary
        update_cost_summary(entry)
        
        return entry

def load_cost_logs() -> List[dict]:
    """Load cost logs from file"""
    if os.path.exists(COST_LOG_FILE):
        try:
            with open(COST_LOG_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except:
            return []
    return []

def save_cost_log(entry: dict):
    """Save a cost log entry"""
    logs = load_cost_logs()
    logs.append(entry)
    # Keep last 500 entries
    if len(logs) > 500:
        logs = logs[-500:]
    with open(COST_LOG_FILE, 'w', encoding='utf-8') as f:
        json.dump(logs, f, indent=2)

def load_cost_summary() -> dict:
    """Load cost summary"""
    if os.path.exists(COST_SUMMARY_FILE):
        try:
            with open(COST_SUMMARY_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except:
            pass
    return {
        "total_cost": 0,
        "by_service": {},
        "by_operation": {},
        "by_day": {},
        "last_updated": None
    }

def save_cost_summary(summary: dict):
    """Save cost summary"""
    with open(COST_SUMMARY_FILE, 'w', encoding='utf-8') as f:
        json.dump(summary, f, indent=2)

def update_cost_summary(entry: CostEntry):
    """Update cost summary with new entry"""
    summary = load_cost_summary()
    
    # Total cost
    summary["total_cost"] = round(summary.get("total_cost", 0) + entry.total_cost, 8)
    
    # By service
    if entry.service not in summary["by_service"]:
        summary["by_service"][entry.service] = {"cost": 0, "count": 0}
    summary["by_service"][entry.service]["cost"] = round(
        summary["by_service"][entry.service]["cost"] + entry.total_cost, 8
    )
    summary["by_service"][entry.service]["count"] += 1
    
    # By operation
    if entry.operation not in summary["by_operation"]:
        summary["by_operation"][entry.operation] = {"cost": 0, "count": 0}
    summary["by_operation"][entry.operation]["cost"] = round(
        summary["by_operation"][entry.operation]["cost"] + entry.total_cost, 8
    )
    summary["by_operation"][entry.operation]["count"] += 1
    
    # By day
    day = entry.timestamp[:10]  # YYYY-MM-DD
    if day not in summary["by_day"]:
        summary["by_day"][day] = {"cost": 0, "count": 0}
    summary["by_day"][day]["cost"] = round(
        summary["by_day"][day]["cost"] + entry.total_cost, 8
    )
    summary["by_day"][day]["count"] += 1
    
    summary["last_updated"] = datetime.now().isoformat()
    
    save_cost_summary(summary)

# ==================== Semantic Chunking ====================

def smart_chunk_text(text: str, min_chunk_size: int = 100, config: AgentConfig = None) -> tuple[List[str], str]:
    """
    Split text into semantic chunks with improved settings.
    Returns: (chunks, chunker_type) where chunker_type is 'semantic' or 'fallback'
    """
    if config is None:
        config = AgentConfig()
    
    if len(text) < min_chunk_size:
        return [text], "semantic"
    
    try:
        chunks = semantic_chunker.split_text(text)
        
        merged_chunks = []
        current_chunk = ""
        max_chunk_size = config.chunk_size  # Use config value (3500)
        
        for chunk in chunks:
            if len(current_chunk) + len(chunk) < max_chunk_size:
                current_chunk += "\n\n" + chunk if current_chunk else chunk
            else:
                if current_chunk:
                    merged_chunks.append(current_chunk.strip())
                current_chunk = chunk
        
        if current_chunk:
            merged_chunks.append(current_chunk.strip())
        
        # If we got very few chunks from a large text, use fallback
        if len(merged_chunks) < 2 and len(text) > 3000:
            return fallback_splitter.split_text(text), "fallback"
        
        return merged_chunks, "semantic"
    
    except Exception as e:
        print(f"Semantic chunking failed: {e}, using fallback")
        return fallback_splitter.split_text(text), "fallback"

# ==================== Unstructured.io Smart Parsing ====================

def parse_with_unstructured(content, content_type: str = "html", source_url: str = None) -> tuple[List[dict], str]:
    """
    Parse content using Unstructured.io for smart element extraction.
    Returns: (chunks_with_metadata, chunker_type)
    Each chunk has: {"text": str, "element_type": str, "metadata": dict}
    """
    if not UNSTRUCTURED_AVAILABLE:
        if isinstance(content, str):
            chunks, chunk_type = smart_chunk_text(content)
            return [{"text": c, "element_type": "NarrativeText", "metadata": {}} for c in chunks], chunk_type
        return [], "fallback"
    
    try:
        elements = []
        
        if content_type == "html":
            if isinstance(content, str):
                elements = partition_html(text=content)
            else:
                elements = partition_html(file=content)
        elif content_type == "pdf":
            elements = partition_pdf(filename=content)
        elif content_type == "docx":
            elements = partition_docx(filename=content)
        elif content_type == "xlsx":
            elements = partition_xlsx(filename=content)
        elif content_type == "pptx":
            elements = partition_pptx(filename=content)
        elif content_type == "text":
            if isinstance(content, str):
                elements = partition_text(text=content)
            else:
                elements = partition_text(filename=content)
        else:
            if isinstance(content, str):
                chunks, chunk_type = smart_chunk_text(content)
                return [{"text": c, "element_type": "NarrativeText", "metadata": {}} for c in chunks], chunk_type
        
        if not elements:
            return [], "unstructured"
        
        # Use chunk_by_title for structure-aware chunking
        chunked_elements = chunk_by_title(
            elements,
            max_characters=4000,  # Larger chunks for better context
            combine_text_under_n_chars=300,  # Combine small elements
            new_after_n_chars=3500
        )
        
        chunks_with_metadata = []
        for i, chunk in enumerate(chunked_elements):
            element_type = type(chunk).__name__
            text = str(chunk)
            
            chunk_metadata = {
                "element_type": element_type,
                "chunk_index": i,
                "has_table": element_type == "Table",
            }
            
            if hasattr(chunk, 'metadata') and hasattr(chunk.metadata, 'page_number'):
                chunk_metadata["page_number"] = chunk.metadata.page_number
            
            chunks_with_metadata.append({
                "text": text,
                "element_type": element_type,
                "metadata": chunk_metadata
            })
        
        print(f"ðŸ“¦ Unstructured created {len(chunks_with_metadata)} chunks")
        
        type_counts = {}
        for c in chunks_with_metadata:
            et = c["element_type"]
            type_counts[et] = type_counts.get(et, 0) + 1
        print(f"ðŸ“Š Element types: {type_counts}")
        
        return chunks_with_metadata, "unstructured"
        
    except Exception as e:
        print(f"âš ï¸ Unstructured parsing failed: {e}, using fallback")
        if isinstance(content, str):
            chunks, chunk_type = smart_chunk_text(content)
            return [{"text": c, "element_type": "NarrativeText", "metadata": {}} for c in chunks], chunk_type
        return [], "fallback"

def parse_html_with_unstructured(html_content: str, url: str = None) -> tuple[str, List[dict], str]:
    """
    Parse HTML content using Unstructured.io.
    Returns: (full_text, chunks_with_metadata, chunker_type)
    """
    if not UNSTRUCTURED_AVAILABLE:
        soup = BeautifulSoup(html_content, 'html.parser')
        for element in soup(['script', 'style', 'nav', 'footer', 'header']):
            element.decompose()
        text = soup.get_text(separator='\n', strip=True)
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        full_text = '\n'.join(lines)
        chunks, chunk_type = smart_chunk_text(full_text)
        chunks_with_meta = [{"text": c, "element_type": "NarrativeText", "metadata": {}} for c in chunks]
        return full_text, chunks_with_meta, chunk_type
    
    try:
        elements = partition_html(text=html_content)
        
        if not elements:
            return "", [], "unstructured"
        
        full_text_parts = []
        for el in elements:
            el_type = type(el).__name__
            el_text = str(el)
            
            if el_type == "Table":
                full_text_parts.append(f"\n[TABLE]\n{el_text}\n[/TABLE]\n")
            else:
                full_text_parts.append(el_text)
        
        full_text = "\n\n".join(full_text_parts)
        
        chunked_elements = chunk_by_title(
            elements,
            max_characters=4000,  # Larger chunks for better context
            combine_text_under_n_chars=300,  # Combine small elements
            new_after_n_chars=3500
        )
        
        chunks_with_metadata = []
        for i, chunk in enumerate(chunked_elements):
            element_type = type(chunk).__name__
            text = str(chunk)
            
            chunks_with_metadata.append({
                "text": text,
                "element_type": element_type,
                "metadata": {"chunk_index": i, "source_url": url}
            })
        
        print(f"Ã¢Å“â€¦ Unstructured HTML: {len(elements)} elements Ã¢â€ â€™ {len(chunks_with_metadata)} chunks")
        
        return full_text, chunks_with_metadata, "unstructured"
        
    except Exception as e:
        print(f"âš ï¸ Unstructured HTML parsing failed: {e}")
        soup = BeautifulSoup(html_content, 'html.parser')
        for element in soup(['script', 'style', 'nav', 'footer', 'header']):
            element.decompose()
        text = soup.get_text(separator='\n', strip=True)
        chunks, chunk_type = smart_chunk_text(text)
        chunks_with_meta = [{"text": c, "element_type": "NarrativeText", "metadata": {}} for c in chunks]
        return text, chunks_with_meta, chunk_type

# ==================== Playwright for Dynamic Pages ====================

async def fetch_page_with_playwright(url: str, wait_time: int = 5000, scroll: bool = True, progress_callback=None) -> str:
    """
    Fetch a page using Playwright headless browser (async version).
    This renders JavaScript and returns the fully loaded HTML.
    
    Args:
        url: The URL to fetch
        wait_time: Time to wait for JS to render (ms)
        scroll: Whether to scroll the page to trigger lazy loading
        progress_callback: Optional async function to report progress
    
    Returns:
        The fully rendered HTML content
    """
    async def log_progress(message):
        print(message)
        if progress_callback:
            await progress_callback(message)
    
    if not PLAYWRIGHT_AVAILABLE:
        await log_progress("âš ï¸ Playwright not available, falling back to requests")
        response = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'}, timeout=15)
        return response.text
    
    try:
        await log_progress(f"ðŸŽ­ [1/7] Starting Playwright browser...")
        
        async with async_playwright() as p:
            await log_progress(f"ðŸŽ­ [2/7] Launching Chromium...")
            browser = await p.chromium.launch(headless=True)
            
            await log_progress(f"ðŸŽ­ [3/7] Creating browser context...")
            context = await browser.new_context(
                user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                viewport={'width': 1920, 'height': 1080}
            )
            page = await context.new_page()
            
            await log_progress(f"ðŸŽ­ [4/7] Navigating to {url[:50]}...")
            await page.goto(url, wait_until='load', timeout=60000)
            
            await log_progress(f"ðŸŽ­ [5/7] Waiting for JavaScript to render ({wait_time/1000}s)...")
            await page.wait_for_timeout(wait_time)
            
            # Try to click "Expand All" or similar buttons if they exist
            # Try to expand any collapsed sections using generic selectors
            try:
                all_expand_buttons = []
                for selector in EXPAND_BUTTON_SELECTORS[:15]:  # Try first 15 selectors
                    try:
                        buttons = await page.query_selector_all(selector)
                        all_expand_buttons.extend(buttons)
                    except:
                        pass
                
                if all_expand_buttons:
                    await log_progress(f"ðŸŽ­ [5.5/7] Found {len(all_expand_buttons)} expandable sections, clicking...")
                    for btn in all_expand_buttons[:15]:  # Click up to 15 buttons
                        try:
                            await btn.click()
                            await page.wait_for_timeout(300)
                        except:
                            pass
            except:
                pass
            
            # Wait for tables specifically if present
            try:
                await page.wait_for_selector('table', timeout=3000)
                await log_progress(f"ðŸŽ­ [5.6/7] Tables detected, waiting for data...")
                await page.wait_for_timeout(1500)
            except:
                pass  # No tables or timeout - continue
            
            if scroll:
                await log_progress(f"ðŸŽ­ [6/7] Scrolling page to load dynamic content...")
                for i in range(8):
                    await page.evaluate('window.scrollTo(0, document.body.scrollHeight * {})'.format((i + 1) / 8))
                    await page.wait_for_timeout(800)
                
                await page.evaluate('window.scrollTo(0, 0)')
                await page.wait_for_timeout(500)
            
            await log_progress(f"ðŸŽ­ [6.5/7] Final wait for lazy-loaded content...")
            await page.wait_for_timeout(2000)
            
            await log_progress(f"ðŸŽ­ [7/7] Extracting rendered HTML...")
            html_content = await page.content()
            
            await browser.close()
            
            await log_progress(f"Ã¢Å“â€¦ Playwright fetched {len(html_content):,} bytes from {url[:50]}")
            return html_content
            
    except Exception as e:
        await log_progress(f"âš ï¸ Playwright fetch failed: {e}, falling back to requests")
        try:
            response = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'}, timeout=15)
            return response.text
        except Exception as req_e:
            await log_progress(f"Ã¢ÂÅ’ Requests also failed: {req_e}")
            raise

def should_use_playwright(url: str, force_playwright: bool = None) -> bool:
    """
    Determine if a URL should use Playwright.
    
    Args:
        url: The URL to check
        force_playwright: If True, always use. If False, never use. If None, auto-detect.
    
    Returns:
        Boolean indicating whether to use Playwright
    """
    # User override
    if force_playwright is True:
        return True
    if force_playwright is False:
        return False
    
    # Auto-detect based on configurable domain list
    dynamic_domains = load_dynamic_domains()
    
    for domain in dynamic_domains:
        if domain in url.lower():
            return True
    
    return False

# ==================== BM25 Index Management ====================

def load_bm25_index() -> dict:
    if os.path.exists(BM25_INDEX_FILE):
        with open(BM25_INDEX_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    return {"documents": [], "metadata": []}

def save_bm25_index(bm25_data: dict):
    with open(BM25_INDEX_FILE, 'w', encoding='utf-8') as f:
        json.dump(bm25_data, f, ensure_ascii=False)

def add_to_bm25_index(chunks: List[str], metadata_list: List[dict]):
    bm25_data = load_bm25_index()
    
    for chunk, meta in zip(chunks, metadata_list):
        bm25_data["documents"].append(chunk)
        bm25_data["metadata"].append(meta)
    
    save_bm25_index(bm25_data)

def tokenize_for_bm25(text: str) -> List[str]:
    """
    Improved tokenization for BM25 with stemming-like behavior and stop word removal.
    """
    # Remove punctuation and convert to lowercase
    text = re.sub(r'[^\w\s]', '', text.lower())
    
    # Tokenize
    tokens = text.split()
    
    # Remove stop words and very short tokens
    filtered_tokens = [t for t in tokens if t not in STOP_WORDS and len(t) > 2]
    
    # Simple suffix removal (poor man's stemming)
    stemmed_tokens = []
    for token in filtered_tokens:
        # Remove common suffixes
        if token.endswith('ing'):
            token = token[:-3]
        elif token.endswith('tion'):
            token = token[:-4]
        elif token.endswith('ed') and len(token) > 4:
            token = token[:-2]
        elif token.endswith('ly') and len(token) > 4:
            token = token[:-2]
        elif token.endswith('ment') and len(token) > 5:
            token = token[:-4]
        elif token.endswith('ness') and len(token) > 5:
            token = token[:-4]
        elif token.endswith('able') and len(token) > 5:
            token = token[:-4]
        elif token.endswith('ies'):
            token = token[:-3] + 'y'
        elif token.endswith('es') and len(token) > 3:
            token = token[:-2]
        elif token.endswith('s') and len(token) > 3:
            token = token[:-1]
        
        if len(token) > 2:
            stemmed_tokens.append(token)
    
    return stemmed_tokens


def search_bm25(query: str, top_k: int = 10) -> List[dict]:
    """Search using BM25 with improved tokenization."""
    bm25_data = load_bm25_index()
    
    if not bm25_data["documents"]:
        return []
    
    # Use improved tokenization
    tokenized_docs = [tokenize_for_bm25(doc) for doc in bm25_data["documents"]]
    tokenized_query = tokenize_for_bm25(query)
    
    # Handle empty tokenization
    if not tokenized_query:
        tokenized_query = query.lower().split()
    
    bm25 = BM25Okapi(tokenized_docs)
    scores = bm25.get_scores(tokenized_query)
    
    top_indices = np.argsort(scores)[::-1][:top_k]
    
    results = []
    for idx in top_indices:
        if scores[idx] > 0:
            results.append({
                "text": bm25_data["documents"][idx],
                "metadata": bm25_data["metadata"][idx],
                "score": float(scores[idx])
            })
    
    return results

# ==================== Document Storage ====================

def save_full_document(doc_id: str, content: str, metadata: dict, images: List[dict] = None):
    doc_path = os.path.join(DOCUMENTS_FOLDER, f"{doc_id}.json")
    doc_data = {
        "id": doc_id,
        "content": content,
        "metadata": metadata,
        "images": images or [],
        "created_at": datetime.now().isoformat()
    }
    with open(doc_path, 'w', encoding='utf-8') as f:
        json.dump(doc_data, f, ensure_ascii=False, indent=2)
    return doc_id

def load_full_document(doc_id: str) -> Optional[dict]:
    doc_path = os.path.join(DOCUMENTS_FOLDER, f"{doc_id}.json")
    if os.path.exists(doc_path):
        with open(doc_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    return None

def get_full_documents_by_ids(doc_ids: List[str]) -> List[dict]:
    """Load full documents by IDs with content deduplication"""
    documents = []
    seen_ids = set()
    seen_content_hashes = set()  # Deduplicate by content
    
    for doc_id in doc_ids:
        if doc_id not in seen_ids:
            doc = load_full_document(doc_id)
            if doc:
                # Create content hash for deduplication
                content = doc.get('content', '')
                content_hash = hash(content[:1000])  # Hash first 1000 chars
                
                if content_hash not in seen_content_hashes:
                    documents.append(doc)
                    seen_content_hashes.add(content_hash)
                    print(f"Ã¢Å“â€¦ Loaded unique document: {doc.get('metadata', {}).get('title', doc_id)[:50]}")
                else:
                    print(f"Ã¢Â­ï¸ Skipping duplicate content: {doc.get('metadata', {}).get('title', doc_id)[:50]}")
                
                seen_ids.add(doc_id)
    
    print(f"ðŸ“Å¡ Unique documents after dedup: {len(documents)}")
    return documents

# ==================== Conversation Functions ====================

def get_conversation_path(conversation_id: str) -> str:
    return os.path.join(CONVERSATIONS_FOLDER, f"{conversation_id}.json")

def load_conversation(conversation_id: str) -> Optional[dict]:
    path = get_conversation_path(conversation_id)
    if os.path.exists(path):
        with open(path, 'r') as f:
            return json.load(f)
    return None

def save_conversation(conversation: dict):
    path = get_conversation_path(conversation['id'])
    with open(path, 'w') as f:
        json.dump(conversation, f, indent=2)

def create_conversation(first_message: str) -> dict:
    conversation_id = uuid.uuid4().hex[:12]
    title = first_message[:50] + "..." if len(first_message) > 50 else first_message
    
    conversation = {
        "id": conversation_id,
        "title": title,
        "messages": [],
        "created_at": datetime.now().isoformat(),
        "updated_at": datetime.now().isoformat()
    }
    save_conversation(conversation)
    return conversation

def add_message_to_conversation(conversation_id: str, role: str, content: str, sources: list = []):
    conversation = load_conversation(conversation_id)
    if conversation:
        conversation['messages'].append({
            "role": role,
            "content": content,
            "timestamp": datetime.now().isoformat(),
            "sources": sources
        })
        conversation['updated_at'] = datetime.now().isoformat()
        save_conversation(conversation)

def list_all_conversations() -> List[dict]:
    conversations = []
    for filename in os.listdir(CONVERSATIONS_FOLDER):
        if filename.endswith('.json'):
            path = os.path.join(CONVERSATIONS_FOLDER, filename)
            try:
                with open(path, 'r') as f:
                    content = f.read()
                    if not content.strip():  # Empty file
                        print(f"âš ï¸ Skipping empty conversation file: {filename}")
                        continue
                    conv = json.loads(content)
                    conversations.append({
                        "id": conv['id'],
                        "title": conv['title'],
                        "created_at": conv['created_at'],
                        "updated_at": conv['updated_at'],
                        "message_count": len(conv['messages'])
                    })
            except (json.JSONDecodeError, KeyError) as e:
                print(f"âš ï¸ Skipping corrupted conversation file {filename}: {e}")
                continue
    conversations.sort(key=lambda x: x['updated_at'], reverse=True)
    return conversations

# ==================== Retrieval Functions ====================

def rerank_results(query: str, documents: List[dict], top_n: int = 5, source: str = "") -> List[dict]:
    if not documents:
        return []
    
    try:
        # First, deduplicate by content
        seen_content = set()
        unique_docs = []
        for doc in documents:
            text = doc.get('text', doc.get('metadata', {}).get('text', ''))
            if text:
                content_hash = hash(text[:500])  # Hash first 500 chars
                if content_hash not in seen_content:
                    unique_docs.append(doc)
                    seen_content.add(content_hash)
        
        if len(unique_docs) < len(documents):
            print(f"ðŸ”â€ž Deduplicated: {len(documents)} Ã¢â€ â€™ {len(unique_docs)} documents before reranking")
        
        documents = unique_docs
        
        texts = [doc.get('text', doc.get('metadata', {}).get('text', '')) for doc in documents]
        texts = [t for t in texts if t]
        
        if not texts:
            return documents[:top_n]
        
        response = co.rerank(
            model="rerank-v3.5",
            query=query,
            documents=texts,
            top_n=min(top_n, len(texts))
        )
        
        # Log reranking cost
        rerank_calc = CostTracker.calculate_rerank_cost(1)
        CostTracker.log_cost(
            service="cohere",
            model="rerank-v3.5",
            operation="reranking",
            quantity=1,
            unit="searches",
            unit_price=API_PRICING["cohere"]["rerank-v3.5"]["price_per_search"],
            total_cost=rerank_calc["cost"],
            source=source or query[:50],
            documents_count=len(texts)
        )
        
        reranked = []
        for result in response.results:
            original_doc = documents[result.index]
            reranked.append({
                'text': texts[result.index],
                'metadata': original_doc.get('metadata', original_doc),
                'score': result.relevance_score
            })
        
        return reranked
    
    except Exception as e:
        print(f"Reranking error: {e}")
        return documents[:top_n]

def rewrite_query_with_llm(query: str) -> List[str]:
    """
    Use LLM to understand the query and generate search variations.
    This helps match user's natural language to exact product names in the database.
    Now with better category awareness (compute vs database vs storage).
    """
    try:
        prompt = QUERY_REWRITE_PROMPT.format(query=query)

        response = llm.invoke(prompt)
        response_text = response.content.strip()
        
        # Clean up response
        if response_text.startswith("```"):
            response_text = response_text.split("```")[1]
            if response_text.startswith("json"):
                response_text = response_text[4:]
        
        variations = json.loads(response_text.strip())
        
        # Always include original query
        if query not in variations:
            variations.insert(0, query)
        
        print(f"ðŸ§  LLM Query Rewrite: '{query}' Ã¢â€ â€™ {variations}")
        return variations[:5]  # Max 5 variations
        
    except Exception as e:
        print(f"âš ï¸ LLM query rewrite failed: {e}, using original query")
        return [query]


def normalize_query(query: str) -> str:
    """
    Simple query normalization - generic version.
    Just returns the query with basic cleanup.
    """
    # Basic cleanup - remove extra whitespace
    return ' '.join(query.split())


def multi_query_search(query: str, top_k: int = 10, source_filter: dict = None) -> List[dict]:
    """
    Search using multiple query variations generated by LLM.
    Combines results from all variations for better recall.
    """
    # Get query variations from LLM
    query_variations = rewrite_query_with_llm(query)
    
    all_results = {}
    
    for variation in query_variations:
        # Embed each variation
        query_vector = embeddings.embed_query(variation)
        
        # Search Pinecone
        if source_filter:
            results = index.query(
                vector=query_vector,
                top_k=top_k,
                include_metadata=True,
                filter=source_filter
            )
        else:
            results = index.query(
                vector=query_vector,
                top_k=top_k,
                include_metadata=True
            )
        
        # Combine results (keep highest score for each unique result)
        for match in results['matches']:
            result_id = match['id']
            if result_id not in all_results or match['score'] > all_results[result_id]['score']:
                all_results[result_id] = {
                    'id': result_id,
                    'score': match['score'],
                    'metadata': match['metadata'],
                    'matched_variation': variation
                }
    
    # Sort by score and return top results
    sorted_results = sorted(all_results.values(), key=lambda x: x['score'], reverse=True)
    
    print(f"ðŸ”Â Multi-query search: {len(query_variations)} variations Ã¢â€ â€™ {len(sorted_results)} unique results")
    
    return sorted_results[:top_k]


def hybrid_search(query: str, top_k: int = 10, source_filter: dict = None, tracker: PipelineTracker = None) -> List[dict]:
    """Search in ALL sources (documents, websites, youtube, memory) with LLM query understanding"""
    
    # Use LLM-powered multi-query search for better recall
    multi_results = multi_query_search(query, top_k=top_k * 2, source_filter=source_filter)
    
    # Calculate and log query embedding cost (approximate for multiple queries)
    embed_calc = CostTracker.calculate_embedding_cost(query)
    CostTracker.log_cost(
        service="openai",
        model="text-embedding-3-large",
        operation="query_embedding",
        quantity=embed_calc["tokens"] * 3,  # Approximate for multiple variations
        unit="tokens",
        unit_price=API_PRICING["openai"]["text-embedding-3-large"]["price_per_million_tokens"],
        total_cost=embed_calc["cost"] * 3,
        source=query[:50],
        pipeline_id=tracker.id if tracker else ""
    )
    
    if tracker:
        tracker.complete_step("pinecone_search", top_k=top_k)
        tracker.start_step("bm25_search")
    
    # BM25 keyword search (use original query for keyword matching)
    bm25_results = search_bm25(query, top_k=top_k)
    
    if tracker:
        tracker.complete_step("bm25_search", results_count=len(bm25_results))
    
    # Combine multi-query results
    combined = {}
    
    for result in multi_results:
        metadata = result['metadata']
        doc_type = metadata.get('type', '')
        score = result['score']
        
        # Handle memory type differently (no parent_id)
        if doc_type == 'memory':
            memory_id = result['id']
            memory_text = metadata.get('text', '')
            memory_answer = metadata.get('answer', '')
            
            # Skip negative/failed memories
            negative_patterns = [
                "i don't have specific information",
                "i don't have information",
                "unfortunately, i don't have",
                "i couldn't find",
                "i don't see any",
                "no relevant information"
            ]
            
            is_negative = any(p in memory_answer.lower() or p in memory_text.lower() for p in negative_patterns)
            if is_negative:
                print(f"Ã¢Â­ï¸ Skipping negative memory result: {memory_id}")
                continue  # Skip this memory result
            
            combined[memory_id] = {
                'text': memory_text,
                'metadata': metadata,
                'score': score,
                'source': 'semantic',
                'is_memory': True
            }
        else:
            parent_id = metadata.get('parent_id')
            if parent_id:
                if parent_id not in combined or score > combined[parent_id]['score']:
                    combined[parent_id] = {
                        'text': metadata.get('content', '') or metadata.get('text', ''),
                        'metadata': metadata,
                        'score': score,
                        'source': 'semantic',
                        'is_memory': False
                    }
    
    # Add BM25 results
    for result in bm25_results:
        parent_id = result['metadata'].get('parent_id')
        if parent_id:
            bm25_score = min(result['score'] / 10, 1.0)
            
            if parent_id in combined:
                combined[parent_id]['score'] = min(combined[parent_id]['score'] + 0.1, 1.0)
                combined[parent_id]['source'] = 'hybrid'
            else:
                combined[parent_id] = {
                    'text': result['text'],
                    'metadata': result['metadata'],
                    'score': bm25_score,
                    'source': 'bm25',
                    'is_memory': False
                }
    
    sorted_results = sorted(combined.values(), key=lambda x: x['score'], reverse=True)
    return sorted_results[:top_k]

def get_document_context_with_sources(query: str, top_k: int = 10, display_min_score: float = 0.30, track_pipeline: bool = True):
    """Get context from ALL sources including memory and images"""
    
    # Start pipeline tracking
    tracker = None
    if track_pipeline:
        tracker = PipelineTracker("retrieval", query[:50])
        active_pipelines[tracker.id] = tracker
    
    try:
        # Step 1: Embed Query
        if tracker:
            tracker.start_step("embed_query")
        
        # Ã˜Â¯Ã™Ë†Ã˜Â± Ã™ÂÃ™Å  Ã™Æ’Ã™â€ž Ã˜Â§Ã™â€žÃ™â€¦Ã˜ÂµÃ˜Â§Ã˜Â¯Ã˜Â± (Ã˜Â¨Ã˜Â¯Ã™Ë†Ã™â€  filter)
        # Step 2 & 3: Pinecone + BM25 search (inside hybrid_search)
        if tracker:
            tracker.complete_step("embed_query", model="text-embedding-3-large")
            tracker.start_step("pinecone_search")
        
        search_results = hybrid_search(query, top_k=top_k, source_filter=None, tracker=tracker)
        
        if tracker:
            tracker.start_step("rrf_merge")
            tracker.complete_step("rrf_merge", unique_results=len(search_results))
        
        # Step 5: Reranking
        if tracker:
            tracker.start_step("rerank")
        reranked_docs = rerank_results(query, search_results, top_n=8)
        if tracker:
            tracker.complete_step("rerank", model="rerank-v3.5", input_docs=len(search_results), output_docs=len(reranked_docs))
        
        print(f"ðŸ”Â Search results count: {len(search_results)}")
        print(f"ðŸ”Â Reranked docs count: {len(reranked_docs)}")
        
        parent_ids = []
        sources_info = []
        memory_contexts = []
        chunk_texts = {}  # Store chunk texts as fallback
        all_images = []  # Store images from relevant documents
        
        for doc in reranked_docs:
            metadata = doc.get('metadata', {})
            score = doc.get('score', 0)
            doc_type = metadata.get('type', 'document')
            chunk_text = doc.get('text', '') or metadata.get('text', '')  # Try both places
            
            print(f"ðŸ“„ Doc type: {doc_type}, Score: {score}, Has text: {bool(chunk_text)}, Text length: {len(chunk_text) if chunk_text else 0}")
            
            # Handle Memory
            if doc_type == 'memory':
                if score >= display_min_score:
                    question = metadata.get('question', '')
                    answer = metadata.get('answer', '')
                    memory_contexts.append({
                        'question': question,
                        'answer': answer,
                        'score': score
                    })
                    sources_info.append({
                        "type": "memory",
                        "title": f"Previous Q&A: {question[:50]}...",
                        "source": "Memory",
                        "score": round(score, 2),
                        "download_url": None
                    })
            else:
                # Handle Documents/Websites/YouTube
                parent_id = metadata.get('parent_id')
                
                if parent_id and parent_id not in parent_ids:
                    parent_ids.append(parent_id)
                    
                    # Store chunk text as fallback
                    if chunk_text:
                        if parent_id not in chunk_texts:
                            chunk_texts[parent_id] = {
                                'title': metadata.get('title', metadata.get('filename', 'Document')),
                                'texts': []
                            }
                        chunk_texts[parent_id]['texts'].append(chunk_text)
                    
                    if score >= display_min_score:
                        source_url = metadata.get('source', 'Unknown')
                        source_title = metadata.get('title', metadata.get('filename', source_url))
                        domain = metadata.get('domain', '')
                        file_path = metadata.get('file_path', '')
                        
                        # Add domain to title for clarity
                        if domain and domain.lower() not in source_title.lower():
                            source_title = f"{domain} - {source_title}"
                        
                        download_url = None
                        if file_path and os.path.exists(file_path):
                            filename = os.path.basename(file_path)
                            download_url = f"/download/{filename}"
                        
                        sources_info.append({
                            "type": doc_type,
                            "title": source_title,
                            "source": source_url,
                            "score": round(score, 2),
                            "download_url": download_url
                        })
                elif parent_id and chunk_text:
                    # Add more chunks for existing parent
                    if parent_id in chunk_texts:
                        chunk_texts[parent_id]['texts'].append(chunk_text)
        
        print(f"ðŸ“Å¡ Parent IDs found: {len(parent_ids)}")
        print(f"ðŸ“Å¡ Chunk texts collected: {len(chunk_texts)}")
        
        # Build context from full documents OR chunk texts
        full_documents = get_full_documents_by_ids(parent_ids)
        
        print(f"ðŸ“Å¡ Full documents loaded: {len(full_documents)}")
        
        # TOKEN LIMIT to prevent 429 errors
        MAX_CONTEXT_TOKENS = 15000  # Leave room for system prompt, history, and response
        
        def estimate_tokens(text):
            return max(1, len(text) // 4)
        
        context = ""
        context_tokens = 0
        used_parent_ids = set()
        truncated = False
        
        # PRIORITY 1: Use chunk texts FIRST (they contain accurate extracted data like tables)
        # This is more reliable than full documents which might miss table content
        seen_chunk_content = set()
        for parent_id, chunk_data in chunk_texts.items():
            if context_tokens >= MAX_CONTEXT_TOKENS:
                truncated = True
                break
            
            # Deduplicate by content
            content_preview = "".join(chunk_data['texts'])[:500]
            content_hash = hash(content_preview)
            
            if content_hash in seen_chunk_content:
                print(f"Ã¢Â­ï¸ Skipping duplicate chunk content: {chunk_data['title'][:40]}")
                continue
            
            seen_chunk_content.add(content_hash)
            
            print(f"ðŸ“Â Using chunk texts for: {chunk_data['title'][:50]}")
            chunk_content = f"=== {chunk_data['title']} ===\n"
            
            # Join unique chunks (these contain the actual extracted data)
            unique_texts = list(dict.fromkeys(chunk_data['texts']))  # Preserve order, remove dups
            chunk_content += "\n\n".join(unique_texts) + "\n\n"
            
            chunk_tokens = estimate_tokens(chunk_content)
            if context_tokens + chunk_tokens <= MAX_CONTEXT_TOKENS:
                context += chunk_content
                context_tokens += chunk_tokens
                used_parent_ids.add(parent_id)
            else:
                # Try to add partial
                remaining_tokens = MAX_CONTEXT_TOKENS - context_tokens
                if remaining_tokens > 500:
                    max_chars = remaining_tokens * 4
                    chunk_content = chunk_content[:max_chars] + "\n... [truncated]"
                    context += chunk_content
                    context_tokens = MAX_CONTEXT_TOKENS
                    used_parent_ids.add(parent_id)
                truncated = True
                break
        
        # PRIORITY 2: Supplement with full documents (for additional context)
        # Only add if we have room and the parent wasn't already covered by chunks
        for doc in full_documents:
            if context_tokens >= MAX_CONTEXT_TOKENS:
                break
                
            parent_id = doc.get('id', '')
            
            # Skip if already covered by chunk texts
            if parent_id in used_parent_ids:
                continue
            
            doc_content = doc['content']
            doc_tokens = estimate_tokens(doc_content)
            
            # Check if adding this doc would exceed limit
            if context_tokens + doc_tokens > MAX_CONTEXT_TOKENS:
                remaining_tokens = MAX_CONTEXT_TOKENS - context_tokens
                if remaining_tokens > 500:
                    max_chars = remaining_tokens * 4
                    doc_content = doc_content[:max_chars] + "\n... [truncated]"
                    truncated = True
                else:
                    truncated = True
                    continue
            
            used_parent_ids.add(parent_id)
            context += f"=== {doc['metadata'].get('title', 'Document')} (Full) ===\n"
            context += doc_content + "\n\n"
            context_tokens = estimate_tokens(context)
            
            # Collect images from document
            doc_images = doc.get('images', [])
            if doc_images:
                for img in doc_images:
                    img['doc_title'] = doc['metadata'].get('title', 'Document')
                    all_images.append(img)
            
            # Stop if we've reached the limit
            if context_tokens >= MAX_CONTEXT_TOKENS:
                break
        
        # Add Memory Context (if we have room)
        if memory_contexts and context_tokens < MAX_CONTEXT_TOKENS:
            memory_content = "=== Previous Conversations (Memory) ===\n"
            for mem in memory_contexts:
                memory_content += f"Q: {mem['question']}\n"
                memory_content += f"A: {mem['answer']}\n\n"
            
            memory_tokens = estimate_tokens(memory_content)
            if context_tokens + memory_tokens <= MAX_CONTEXT_TOKENS:
                context += memory_content
                context_tokens += memory_tokens
        
        if truncated:
            print(f"âš ï¸ Context was truncated to fit within {MAX_CONTEXT_TOKENS} token limit")
        
        print(f"ðŸ“Â Final context length: {len(context)} chars (~{context_tokens} tokens)")
        print(f"ðŸ–¼ï¸ Images found: {len(all_images)}")
        
        # Remove duplicate sources
        seen_sources = set()
        unique_sources = []
        for source in sources_info:
            source_key = f"{source['type']}_{source['source']}"
            if source_key not in seen_sources:
                seen_sources.add(source_key)
                unique_sources.append(source)
        
        # Return tracker for chat to complete
        return context, unique_sources[:5], all_images[:10], tracker  # Limit to 10 images
    
    except Exception as e:
        if tracker:
            tracker.error_step("rerank", str(e))
            tracker.finish(error=str(e))
            if tracker.id in active_pipelines:
                del active_pipelines[tracker.id]
        raise

def save_to_memory(question: str, answer: str, user_id: str):
    """Save Q&A to memory - but SKIP negative/failed answers"""
    
    # Don't save answers that indicate failure or lack of information
    negative_patterns = [
        "i don't have specific information",
        "i don't have information",
        "unfortunately, i don't have",
        "i couldn't find",
        "i don't see any",
        "no relevant information",
        "i'm not able to find",
        "i cannot find",
        "not available in",
        "don't have details about"
    ]
    
    answer_lower = answer.lower()
    for pattern in negative_patterns:
        if pattern in answer_lower:
            print(f"Ã¢Â­ï¸ Skipping memory save - negative answer detected: '{pattern[:30]}...'")
            return  # Don't save this to memory
    
    # Also skip very short answers (likely not useful)
    if len(answer) < 50:
        print(f"Ã¢Â­ï¸ Skipping memory save - answer too short ({len(answer)} chars)")
        return
    
    memory_text = f"Question: {question}\nAnswer: {answer}"
    vector = embeddings.embed_query(memory_text)
    memory_id = f"memory_{user_id}_{uuid.uuid4().hex[:8]}"
    index.upsert(vectors=[{
        "id": memory_id,
        "values": vector,
        "metadata": {
            "type": "memory",
            "user_id": user_id,
            "question": question,
            "answer": answer,
            "timestamp": datetime.now().isoformat(),
            "text": memory_text
        }
    }])
    print(f"Ã¢Å“â€¦ Saved to memory: {question[:50]}...")

def extract_metadata(text: str, source_type: str):
    prompt = f"""Analyze the following {source_type} content and extract metadata.

Content:
{text[:3000]}

Respond in JSON format only:
{{
    "title": "A descriptive title (max 10 words)",
    "summary": "Brief summary (2-3 sentences)",
    "keywords": ["keyword1", "keyword2", "keyword3"],
    "language": "detected language"
}}

JSON only:"""

    try:
        response = llm.invoke(prompt)
        json_str = response.content.strip()
        if json_str.startswith("```"):
            json_str = json_str.split("```")[1]
            if json_str.startswith("json"):
                json_str = json_str[4:]
        return json.loads(json_str.strip())
    except:
        return {"title": "Unknown", "summary": "", "keywords": [], "language": "unknown"}

# ==================== Ingestion with Semantic Chunking ====================

def extract_domain_name(url: str) -> str:
    """Extract readable domain name from URL (e.g., docs.oracle.com -> Oracle)"""
    try:
        from urllib.parse import urlparse
        parsed = urlparse(url)
        domain = parsed.netloc.lower()
        
        # Extract main company name
        domain_parts = domain.replace('www.', '').split('.')
        
        # Common patterns
        company_mappings = {
            'oracle': 'Oracle',
            'boomi': 'Boomi',
            'opentext': 'OpenText',
            'microsoft': 'Microsoft',
            'google': 'Google',
            'aws': 'AWS',
            'amazon': 'Amazon',
            'salesforce': 'Salesforce',
            'ibm': 'IBM',
            'sap': 'SAP',
        }
        
        for part in domain_parts:
            for key, value in company_mappings.items():
                if key in part:
                    return value
        
        # Default: capitalize first meaningful part
        for part in domain_parts:
            if part not in ['docs', 'www', 'help', 'support', 'com', 'org', 'net', 'io']:
                return part.capitalize()
        
        return ""
    except:
        return ""

def ingest_document_with_semantic_chunks(full_text: str, source: str, doc_type: str, file_path: str, metadata: dict, images: List[dict] = None):
    # Start pipeline tracking
    tracker = PipelineTracker("ingestion", source)
    active_pipelines[tracker.id] = tracker
    
    try:
        parent_id = f"parent_{uuid.uuid4().hex[:12]}"
        
        # Step 1: Extract (already done before calling this function)
        tracker.start_step("extract")
        domain_name = extract_domain_name(source) if doc_type == "website" else ""
        tracker.complete_step("extract", 
            doc_type=doc_type, 
            text_length=len(full_text),
            images_count=len(images) if images else 0
        )
        
        # Save full document
        save_full_document(parent_id, full_text, {
            "source": source,
            "type": doc_type,
            "file_path": file_path,
            "title": metadata.get("title", source),
            "summary": metadata.get("summary", ""),
            "domain": domain_name
        }, images=images)
        
        # Step 2: Chunking
        tracker.start_step("chunk")
        chunks, chunker_type = smart_chunk_text(full_text)
        tracker.complete_step("chunk", 
            chunks_count=len(chunks), 
            chunker_type=chunker_type,
            avg_chunk_size=round(len(full_text) / len(chunks)) if chunks else 0
        )
        
        chunker_emoji = "ðŸ§ " if chunker_type == "semantic" else "ðŸ“"
        print(f"ðŸ“¦ Created {len(chunks)} chunks using {chunker_emoji} {chunker_type} chunker for {source}")
        
        # Step 3: Embedding
        tracker.start_step("embed")
        metadata_list = []
        total_embed_cost = 0
        total_tokens = 0
        
        for i, chunk in enumerate(chunks):
            enhanced_text = f"{domain_name} - {chunk}" if domain_name else chunk
            
            chunk_metadata = {
                "text": enhanced_text,
                "source": source,
                "type": doc_type,
                "filename": source,
                "file_path": file_path,
                "title": metadata.get("title", source),
                "summary": metadata.get("summary", ""),
                "parent_id": parent_id,
                "chunk_index": i,
                "chunk_type": chunker_type,
                "domain": domain_name
            }
            
            # Calculate embedding cost
            embed_calc = CostTracker.calculate_embedding_cost(chunk)
            total_embed_cost += embed_calc["cost"]
            total_tokens += embed_calc["tokens"]
            
            vector = embeddings.embed_query(chunk)
            
            # Step 4: Store vectors (inside loop for first chunk only to mark timing)
            if i == 0:
                tracker.complete_step("embed", model="text-embedding-3-large", dimensions=3072)
                tracker.start_step("store_vectors")
            
            index.upsert(vectors=[{
                "id": f"{parent_id}_chunk_{i}",
                "values": vector,
                "metadata": chunk_metadata
            }])
            
            metadata_list.append(chunk_metadata)
        
        # Log embedding cost
        CostTracker.log_cost(
            service="openai",
            model="text-embedding-3-large",
            operation="embedding",
            quantity=total_tokens,
            unit="tokens",
            unit_price=API_PRICING["openai"]["text-embedding-3-large"]["price_per_million_tokens"],
            total_cost=total_embed_cost,
            source=source,
            pipeline_id=tracker.id,
            chunks=len(chunks)
        )
        
        tracker.complete_step("store_vectors", vectors_stored=len(chunks), index="pinecone")
        
        # Step 5: Store BM25
        tracker.start_step("store_bm25")
        add_to_bm25_index(chunks, metadata_list)
        tracker.complete_step("store_bm25", entries_added=len(chunks))
        
        # Finish tracking
        tracker.finish(
            total_chunks=len(chunks),
            chunker_type=chunker_type,
            images_count=len(images) if images else 0,
            embedding_cost=total_embed_cost
        )
        
        return len(chunks), chunker_type
        
    except Exception as e:
        tracker.error_step("embed", str(e))
        tracker.finish(error=str(e))
        raise
    finally:
        # Remove from active pipelines
        if tracker.id in active_pipelines:
            del active_pipelines[tracker.id]

def ingest_with_unstructured(chunks_with_metadata: List[dict], full_text: str, source: str, doc_type: str, file_path: str, metadata: dict, chunker_type: str = "unstructured", images: List[dict] = None):
    """
    Ingest document using pre-parsed Unstructured.io chunks.
    chunks_with_metadata: List of {"text": str, "element_type": str, "metadata": dict}
    """
    tracker = PipelineTracker("ingestion", source)
    active_pipelines[tracker.id] = tracker
    
    try:
        parent_id = f"parent_{uuid.uuid4().hex[:12]}"
        
        tracker.start_step("extract")
        domain_name = extract_domain_name(source) if doc_type == "website" else ""
        
        element_type_counts = {}
        for c in chunks_with_metadata:
            et = c.get("element_type", "Unknown")
            element_type_counts[et] = element_type_counts.get(et, 0) + 1
        
        tracker.complete_step("extract", 
            doc_type=doc_type, 
            text_length=len(full_text),
            images_count=len(images) if images else 0,
            element_types=element_type_counts
        )
        
        doc_metadata = {
            "source": source,
            "type": doc_type,
            "file_path": file_path,
            "title": metadata.get("title", source),
            "summary": metadata.get("summary", ""),
            "domain": domain_name,
            "chunker_type": chunker_type,
            "element_type_counts": element_type_counts
        }
        save_full_document(parent_id, full_text, doc_metadata, images=images)
        
        tracker.start_step("chunk")
        tracker.complete_step("chunk", 
            chunks_count=len(chunks_with_metadata), 
            chunker_type=chunker_type,
            element_types=element_type_counts,
            avg_chunk_size=round(len(full_text) / len(chunks_with_metadata)) if chunks_with_metadata else 0
        )
        
        chunker_emoji = "ðŸ”Â§" if chunker_type == "unstructured" else ("ðŸ§ " if chunker_type == "semantic" else "ðŸ“")
        print(f"ðŸ“¦ Created {len(chunks_with_metadata)} chunks using {chunker_emoji} {chunker_type} chunker")
        print(f"ðŸ“Š Element types: {element_type_counts}")
        
        tracker.start_step("embed")
        metadata_list = []
        plain_chunks = []
        total_embed_cost = 0
        total_tokens = 0
        
        for i, chunk_data in enumerate(chunks_with_metadata):
            chunk_text = chunk_data.get("text", "")
            element_type = chunk_data.get("element_type", "Unknown")
            chunk_extra_meta = chunk_data.get("metadata", {})
            
            if not chunk_text.strip():
                continue
            
            enhanced_text = f"{domain_name} - {chunk_text}" if domain_name else chunk_text
            
            chunk_metadata = {
                "text": enhanced_text,
                "source": source,
                "type": doc_type,
                "filename": source,
                "file_path": file_path,
                "title": metadata.get("title", source),
                "summary": metadata.get("summary", ""),
                "parent_id": parent_id,
                "chunk_index": i,
                "chunk_type": chunker_type,
                "element_type": element_type,
                "domain": domain_name,
                "has_table": element_type == "Table" or chunk_extra_meta.get("has_table", False)
            }
            
            if "page_number" in chunk_extra_meta:
                chunk_metadata["page_number"] = chunk_extra_meta["page_number"]
            
            embed_calc = CostTracker.calculate_embedding_cost(chunk_text)
            total_embed_cost += embed_calc["cost"]
            total_tokens += embed_calc["tokens"]
            
            vector = embeddings.embed_query(chunk_text)
            
            if i == 0:
                tracker.complete_step("embed", model="text-embedding-3-large", dimensions=3072)
                tracker.start_step("store_vectors")
            
            index.upsert(vectors=[{
                "id": f"{parent_id}_chunk_{i}",
                "values": vector,
                "metadata": chunk_metadata
            }])
            
            metadata_list.append(chunk_metadata)
            plain_chunks.append(chunk_text)
        
        CostTracker.log_cost(
            service="openai",
            model="text-embedding-3-large",
            operation="embedding",
            quantity=total_tokens,
            unit="tokens",
            unit_price=API_PRICING["openai"]["text-embedding-3-large"]["price_per_million_tokens"],
            total_cost=total_embed_cost,
            source=source,
            pipeline_id=tracker.id,
            chunks=len(plain_chunks)
        )
        
        tracker.complete_step("store_vectors", vectors_stored=len(plain_chunks), index="pinecone")
        
        tracker.start_step("store_bm25")
        add_to_bm25_index(plain_chunks, metadata_list)
        tracker.complete_step("store_bm25", entries_added=len(plain_chunks))
        
        tracker.finish(
            total_chunks=len(plain_chunks),
            chunker_type=chunker_type,
            element_types=element_type_counts,
            images_count=len(images) if images else 0,
            embedding_cost=total_embed_cost
        )
        
        return len(plain_chunks), chunker_type, element_type_counts
        
    except Exception as e:
        tracker.error_step("embed", str(e))
        tracker.finish(error=str(e))
        raise
    finally:
        if tracker.id in active_pipelines:
            del active_pipelines[tracker.id]

# ==================== API Endpoints ====================

@app.get("/")
def root():
    return {"message": "ðŸ¤– RAG Agent API v2.4 - With Progress Tracking"}

# ==================== Progress Streaming Helpers ====================

def progress_event(step: str, percent: int, message: str, **extra):
    """Create a progress event for SSE"""
    data = {"type": "progress", "step": step, "percent": percent, "message": message, **extra}
    return f"data: {json.dumps(data)}\n\n"

def done_event(message: str, chunks: int = 0, **extra):
    """Create a done event for SSE"""
    data = {"type": "done", "percent": 100, "message": message, "chunks": chunks, **extra}
    return f"data: {json.dumps(data)}\n\n"

def error_event(message: str):
    """Create an error event for SSE"""
    data = {"type": "error", "message": message}
    return f"data: {json.dumps(data)}\n\n"

def ingest_with_progress(full_text: str, source: str, doc_type: str, file_path: str, metadata: dict, images: List[dict] = None):
    """Generator that yields progress updates during ingestion - with BATCH embedding"""
    
    parent_id = f"parent_{uuid.uuid4().hex[:12]}"
    domain_name = extract_domain_name(source) if doc_type == "website" else ""
    
    yield progress_event("saving_doc", 40, "ðŸ’¾ Saving document...")
    
    save_full_document(parent_id, full_text, {
        "source": source,
        "type": doc_type,
        "file_path": file_path,
        "title": metadata.get("title", source),
        "summary": metadata.get("summary", ""),
        "domain": domain_name
    }, images=images)
    
    yield progress_event("chunking", 45, "Ã¢Å“â€šï¸ Creating semantic chunks...")
    
    chunks, chunker_type = smart_chunk_text(full_text)
    total_chunks = len(chunks)
    
    if not chunks:
        yield error_event("No valid chunks to embed")
        return
    
    # Show chunker type in progress
    chunker_display = "ðŸ§  Semantic Chunker" if chunker_type == "semantic" else "ðŸ“ Fallback Splitter"
    img_count = len(images) if images else 0
    yield progress_event("chunking", 50, f"Created {total_chunks} chunks using {chunker_display}", chunker_type=chunker_type)
    
    # Prepare metadata
    metadata_list = []
    for i, chunk in enumerate(chunks):
        enhanced_text = f"{domain_name} - {chunk}" if domain_name else chunk
        
        chunk_metadata = {
            "text": enhanced_text,
            "source": source,
            "type": doc_type,
            "filename": source,
            "file_path": file_path,
            "title": metadata.get("title", source),
            "summary": metadata.get("summary", ""),
            "parent_id": parent_id,
            "chunk_index": i,
            "chunk_type": chunker_type,
            "domain": domain_name
        }
        metadata_list.append(chunk_metadata)
    
    # BATCH EMBEDDING - much faster!
    yield progress_event("embedding", 55, f"ðŸ§  Batch embedding {total_chunks} chunks (this is fast!)...")
    
    try:
        print(f"ðŸš€ Starting batch embedding for {total_chunks} chunks...")
        vectors = embeddings.embed_documents(chunks)
        print(f"Ã¢Å“â€¦ Batch embedding complete: {len(vectors)} vectors")
        
        yield progress_event("embedding_done", 75, f"Ã¢Å“â€¦ Embedded {len(vectors)} chunks in one batch!")
        
    except Exception as e:
        print(f"Ã¢ÂÅ’ Batch embedding failed: {e}")
        yield error_event(f"Embedding failed: {str(e)}")
        return
    
    # Store vectors in Pinecone (batch upsert)
    yield progress_event("storing", 80, f"ðŸ’¾ Storing {len(vectors)} vectors in Pinecone...")
    
    batch_size = 100
    vectors_to_upsert = []
    
    for i, (vector, meta) in enumerate(zip(vectors, metadata_list)):
        vectors_to_upsert.append({
            "id": f"{parent_id}_chunk_{i}",
            "values": vector,
            "metadata": meta
        })
        
        if len(vectors_to_upsert) >= batch_size:
            index.upsert(vectors=vectors_to_upsert)
            vectors_to_upsert = []
            progress = 80 + int((i / len(vectors)) * 15)
            yield progress_event("storing", progress, f"ðŸ’¾ Stored {i+1}/{len(vectors)} vectors...")
    
    if vectors_to_upsert:
        index.upsert(vectors=vectors_to_upsert)
    
    yield progress_event("bm25", 97, "ðŸ“Å¡ Adding to BM25 index...")
    add_to_bm25_index(chunks, metadata_list)
    
    # Include chunker type in done message
    chunker_emoji = "ðŸ§ " if chunker_type == "semantic" else "ðŸ“"
    yield done_event(f"Ã¢Å“â€¦ Ingested with {total_chunks} chunks {chunker_emoji}", total_chunks, chunker_type=chunker_type, images_count=img_count)

def ingest_with_progress_unstructured(chunks_with_metadata: List[dict], full_text: str, source: str, doc_type: str, file_path: str, metadata: dict, chunker_type: str = "unstructured", images: List[dict] = None):
    """Generator that yields progress updates during Unstructured ingestion - with BATCH embedding"""
    
    parent_id = f"parent_{uuid.uuid4().hex[:12]}"
    domain_name = extract_domain_name(source) if doc_type == "website" else ""
    
    element_type_counts = {}
    for c in chunks_with_metadata:
        et = c.get("element_type", "Unknown")
        element_type_counts[et] = element_type_counts.get(et, 0) + 1
    
    yield progress_event("saving_doc", 40, "ðŸ’¾ Saving document...")
    
    save_full_document(parent_id, full_text, {
        "source": source,
        "type": doc_type,
        "file_path": file_path,
        "title": metadata.get("title", source),
        "summary": metadata.get("summary", ""),
        "domain": domain_name,
        "chunker_type": chunker_type,
        "element_type_counts": element_type_counts
    }, images=images)
    
    total_chunks = len(chunks_with_metadata)
    element_summary = ", ".join([f"{v} {k}" for k, v in element_type_counts.items()])
    img_count = len(images) if images else 0
    yield progress_event("chunking", 45, f"ðŸ”Â§ Unstructured created {total_chunks} chunks ({element_summary})", 
                        chunker_type=chunker_type, element_types=element_type_counts)
    
    # Prepare all chunks first
    yield progress_event("preparing", 50, f"ðŸ“Â Preparing {total_chunks} chunks for embedding...")
    
    metadata_list = []
    plain_chunks = []
    chunk_indices = []
    
    for i, chunk_data in enumerate(chunks_with_metadata):
        chunk_text = chunk_data.get("text", "")
        element_type = chunk_data.get("element_type", "Unknown")
        
        if not chunk_text.strip():
            continue
        
        enhanced_text = f"{domain_name} - {chunk_text}" if domain_name else chunk_text
        
        chunk_metadata = {
            "content": enhanced_text,
            "text": enhanced_text,  # Keep both for backward compatibility
            "source": source,
            "type": doc_type,
            "filename": source,
            "file_path": file_path,
            "title": metadata.get("title", source),
            "summary": metadata.get("summary", ""),
            "parent_id": parent_id,
            "chunk_index": i,
            "chunk_type": chunker_type,
            "element_type": element_type,
            "domain": domain_name,
            "has_table": element_type == "Table"
        }
        
        metadata_list.append(chunk_metadata)
        plain_chunks.append(chunk_text)
        chunk_indices.append(i)
    
    if not plain_chunks:
        yield error_event("No valid chunks to embed")
        return
    
    # BATCH EMBEDDING - much faster!
    yield progress_event("embedding", 55, f"ðŸ§  Batch embedding {len(plain_chunks)} chunks (this is fast!)...")
    
    try:
        # Use embed_documents for batch embedding (single API call!)
        print(f"ðŸš€ Starting batch embedding for {len(plain_chunks)} chunks...")
        vectors = embeddings.embed_documents(plain_chunks)
        print(f"Ã¢Å“â€¦ Batch embedding complete: {len(vectors)} vectors")
        
        yield progress_event("embedding_done", 75, f"Ã¢Å“â€¦ Embedded {len(vectors)} chunks in one batch!")
        
    except Exception as e:
        print(f"Ã¢ÂÅ’ Batch embedding failed: {e}")
        yield error_event(f"Embedding failed: {str(e)}")
        return
    
    # Store vectors in Pinecone
    yield progress_event("storing", 80, f"ðŸ’¾ Storing {len(vectors)} vectors in Pinecone...")
    
    # Batch upsert to Pinecone (in batches of 100)
    batch_size = 100
    vectors_to_upsert = []
    
    for i, (vector, meta, chunk_idx) in enumerate(zip(vectors, metadata_list, chunk_indices)):
        vectors_to_upsert.append({
            "id": f"{parent_id}_chunk_{chunk_idx}",
            "values": vector,
            "metadata": meta
        })
        
        # Upsert in batches
        if len(vectors_to_upsert) >= batch_size:
            index.upsert(vectors=vectors_to_upsert)
            vectors_to_upsert = []
            progress = 80 + int((i / len(vectors)) * 15)
            yield progress_event("storing", progress, f"ðŸ’¾ Stored {i+1}/{len(vectors)} vectors...")
    
    # Upsert remaining vectors
    if vectors_to_upsert:
        index.upsert(vectors=vectors_to_upsert)
    
    yield progress_event("bm25", 97, "ðŸ“Å¡ Adding to BM25 index...")
    add_to_bm25_index(plain_chunks, metadata_list)
    
    print(f"ðŸ“¦ Created {len(plain_chunks)} chunks using ðŸ”Â§ {chunker_type} chunker")
    print(f"ðŸ“Š Element types: {element_type_counts}")
    
    yield done_event(
        f"Ã¢Å“â€¦ Ingested with {len(plain_chunks)} chunks ðŸ”Â§ ({element_summary})", 
        len(plain_chunks), 
        chunker_type=chunker_type, 
        element_types=element_type_counts,
        images_count=img_count
    )

# ==================== Streaming Scrape Endpoint ====================

@app.get("/scrape/stream")
async def scrape_website_stream(url: str):
    """Scrape a website with progress streaming"""
    
    async def generate():
        try:
            # Check if we should use Playwright for this URL
            use_playwright = PLAYWRIGHT_AVAILABLE and should_use_playwright(url)
            
            if use_playwright:
                # Detailed Playwright progress
                yield progress_event("playwright_start", 5, f"ðŸŽ­ Detected dynamic site, using Playwright...")
                yield progress_event("playwright_browser", 8, f"ðŸŽ­ [1/7] Starting headless browser...")
                yield progress_event("playwright_navigate", 12, f"ðŸŽ­ [2/7] Navigating to {url[:40]}...")
                
                html_content = await fetch_page_with_playwright(url)
                html_size_kb = len(html_content) / 1024
                
                yield progress_event("playwright_done", 18, f"ðŸŽ­ âœ“ Rendered page ({html_size_kb:.1f} KB)")
            else:
                yield progress_event("fetching", 10, f"ðŸ“„ Fetching {url}...")
                headers = {'User-Agent': 'Mozilla/5.0'}
                response = requests.get(url, headers=headers, timeout=15)
                response.raise_for_status()
                html_content = response.text
                html_size_kb = len(html_content) / 1024
                yield progress_event("fetched", 18, f"ðŸ“„ âœ“ Fetched ({html_size_kb:.1f} KB)")
            
            yield progress_event("parsing", 20, "ðŸ”Â Analyzing HTML structure...")
            
            # Extract images BEFORE parsing
            images = extract_images_from_page(url, html_content)
            yield progress_event("images", 25, f"ðŸ–¼ï¸ Found {len(images)} images")
            
            # Use Unstructured.io if available
            if UNSTRUCTURED_AVAILABLE:
                yield progress_event("unstructured", 28, "ðŸ”Â§ Starting Unstructured.io smart parsing...")
                yield progress_event("unstructured_parse", 32, "ðŸ”Â§ Detecting tables, text blocks, headers...")
                
                full_text, chunks_with_metadata, chunker_type = parse_html_with_unstructured(html_content, url)
                
                if not full_text or not chunks_with_metadata:
                    yield error_event("No content found on page")
                    return
                
                # Count element types
                element_types = {}
                for c in chunks_with_metadata:
                    et = c.get("element_type", "Unknown")
                    element_types[et] = element_types.get(et, 0) + 1
                element_summary = ", ".join([f"{v} {k}" for k, v in element_types.items()])
                
                yield progress_event("unstructured_done", 38, f"ðŸ”Â§ âœ“ Found {len(chunks_with_metadata)} chunks ({element_summary})")
                
                yield progress_event("metadata", 40, "ðŸ“‹ Extracting metadata...")
                metadata = extract_metadata(full_text[:3000], "website")
                
                yield progress_event("embedding_start", 42, f"ðŸ§  Starting embedding of {len(chunks_with_metadata)} chunks...")
                
                for event in ingest_with_progress_unstructured(chunks_with_metadata, full_text, url, "website", "", metadata, chunker_type, images=images):
                    yield event
            else:
                # Fallback to BeautifulSoup
                yield progress_event("parsing", 30, "ðŸ“Â Parsing with BeautifulSoup...")
                soup = BeautifulSoup(html_content, 'html.parser')
                for element in soup(['script', 'style', 'nav', 'footer', 'header']):
                    element.decompose()
                
                text = soup.get_text(separator='\n', strip=True)
                lines = [line.strip() for line in text.splitlines() if line.strip()]
                full_text = '\n'.join(lines)
                
                if not full_text:
                    yield error_event("No content found on page")
                    return
                
                yield progress_event("metadata", 35, "ðŸ“‹ Extracting metadata...")
                metadata = extract_metadata(full_text, "website")
                
                for event in ingest_with_progress(full_text, url, "website", "", metadata, images=images):
                    yield event
                
        except Exception as e:
            yield error_event(str(e))
    
    return StreamingResponse(generate(), media_type="text/event-stream")

# ==================== Streaming Recursive Scrape ====================

@app.get("/scrape-recursive/stream")
async def scrape_recursive_stream(url: str, max_depth: int = 2, max_pages: int = 20, same_path_only: bool = True):
    """Recursive scrape with progress streaming"""
    
    async def generate():
        try:
            visited = set()
            to_visit = [(url, 0)]  # (url, depth)
            pages_scraped = []
            total_chunks = 0
            total_images = 0
            used_playwright = False
            
            # Check if this is a dynamic site
            is_dynamic_site = PLAYWRIGHT_AVAILABLE and should_use_playwright(url)
            
            if is_dynamic_site:
                yield progress_event("starting", 2, f"ðŸŽ­ Detected dynamic site - will use Playwright for JavaScript rendering")
                yield progress_event("starting", 5, f"ðŸš€ Starting recursive scrape from {url}")
            else:
                yield progress_event("starting", 5, f"ðŸš€ Starting recursive scrape from {url}")
            
            while to_visit and len(visited) < max_pages:
                current_url, depth = to_visit.pop(0)
                
                if current_url in visited:
                    continue
                    
                visited.add(current_url)
                page_num = len(visited)
                
                # Calculate overall progress based on pages
                overall_progress = min(5 + int((page_num / max_pages) * 90), 95)
                
                try:
                    # Use Playwright for dynamic sites, requests for static
                    use_playwright = PLAYWRIGHT_AVAILABLE and should_use_playwright(current_url)
                    
                    if use_playwright:
                        used_playwright = True
                        yield progress_event("playwright", overall_progress, 
                            f"ðŸŽ­ [{page_num}/{max_pages}] Starting browser for: {current_url[:35]}...",
                            current_page=page_num, total_pages=max_pages, current_url=current_url)
                        
                        yield progress_event("playwright_render", overall_progress, 
                            f"ðŸŽ­ [{page_num}/{max_pages}] Rendering JavaScript (this may take 20-40 seconds)...",
                            current_page=page_num, total_pages=max_pages)
                        
                        html_content = await fetch_page_with_playwright(current_url)
                        html_size_kb = len(html_content) / 1024
                        
                        yield progress_event("playwright_done", overall_progress, 
                            f"ðŸŽ­ [{page_num}/{max_pages}] âœ“ Rendered {html_size_kb:.0f}KB",
                            current_page=page_num, total_pages=max_pages)
                    else:
                        yield progress_event("scraping", overall_progress, 
                            f"ðŸ“„ [{page_num}/{max_pages}] Fetching: {current_url[:45]}...",
                            current_page=page_num, total_pages=max_pages, current_url=current_url)
                        
                        headers = {'User-Agent': 'Mozilla/5.0'}
                        response = requests.get(current_url, headers=headers, timeout=10)
                        response.raise_for_status()
                        html_content = response.text
                    
                    # Extract links if we haven't reached max depth
                    if depth < max_depth:
                        new_links = extract_links_from_page(current_url, html_content, same_path_only)
                        if new_links:
                            yield progress_event("links", overall_progress, 
                                f"ðŸ”â€” [{page_num}/{max_pages}] Found {len(new_links)} links to explore")
                        for link in new_links:
                            if link not in visited and len(to_visit) + len(visited) < max_pages * 2:
                                to_visit.append((link, depth + 1))
                    
                    # Extract images BEFORE removing elements
                    images = extract_images_from_page(current_url, html_content)
                    total_images += len(images)
                    
                    # Parse content using Unstructured if available
                    if UNSTRUCTURED_AVAILABLE:
                        yield progress_event("parsing", overall_progress, 
                            f"ðŸ”Â§ [{page_num}/{max_pages}] Parsing with Unstructured.io...")
                        
                        full_text, chunks_with_metadata, chunker_type = parse_html_with_unstructured(html_content, current_url)
                        
                        if full_text and len(full_text) > 100:
                            metadata = extract_metadata(full_text[:3000], "website")
                            chunks = ingest_with_unstructured(chunks_with_metadata, full_text, current_url, "website", "", metadata, chunker_type, images=images)
                            total_chunks += chunks
                            pages_scraped.append(current_url)
                            
                            # Count element types for display
                            element_types = {}
                            for c in chunks_with_metadata:
                                et = c.get("element_type", "Unknown")
                                element_types[et] = element_types.get(et, 0) + 1
                            element_summary = ", ".join([f"{v} {k}" for k, v in element_types.items()])
                            
                            yield progress_event("page_done", overall_progress,
                                f"âœ“ Page {page_num}: {chunks} chunks ðŸ”Â§ ({element_summary}), {len(images)} images",
                                page_url=current_url, page_chunks=chunks, page_images=len(images), chunker_type=chunker_type)
                    else:
                        # Fallback to BeautifulSoup
                        soup = BeautifulSoup(html_content, 'html.parser')
                        for element in soup(['script', 'style', 'nav', 'footer', 'header']):
                            element.decompose()
                        
                        text = soup.get_text(separator='\n', strip=True)
                        lines = [line.strip() for line in text.splitlines() if line.strip()]
                        full_text = '\n'.join(lines)
                        
                        if full_text and len(full_text) > 100:
                            metadata = extract_metadata(full_text[:3000], "website")
                            chunks, chunker_type = ingest_document_with_semantic_chunks(full_text, current_url, "website", "", metadata, images=images)
                            total_chunks += chunks
                            pages_scraped.append(current_url)
                            chunker_emoji = "ðŸ§ " if chunker_type == "semantic" else "ðŸ“"
                            
                            yield progress_event("page_done", overall_progress,
                                f"âœ“ Page {page_num}: {chunks} chunks {chunker_emoji}, {len(images)} images",
                                page_url=current_url, page_chunks=chunks, page_images=len(images), chunker_type=chunker_type)
                    
                except Exception as e:
                    yield progress_event("page_error", overall_progress,
                        f"âœ— Failed: {current_url[:30]}... - {str(e)[:30]}")
                    continue
            
            # Final message with indicators
            unstructured_indicator = " ðŸ”Â§" if UNSTRUCTURED_AVAILABLE else ""
            playwright_indicator = " ðŸŽ­" if used_playwright else ""
            yield done_event(
                f"Ã¢Å“â€¦ Scraped {len(pages_scraped)} pages with {total_chunks} chunks{unstructured_indicator}{playwright_indicator}, {total_images} images",
                total_chunks,
                pages_scraped=pages_scraped,
                total_pages=len(pages_scraped),
                total_images=total_images
            )
            
        except Exception as e:
            yield error_event(str(e))
    
    return StreamingResponse(generate(), media_type="text/event-stream")

# ==================== Streaming YouTube Endpoint ====================

@app.get("/youtube/stream")
async def youtube_stream(url: str):
    """Add YouTube video with progress streaming"""
    
    async def generate():
        try:
            yield progress_event("fetching", 10, "Fetching YouTube transcript...")
            
            from youtube_transcript_api import YouTubeTranscriptApi
            import re
            
            video_id_match = re.search(r"(?:v=|youtu\.be/)([a-zA-Z0-9_-]{11})", url)
            if not video_id_match:
                yield error_event("Invalid YouTube URL")
                return
            
            video_id = video_id_match.group(1)
            
            yield progress_event("transcript", 30, "Downloading transcript...")
            
            try:
                transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
            except Exception as e:
                yield error_event(f"Could not get transcript: {str(e)}")
                return
            
            yield progress_event("processing", 40, "Processing transcript...")
            
            full_text = " ".join([entry['text'] for entry in transcript_list])
            
            if not full_text:
                yield error_event("No transcript content found")
                return
            
            yield progress_event("metadata", 35, "Extracting metadata...")
            metadata = extract_metadata(full_text, "youtube")
            
            # Use the progress generator for ingestion
            for event in ingest_with_progress(full_text, url, "youtube", "", metadata):
                yield event
                
        except Exception as e:
            yield error_event(str(e))
    
    return StreamingResponse(generate(), media_type="text/event-stream")

# ==================== Debug Search Endpoint ====================

@app.get("/debug-search")
def debug_search(q: str):
    """Debug endpoint to see raw search results"""
    
    # 1. Raw Semantic Search from Pinecone
    query_vector = embeddings.embed_query(q)
    semantic_results = index.query(
        vector=query_vector,
        top_k=10,
        include_metadata=True
    )
    
    semantic_docs = []
    for match in semantic_results['matches']:
        semantic_docs.append({
            "id": match['id'],
            "score": round(match['score'], 4),
            "title": match['metadata'].get('title', 'N/A'),
            "type": match['metadata'].get('type', 'N/A'),
            "source": match['metadata'].get('source', 'N/A')[:50],
            "text_preview": match['metadata'].get('text', '')[:100] + "..."
        })
    
    # 2. BM25 Results
    bm25_results = search_bm25(q, top_k=10)
    bm25_docs = []
    for result in bm25_results:
        bm25_docs.append({
            "score": round(result['score'], 4),
            "title": result['metadata'].get('title', 'N/A'),
            "type": result['metadata'].get('type', 'N/A'),
        })
    
    # 3. Hybrid Results
    hybrid_results = hybrid_search(q, top_k=10)
    hybrid_docs = []
    for doc in hybrid_results:
        hybrid_docs.append({
            "score": round(doc['score'], 4),
            "title": doc['metadata'].get('title', 'N/A'),
            "type": doc['metadata'].get('type', 'N/A'),
            "source": doc['metadata'].get('source', 'N/A')[:50],
        })
    
    # 4. After Reranking
    reranked = rerank_results(q, hybrid_results, top_n=5)
    reranked_docs = []
    for doc in reranked:
        reranked_docs.append({
            "score": round(doc['score'], 4),
            "title": doc['metadata'].get('title', 'N/A'),
            "type": doc['metadata'].get('type', 'N/A'),
        })
    
    return {
        "query": q,
        "1_semantic_search": semantic_docs,
        "2_bm25_search": bm25_docs,
        "3_hybrid_combined": hybrid_docs,
        "4_after_reranking": reranked_docs,
        "display_threshold": 0.30
    }

@app.get("/debug-chunks")
def debug_chunks(doc_id: str = None, source_url: str = None, limit: int = 50):
    """Debug endpoint to see stored chunks with their element types and structure"""
    
    results = {
        "filter": {"doc_id": doc_id, "source_url": source_url},
        "chunks": [],
        "element_type_summary": {},
        "chunker_type_summary": {},
        "total_found": 0
    }
    
    try:
        dummy_vector = embeddings.embed_query("test")
        
        filter_dict = {}
        if doc_id:
            filter_dict["parent_id"] = {"$eq": doc_id}
        
        pinecone_results = index.query(
            vector=dummy_vector,
            top_k=limit,
            include_metadata=True,
            filter=filter_dict if filter_dict else None
        )
        
        for match in pinecone_results['matches']:
            meta = match['metadata']
            
            if source_url and source_url not in meta.get('source', ''):
                continue
            
            chunk_info = {
                "id": match['id'],
                "chunk_index": meta.get('chunk_index', 0),
                "element_type": meta.get('element_type', 'Unknown'),
                "chunk_type": meta.get('chunk_type', 'Unknown'),
                "has_table": meta.get('has_table', False),
                "source": meta.get('source', 'N/A')[:80],
                "title": meta.get('title', 'N/A'),
                "text_preview": meta.get('text', '')[:300] + "..." if len(meta.get('text', '')) > 300 else meta.get('text', ''),
                "text_length": len(meta.get('text', '')),
                "page_number": meta.get('page_number', None)
            }
            
            results["chunks"].append(chunk_info)
            
            et = chunk_info["element_type"]
            results["element_type_summary"][et] = results["element_type_summary"].get(et, 0) + 1
            
            ct = chunk_info["chunk_type"]
            results["chunker_type_summary"][ct] = results["chunker_type_summary"].get(ct, 0) + 1
        
        results["chunks"].sort(key=lambda x: x.get('chunk_index', 0))
        results["total_found"] = len(results["chunks"])
        
        return results
        
    except Exception as e:
        return {"error": str(e)}

@app.get("/debug-query")
def debug_query(query: str, top_k: int = 10):
    """
    Debug endpoint to see exactly what chunks are retrieved for a query.
    Shows the raw retrieval results before any processing.
    """
    try:
        # Step 1: Embed the query
        query_vector = embeddings.embed_query(query)
        
        # Step 2: Search Pinecone
        pinecone_results = index.query(
            vector=query_vector,
            top_k=top_k,
            include_metadata=True
        )
        
        # Step 3: Format results
        chunks = []
        for i, match in enumerate(pinecone_results['matches']):
            meta = match['metadata']
            text = meta.get('text', '')
            
            chunks.append({
                "rank": i + 1,
                "score": round(match['score'], 4),
                "id": match['id'],
                "element_type": meta.get('element_type', 'Unknown'),
                "chunk_type": meta.get('chunk_type', 'Unknown'),
                "source": meta.get('source', 'N/A'),
                "title": meta.get('title', 'N/A'),
                "text_length": len(text),
                "text_preview": text[:500] + "..." if len(text) > 500 else text,
                "contains_query_terms": {
                    term: term.lower() in text.lower() 
                    for term in query.split()[:5]  # Check first 5 words
                }
            })
        
        # Step 4: Also do BM25 search for comparison
        bm25_results = []
        try:
            bm25_index = load_bm25_index()
            if bm25_index.get("documents"):
                query_tokens = query.lower().split()
                bm25 = BM25Okapi([doc.get("tokens", []) for doc in bm25_index["documents"]])
                scores = bm25.get_scores(query_tokens)
                top_indices = sorted(range(len(scores)), key=lambda i: scores[i], reverse=True)[:5]
                
                for idx in top_indices:
                    if scores[idx] > 0:
                        doc = bm25_index["documents"][idx]
                        bm25_results.append({
                            "bm25_score": round(scores[idx], 4),
                            "source": doc.get("metadata", {}).get("source", "N/A"),
                            "text_preview": doc.get("text", "")[:200]
                        })
        except:
            pass
        
        return {
            "query": query,
            "embedding_model": "text-embedding-3-small",
            "top_k": top_k,
            "pinecone_results_count": len(chunks),
            "chunks": chunks,
            "bm25_results": bm25_results,
            "analysis": {
                "highest_score": chunks[0]["score"] if chunks else 0,
                "lowest_score": chunks[-1]["score"] if chunks else 0,
                "element_types_found": list(set(c["element_type"] for c in chunks)),
                "sources_found": list(set(c["source"] for c in chunks))[:5]
            }
        }
        
    except Exception as e:
        return {"error": str(e)}

@app.get("/unstructured-status")
def unstructured_status():
    """Check if Unstructured.io and Playwright are available and working"""
    return {
        "unstructured_available": UNSTRUCTURED_AVAILABLE,
        "playwright_available": PLAYWRIGHT_AVAILABLE,
        "message": "Ã¢Å“â€¦ Unstructured.io is loaded and ready" if UNSTRUCTURED_AVAILABLE else "âš ï¸ Unstructured.io is not available, using fallback parsing",
        "capabilities": {
            "html_tables": UNSTRUCTURED_AVAILABLE,
            "pdf_structure": UNSTRUCTURED_AVAILABLE,
            "docx_structure": UNSTRUCTURED_AVAILABLE,
            "xlsx_structure": UNSTRUCTURED_AVAILABLE,
            "pptx_structure": UNSTRUCTURED_AVAILABLE,
            "smart_chunking": UNSTRUCTURED_AVAILABLE,
            "javascript_rendering": PLAYWRIGHT_AVAILABLE
        }
    }

@app.get("/debug-html")
async def debug_html(url: str, search: str = None):
    """
    Debug endpoint to see raw HTML and what Unstructured extracts.
    Use ?search=keyword to highlight/filter content containing that keyword.
    """
    try:
        # Determine if we should use Playwright
        use_playwright = PLAYWRIGHT_AVAILABLE and should_use_playwright(url)
        
        if use_playwright:
            html_content = await fetch_page_with_playwright(url)
            fetch_method = "playwright"
        else:
            response = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'}, timeout=15)
            html_content = response.text
            fetch_method = "requests"
        
        # Parse with Unstructured if available
        unstructured_elements = []
        if UNSTRUCTURED_AVAILABLE:
            full_text, chunks, chunker_type = parse_html_with_unstructured(html_content, url)
            
            for i, chunk in enumerate(chunks):
                chunk_text = chunk.get("text", "")
                element_type = chunk.get("element_type", "Unknown")
                
                # Filter by search term if provided
                if search and search.lower() not in chunk_text.lower():
                    continue
                
                unstructured_elements.append({
                    "index": i,
                    "element_type": element_type,
                    "text_length": len(chunk_text),
                    "text_preview": chunk_text[:500] + "..." if len(chunk_text) > 500 else chunk_text,
                    "contains_search": search.lower() in chunk_text.lower() if search else None
                })
        
        # Search in raw HTML
        html_contains_search = False
        html_search_context = None
        if search:
            html_contains_search = search.lower() in html_content.lower()
            if html_contains_search:
                # Find context around the search term
                idx = html_content.lower().find(search.lower())
                start = max(0, idx - 200)
                end = min(len(html_content), idx + len(search) + 200)
                html_search_context = html_content[start:end]
        
        return {
            "url": url,
            "fetch_method": fetch_method,
            "html_size_bytes": len(html_content),
            "html_size_kb": round(len(html_content) / 1024, 2),
            "search_term": search,
            "html_contains_search": html_contains_search,
            "html_search_context": html_search_context,
            "unstructured_available": UNSTRUCTURED_AVAILABLE,
            "total_elements_extracted": len(unstructured_elements),
            "elements": unstructured_elements[:100],  # Limit to first 100
            "raw_html_preview": html_content[:5000] + "..." if len(html_content) > 5000 else html_content
        }
        
    except Exception as e:
        return {"error": str(e)}

@app.get("/frontend")
def serve_frontend():
    return FileResponse("frontend/index.html")

@app.get("/admin")
def serve_admin():
    return FileResponse("frontend/admin.html")

@app.get("/download/{filename}")
def download_file(filename: str):
    file_path = os.path.join(UPLOADS_FOLDER, filename)
    if os.path.exists(file_path):
        return FileResponse(path=file_path, filename=filename, media_type='application/octet-stream')
    return {"error": "File not found"}

# ==================== Conversation Endpoints ====================

@app.get("/conversations")
def get_conversations():
    return {"conversations": list_all_conversations()}

@app.get("/conversations/{conversation_id}")
def get_conversation(conversation_id: str):
    conversation = load_conversation(conversation_id)
    if conversation:
        return conversation
    return {"error": "Conversation not found"}

@app.delete("/conversations/{conversation_id}")
def delete_conversation(conversation_id: str):
    path = get_conversation_path(conversation_id)
    if os.path.exists(path):
        os.remove(path)
        return {"success": True, "message": "Conversation deleted"}
    return {"success": False, "message": "Conversation not found"}

# ==================== AGENTIC RAG SYSTEM ====================
# The LLM decides what to search, evaluates results, and can search multiple times

# Define tools the agent can use
AGENT_TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "search_knowledge_base",
            "description": "Search the internal knowledge base for information. Use this FIRST for any question. The knowledge base contains curated pricing data, documentation, and specific information the user has added.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "The search query. Be specific and use relevant keywords."
                    },
                    "top_k": {
                        "type": "integer",
                        "description": "Number of results to return (default: 10)",
                        "default": 10
                    }
                },
                "required": ["query"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "search_web",
            "description": "Search the internet for current information. Use this when: 1) Knowledge base doesn't have the specific info needed, 2) You need to fill gaps in a comparison, 3) You need the most current/updated information. This searches the live internet.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "The web search query. Be specific."
                    }
                },
                "required": ["query"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "list_available_sources",
            "description": "List all available sources/documents in the knowledge base. Use this to see what information is available before searching.",
            "parameters": {
                "type": "object",
                "properties": {},
                "required": []
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "get_source_content",
            "description": "Get all content from a specific source. Use this when you need comprehensive information from a particular document or website.",
            "parameters": {
                "type": "object",
                "properties": {
                    "source_name": {
                        "type": "string",
                        "description": "The name or URL of the source to retrieve"
                    }
                },
                "required": ["source_name"]
            }
        }
    }
]

def execute_tool(tool_name: str, arguments: dict) -> dict:
    """Execute a tool and return results"""
    
    if tool_name == "search_knowledge_base":
        query = arguments.get("query", "")
        top_k = arguments.get("top_k", 10)
        
        # Use existing search with reranking
        try:
            query_vector = embeddings.embed_query(query)
            results = index.query(vector=query_vector, top_k=top_k * 3, include_metadata=True)
            
            # Format results - EXCLUDE MEMORY to get actual data
            search_results = []
            seen_content = set()
            
            for match in results.matches:
                # Skip memory type - we want actual documents/websites only
                doc_type = match.metadata.get("type", "")
                if doc_type == "memory":
                    continue
                    
                content = match.metadata.get("content", "") or match.metadata.get("text", "")
                if not content or content[:100] in seen_content:
                    continue
                seen_content.add(content[:100])
                
                search_results.append({
                    "content": content[:3000],  # Limit content size
                    "source": match.metadata.get("source", "Unknown"),
                    "title": match.metadata.get("title", ""),
                    "type": doc_type,
                    "score": float(match.score)
                })
            
            # Rerank with Cohere if available
            if search_results and co:
                try:
                    docs_to_rerank = [r["content"] for r in search_results[:20]]
                    rerank_results = co.rerank(
                        model="rerank-v3.5",
                        query=query,
                        documents=docs_to_rerank,
                        top_n=min(top_k, len(docs_to_rerank))
                    )
                    
                    reranked = []
                    for r in rerank_results.results:
                        result = search_results[r.index]
                        result["relevance_score"] = float(r.relevance_score)
                        reranked.append(result)
                    search_results = reranked
                except Exception as e:
                    print(f"Rerank failed: {e}")
                    search_results = search_results[:top_k]
            
            return {
                "success": True,
                "query": query,
                "results_count": len(search_results),
                "results": search_results[:top_k]
            }
            
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    elif tool_name == "list_available_sources":
        try:
            # Get unique sources from recent vectors
            query_vector = embeddings.embed_query("information data content")
            results = index.query(vector=query_vector, top_k=100, include_metadata=True)
            
            sources = {}
            for match in results.matches:
                source = match.metadata.get("source", "Unknown")
                title = match.metadata.get("title", "")
                source_type = match.metadata.get("type", "")
                
                if source not in sources:
                    sources[source] = {
                        "source": source,
                        "title": title,
                        "type": source_type,
                        "chunks_found": 1
                    }
                else:
                    sources[source]["chunks_found"] += 1
            
            return {
                "success": True,
                "sources_count": len(sources),
                "sources": list(sources.values())
            }
            
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    elif tool_name == "get_source_content":
        source_name = arguments.get("source_name", "")
        
        try:
            # Search for content from specific source
            query_vector = embeddings.embed_query(source_name)
            results = index.query(vector=query_vector, top_k=50, include_metadata=True)
            
            # Filter by source
            source_content = []
            for match in results.matches:
                match_source = match.metadata.get("source", "")
                if source_name.lower() in match_source.lower():
                    content = match.metadata.get("content", "") or match.metadata.get("text", "")
                    source_content.append({
                        "content": content[:2000],
                        "title": match.metadata.get("title", ""),
                        "type": match.metadata.get("type", "")
                    })
            
            return {
                "success": True,
                "source": source_name,
                "chunks_found": len(source_content),
                "content": source_content[:20]  # Limit to 20 chunks
            }
            
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    elif tool_name == "search_web":
        query = arguments.get("query", "")
        
        try:
            # Use requests to call a search API (DuckDuckGo or similar)
            import requests
            
            # Use DuckDuckGo Instant Answer API
            ddg_url = "https://api.duckduckgo.com/"
            params = {
                "q": query,
                "format": "json",
                "no_html": 1,
                "skip_disambig": 1
            }
            
            response = requests.get(ddg_url, params=params, timeout=10)
            data = response.json()
            
            results = []
            
            # Get abstract/summary
            if data.get("Abstract"):
                results.append({
                    "content": data["Abstract"],
                    "source": data.get("AbstractURL", "DuckDuckGo"),
                    "title": data.get("Heading", query),
                    "type": "web_search"
                })
            
            # Get related topics
            for topic in data.get("RelatedTopics", [])[:5]:
                if isinstance(topic, dict) and topic.get("Text"):
                    results.append({
                        "content": topic["Text"],
                        "source": topic.get("FirstURL", ""),
                        "title": topic.get("Text", "")[:50],
                        "type": "web_search"
                    })
            
            # If no results from DDG, use a simple fallback
            if not results:
                # Try a more direct search approach
                search_url = f"https://html.duckduckgo.com/html/?q={query}"
                headers = {'User-Agent': 'Mozilla/5.0'}
                search_response = requests.get(search_url, headers=headers, timeout=10)
                
                if search_response.status_code == 200:
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(search_response.text, 'html.parser')
                    
                    for result in soup.select('.result')[:5]:
                        title_elem = result.select_one('.result__title')
                        snippet_elem = result.select_one('.result__snippet')
                        link_elem = result.select_one('.result__url')
                        
                        if title_elem and snippet_elem:
                            results.append({
                                "content": snippet_elem.get_text(strip=True),
                                "source": link_elem.get_text(strip=True) if link_elem else "",
                                "title": title_elem.get_text(strip=True),
                                "type": "web_search"
                            })
            
            return {
                "success": True,
                "query": query,
                "results_count": len(results),
                "results": results,
                "source_type": "web"
            }
            
        except Exception as e:
            print(f"Web search failed: {e}")
            return {"success": False, "error": str(e), "results": []}
    
    return {"success": False, "error": f"Unknown tool: {tool_name}"}


async def run_agent_with_streaming(question: str, user_id: str, conversation_id: str):
    """
    Run the agentic RAG with streaming.
    The LLM decides what to search, evaluates results, and answers.
    """
    
    # System prompt for the agent
    system_prompt = AGENT_SYSTEM_PROMPT

    # Get conversation history
    conv = load_conversation(conversation_id) if conversation_id else None
    messages_history = []
    
    if conv and conv.get('messages'):
        # Add recent conversation for context (last 6 messages)
        for msg in conv['messages'][-15:]:  # More conversation history
            role = "user" if msg['role'] == 'user' else "assistant"
            messages_history.append({
                "role": role,
                "content": msg['content'][:3000]  # Larger content window  # Limit size
            })
    
    # Build messages for the agent
    messages = [
        {"role": "system", "content": system_prompt}
    ] + messages_history + [
        {"role": "user", "content": question}
    ]
    
    # Track sources found
    all_sources = []
    tool_calls_made = []
    
    # Agent loop - max 5 iterations
    max_iterations = 12  # Increased for deeper research
    iteration = 0
    
    while iteration < max_iterations:
        iteration += 1
        
        # Call OpenAI with tools
        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=messages,
            tools=AGENT_TOOLS,
            tool_choice="auto",
            temperature=0.7
        )
        
        assistant_message = response.choices[0].message
        
        # Check if the model wants to use tools
        if assistant_message.tool_calls:
            # Add assistant's message with tool calls
            messages.append({
                "role": "assistant",
                "content": assistant_message.content or "",
                "tool_calls": [
                    {
                        "id": tc.id,
                        "type": "function",
                        "function": {
                            "name": tc.function.name,
                            "arguments": tc.function.arguments
                        }
                    } for tc in assistant_message.tool_calls
                ]
            })
            
            # Execute each tool call
            for tool_call in assistant_message.tool_calls:
                tool_name = tool_call.function.name
                try:
                    arguments = json.loads(tool_call.function.arguments)
                except:
                    arguments = {}
                
                # Stream thinking
                thinking_msg = {
                    "type": "thinking",
                    "step": tool_name,
                    "detail": arguments.get("query", str(arguments))
                }
                yield f"data: {json.dumps(thinking_msg)}\n\n"
                
                # Execute tool
                result = execute_tool(tool_name, arguments)
                
                # Track tool call
                tool_calls_made.append({
                    "tool": tool_name,
                    "arguments": arguments,
                    "results_count": result.get("results_count", result.get("chunks_found", 0))
                })
                
                # Collect sources
                if result.get("results"):
                    for r in result["results"]:
                        source_info = {
                            "source": r.get("source", ""),
                            "title": r.get("title", ""),
                            "type": r.get("type", "knowledge_base"),
                            "score": r.get("relevance_score", r.get("score", 0))
                        }
                        if source_info not in all_sources:
                            all_sources.append(source_info)
                
                # Stream tool result status
                result_msg = {
                    "type": "tool_result",
                    "tool": tool_name,
                    "found": result.get("results_count", result.get("chunks_found", 0)),
                    "success": result.get("success", False)
                }
                yield f"data: {json.dumps(result_msg)}\n\n"
                
                # Add tool result to messages
                messages.append({
                    "role": "tool",
                    "tool_call_id": tool_call.id,
                    "content": json.dumps(result)[:15000]  # Limit size
                })
        
        else:
            # No more tool calls - model is ready to answer
            break
    
    # Send sources before final answer
    sources_msg = {
        "type": "sources",
        "data": all_sources[:10],  # Top 10 sources
        "conversation_id": conversation_id
    }
    yield f"data: {json.dumps(sources_msg)}\n\n"
    
    # Stream the final answer
    thinking_msg = {"type": "thinking", "step": "generating", "detail": "Writing answer..."}
    yield f"data: {json.dumps(thinking_msg)}\n\n"
    
    # Get streaming response for final answer
    stream = openai_client.chat.completions.create(
        model="gpt-4o",
        messages=messages,
        temperature=0.7,
        stream=True
    )
    
    full_answer = ""
    for chunk in stream:
        if chunk.choices[0].delta.content:
            content = chunk.choices[0].delta.content
            full_answer += content
            token_msg = {"type": "token", "content": content}
            yield f"data: {json.dumps(token_msg)}\n\n"
    
    # Save to conversation
    if conversation_id:
        add_message_to_conversation(conversation_id, "assistant", full_answer, all_sources[:5])
    
    # Send done message
    done_msg = {
        "type": "done",
        "tool_calls": tool_calls_made,
        "sources_count": len(all_sources)
    }
    yield f"data: {json.dumps(done_msg)}\n\n"


@app.get("/chat/agent")
async def chat_agent_stream(
    message: str, 
    user_id: str = "default_user", 
    conversation_id: str = None
):
    """
    Agentic RAG chat endpoint with streaming.
    The LLM decides what to search and can make multiple searches.
    Shows thinking process in real-time.
    """
    # Create conversation if needed
    if not conversation_id:
        conversation = create_conversation(message)
        conversation_id = conversation['id']
    
    # Add user message
    add_message_to_conversation(conversation_id, "user", message)
    
    return StreamingResponse(
        run_agent_with_streaming(message, user_id, conversation_id),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no"
        }
    )


# ==================== Chat Endpoint (Original) ====================

@app.post("/chat", response_model=ChatResponse)
def chat(request: ChatRequest):
    question = request.message
    user_id = request.user_id
    conversation_id = request.conversation_id
    
    if not conversation_id:
        conversation = create_conversation(question)
        conversation_id = conversation['id']
    
    add_message_to_conversation(conversation_id, "user", question)
    
    conv = load_conversation(conversation_id)
    conversation_history = ""
    if conv and len(conv['messages']) > 1:
        recent_messages = conv['messages'][-7:-1]
        if recent_messages:
            conversation_history = "Recent Conversation:\n"
            for msg in recent_messages:
                role = "User" if msg['role'] == 'user' else "Assistant"
                content = msg['content'][:1500] + "..." if len(msg['content']) > 1500 else msg['content']
                conversation_history += f"{role}: {content}\n\n"
    
    search_query = question
    if conv and len(conv['messages']) > 2:
        user_messages = [m['content'] for m in conv['messages'] if m['role'] == 'user'][-3:]
        search_query = " ".join(user_messages)
    
    doc_context, sources, images, tracker = get_document_context_with_sources(search_query)
    
    # Add image info to prompt with full URLs for inline embedding
    image_info = ""
    if images:
        image_info = "\n\n=== SCREENSHOTS AVAILABLE (YOU MUST EMBED THESE) ===\n"
        for i, img in enumerate(images[:10], 1):
            url = img.get('url', '')
            image_info += f"SCREENSHOT_{i}: ![Screenshot {i}]({url})\n"
        image_info += "=== END SCREENSHOTS ===\n"
        image_info += "\nIMPORTANT: Copy and paste the SCREENSHOT lines above into your response where relevant!\n"
    
    prompt = f"""You are a friendly and helpful AI assistant with access to a knowledge base.

{conversation_history}

Available Information (Documents + Memory):
{doc_context if doc_context.strip() else "No specific information found."}
{image_info}

User Message: {question}

Instructions:
- Be friendly, helpful, and conversational
- If greeting, respond with a friendly greeting
- Use conversation history to understand context

**IMPORTANT - Be Smart About Matching:**
- Use the information above to answer, but be INTELLIGENT about matching
- Understand that users may use different names for the same product/service
- Example: "Database cloud service" likely refers to "Base Database Service"
- Example: "compute E4" means "Compute - Standard - E4"
- Use your knowledge to connect user's question to the available information
- If information seems related but named differently, use it and answer helpfully

**When to say "I don't have information":**
- ONLY say this if there's truly NO relevant information in the context
- Do NOT say this just because the exact wording doesn't match
- Try to find connections and answer helpfully first

- Format using markdown when helpful (headers, bullet points, bold)
- NEVER mention "documents", "context", "memory" in your response
- NEVER list sources in text - shown separately
- Be thorough - include ALL relevant details
- When asked to list ALL items, include EVERY item found

**SCREENSHOT EMBEDDING RULE**:
If SCREENSHOTS are listed above, you MUST embed them in your response.
Simply COPY the exact line like: ![Screenshot 1](/doc-images/xxx.jpeg)
Place screenshots after the relevant step they illustrate.

Response:"""
    
    # Step 6: Generate response
    if tracker:
        tracker.start_step("generate")
    
    response = llm.invoke(prompt)
    answer = response.content
    
    # Calculate and log LLM cost
    llm_calc = CostTracker.calculate_llm_cost(prompt, answer, "gpt-4o")
    CostTracker.log_cost(
        service="openai",
        model="gpt-4o",
        operation="generation",
        quantity=llm_calc["total_tokens"],
        unit="tokens",
        unit_price=0,  # Complex pricing (input/output different)
        total_cost=llm_calc["total_cost"],
        source=question[:50],
        pipeline_id=tracker.id if tracker else "",
        input_tokens=llm_calc["input_tokens"],
        output_tokens=llm_calc["output_tokens"]
    )
    
    # Complete pipeline tracking
    if tracker:
        tracker.complete_step("generate", 
            model="gpt-4o",
            prompt_length=len(prompt),
            response_length=len(answer),
            cost=llm_calc["total_cost"]
        )
        tracker.finish(
            sources_count=len(sources),
            has_images=len(images) > 0,
            total_cost=llm_calc["total_cost"]
        )
        if tracker.id in active_pipelines:
            del active_pipelines[tracker.id]
    
    add_message_to_conversation(conversation_id, "assistant", answer, sources)
    
    if doc_context.strip():
        save_to_memory(question, answer, user_id)
    
    source_objects = [Source(**s) for s in sources]
    
    return ChatResponse(
        answer=answer,
        memories_found=len([s for s in sources if s['type'] == 'memory']),
        sources=source_objects,
        images=[],  # No longer needed separately since images are inline
        conversation_id=conversation_id
    )

# ==================== Upload Endpoints ====================

# Image storage folder for extracted document images
DOC_IMAGES_FOLDER = os.path.join(DATA_DIR, "doc_images")
os.makedirs(DOC_IMAGES_FOLDER, exist_ok=True)

def extract_images_from_pdf(pdf_path: str, doc_id: str) -> List[dict]:
    """Extract images from PDF using PyMuPDF"""
    images = []
    try:
        import fitz  # PyMuPDF
        
        pdf_doc = fitz.open(pdf_path)
        
        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            image_list = page.get_images()
            
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = pdf_doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # Skip very small images (likely icons)
                    if len(image_bytes) < 5000:  # Less than 5KB
                        continue
                    
                    # Save image
                    img_filename = f"{doc_id}_page{page_num+1}_img{img_index+1}.{image_ext}"
                    img_path = os.path.join(DOC_IMAGES_FOLDER, img_filename)
                    
                    with open(img_path, "wb") as f:
                        f.write(image_bytes)
                    
                    images.append({
                        'url': f"/doc-images/{img_filename}",
                        'alt': f"Image from page {page_num + 1}",
                        'title': f"Page {page_num + 1}, Image {img_index + 1}",
                        'context': f"Screenshot from page {page_num + 1} of the document",
                        'local_path': img_path
                    })
                    
                except Exception as e:
                    print(f"Error extracting image {img_index} from page {page_num}: {e}")
                    continue
        
        pdf_doc.close()
        print(f"ðŸ“Â¸ Extracted {len(images)} images from PDF")
        
    except ImportError:
        print("âš ï¸ PyMuPDF not installed. Install with: pip install PyMuPDF")
    except Exception as e:
        print(f"Error extracting images from PDF: {e}")
    
    return images

def extract_images_from_docx(docx_path: str, doc_id: str) -> List[dict]:
    """Extract images from DOCX file"""
    images = []
    try:
        from docx import Document as DocxDocument
        import zipfile
        
        # DOCX is a zip file, extract images from word/media folder
        with zipfile.ZipFile(docx_path, 'r') as zip_ref:
            for file_info in zip_ref.filelist:
                if file_info.filename.startswith('word/media/'):
                    # Skip small files (likely icons)
                    if file_info.file_size < 5000:
                        continue
                    
                    ext = os.path.splitext(file_info.filename)[1]
                    if ext.lower() in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff']:
                        img_data = zip_ref.read(file_info.filename)
                        
                        img_index = len(images) + 1
                        img_filename = f"{doc_id}_img{img_index}{ext}"
                        img_path = os.path.join(DOC_IMAGES_FOLDER, img_filename)
                        
                        with open(img_path, "wb") as f:
                            f.write(img_data)
                        
                        images.append({
                            'url': f"/doc-images/{img_filename}",
                            'alt': f"Image {img_index} from document",
                            'title': f"Document Image {img_index}",
                            'context': f"Screenshot {img_index} from the document",
                            'local_path': img_path
                        })
        
        print(f"ðŸ“Â¸ Extracted {len(images)} images from DOCX")
        
    except Exception as e:
        print(f"Error extracting images from DOCX: {e}")
    
    return images

def extract_images_from_pptx(pptx_path: str, doc_id: str) -> List[dict]:
    """Extract images from PPTX file"""
    images = []
    try:
        import zipfile
        
        # PPTX is a zip file, extract images from ppt/media folder
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            for file_info in zip_ref.filelist:
                if file_info.filename.startswith('ppt/media/'):
                    # Skip small files
                    if file_info.file_size < 5000:
                        continue
                    
                    ext = os.path.splitext(file_info.filename)[1]
                    if ext.lower() in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff']:
                        img_data = zip_ref.read(file_info.filename)
                        
                        img_index = len(images) + 1
                        img_filename = f"{doc_id}_slide_img{img_index}{ext}"
                        img_path = os.path.join(DOC_IMAGES_FOLDER, img_filename)
                        
                        with open(img_path, "wb") as f:
                            f.write(img_data)
                        
                        images.append({
                            'url': f"/doc-images/{img_filename}",
                            'alt': f"Slide image {img_index}",
                            'title': f"Presentation Image {img_index}",
                            'context': f"Screenshot {img_index} from the presentation",
                            'local_path': img_path
                        })
        
        print(f"ðŸ“Â¸ Extracted {len(images)} images from PPTX")
        
    except Exception as e:
        print(f"Error extracting images from PPTX: {e}")
    
    return images

# Serve document images
@app.get("/doc-images/{filename}")
async def get_doc_image(filename: str):
    """Serve extracted document images"""
    file_path = os.path.join(DOC_IMAGES_FOLDER, filename)
    if os.path.exists(file_path):
        # Determine media type
        ext = os.path.splitext(filename)[1].lower()
        media_types = {
            '.png': 'image/png',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.gif': 'image/gif',
            '.bmp': 'image/bmp',
            '.tiff': 'image/tiff'
        }
        media_type = media_types.get(ext, 'image/png')
        return FileResponse(file_path, media_type=media_type)
    return {"error": "Image not found"}

@app.post("/upload", response_model=IngestResponse)
async def upload_file(file: UploadFile = File(...)):
    try:
        suffix = os.path.splitext(file.filename)[1].lower()
        
        file_path = os.path.join(UPLOADS_FOLDER, file.filename)
        with open(file_path, "wb") as f:
            content = await file.read()
            f.write(content)
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        
        # ==================== PDF ====================
        if suffix == ".pdf":
            doc_id = f"pdf_{uuid.uuid4().hex[:8]}"
            images = extract_images_from_pdf(tmp_path, doc_id)
            
            if UNSTRUCTURED_AVAILABLE:
                print(f"ðŸ”Â§ Using Unstructured.io for PDF: {file.filename}")
                try:
                    chunks_with_metadata, chunker_type = parse_with_unstructured(tmp_path, "pdf")
                    if chunks_with_metadata:
                        full_text = "\n\n".join([c["text"] for c in chunks_with_metadata])
                        metadata = extract_metadata(full_text[:3000], "document")
                        num_chunks, chunker_type, element_types = ingest_with_unstructured(
                            chunks_with_metadata, full_text, file.filename, "document", 
                            file_path, metadata, chunker_type=chunker_type, images=images
                        )
                        element_summary = ", ".join([f"{v} {k}" for k, v in element_types.items()])
                        os.unlink(tmp_path)
                        return IngestResponse(success=True, message=f"Ã¢Å“â€¦ Ingested {file.filename} with {num_chunks} chunks ðŸ”Â§ ({element_summary}), {len(images)} images", chunks=num_chunks, chunker_type=chunker_type, element_types=element_types)
                except Exception as e:
                    print(f"âš ï¸ Unstructured PDF failed: {e}, using fallback")
            
            loader = PyPDFLoader(tmp_path)
            documents = loader.load()
            full_text = "\n".join([doc.page_content for doc in documents])
            metadata = extract_metadata(full_text, "document")
            chunks, chunker_type = ingest_document_with_semantic_chunks(full_text, file.filename, "document", file_path, metadata, images=images)
            chunker_emoji = "ðŸ§ " if chunker_type == "semantic" else "ðŸ“"
            os.unlink(tmp_path)
            return IngestResponse(success=True, message=f"Ã¢Å“â€¦ Ingested {file.filename} with {chunks} chunks {chunker_emoji}, {len(images)} images", chunks=chunks, chunker_type=chunker_type)
        
        # ==================== TXT ====================
        elif suffix == ".txt":
            loader = TextLoader(tmp_path)
            documents = loader.load()
            full_text = "\n".join([doc.page_content for doc in documents])
            metadata = extract_metadata(full_text, "document")
            chunks, chunker_type = ingest_document_with_semantic_chunks(full_text, file.filename, "document", file_path, metadata)
            chunker_emoji = "ðŸ§ " if chunker_type == "semantic" else "ðŸ“"
            os.unlink(tmp_path)
            return IngestResponse(success=True, message=f"Ã¢Å“â€¦ Ingested {file.filename} with {chunks} chunks {chunker_emoji}", chunks=chunks, chunker_type=chunker_type)
        
        # ==================== DOCX ====================
        elif suffix == ".docx":
            doc_id = f"docx_{uuid.uuid4().hex[:8]}"
            images = extract_images_from_docx(tmp_path, doc_id)
            
            if UNSTRUCTURED_AVAILABLE:
                print(f"ðŸ”Â§ Using Unstructured.io for DOCX: {file.filename}")
                try:
                    chunks_with_metadata, chunker_type = parse_with_unstructured(tmp_path, "docx")
                    if chunks_with_metadata:
                        full_text = "\n\n".join([c["text"] for c in chunks_with_metadata])
                        metadata = extract_metadata(full_text[:3000], "document")
                        num_chunks, chunker_type, element_types = ingest_with_unstructured(
                            chunks_with_metadata, full_text, file.filename, "document", 
                            file_path, metadata, chunker_type=chunker_type, images=images
                        )
                        element_summary = ", ".join([f"{v} {k}" for k, v in element_types.items()])
                        os.unlink(tmp_path)
                        return IngestResponse(success=True, message=f"Ã¢Å“â€¦ Ingested {file.filename} with {num_chunks} chunks ðŸ”Â§ ({element_summary}), {len(images)} images", chunks=num_chunks, chunker_type=chunker_type, element_types=element_types)
                except Exception as e:
                    print(f"âš ï¸ Unstructured DOCX failed: {e}, using fallback")
            
            from docx import Document as DocxDocument
            doc = DocxDocument(tmp_path)
            full_text = "\n".join([para.text for para in doc.paragraphs if para.text])
            metadata = extract_metadata(full_text, "document")
            chunks, chunker_type = ingest_document_with_semantic_chunks(full_text, file.filename, "document", file_path, metadata, images=images)
            chunker_emoji = "ðŸ§ " if chunker_type == "semantic" else "ðŸ“"
            os.unlink(tmp_path)
            return IngestResponse(success=True, message=f"Ã¢Å“â€¦ Ingested {file.filename} with {chunks} chunks {chunker_emoji}, {len(images)} images", chunks=chunks, chunker_type=chunker_type)
        
        # ==================== XLSX ====================
        elif suffix in [".xlsx", ".xls"]:
            if UNSTRUCTURED_AVAILABLE:
                print(f"ðŸ”Â§ Using Unstructured.io for Excel: {file.filename}")
                try:
                    chunks_with_metadata, chunker_type = parse_with_unstructured(tmp_path, "xlsx")
                    if chunks_with_metadata:
                        full_text = "\n\n".join([c["text"] for c in chunks_with_metadata])
                        metadata = extract_metadata(full_text[:3000], "spreadsheet")
                        num_chunks, chunker_type, element_types = ingest_with_unstructured(
                            chunks_with_metadata, full_text, file.filename, "spreadsheet", 
                            file_path, metadata, chunker_type=chunker_type
                        )
                        element_summary = ", ".join([f"{v} {k}" for k, v in element_types.items()])
                        os.unlink(tmp_path)
                        return IngestResponse(success=True, message=f"Ã¢Å“â€¦ Ingested {file.filename} with {num_chunks} chunks ðŸ”Â§ ({element_summary})", chunks=num_chunks, chunker_type=chunker_type, element_types=element_types)
                except Exception as e:
                    print(f"âš ï¸ Unstructured Excel failed: {e}, using fallback")
            
            from openpyxl import load_workbook
            wb = load_workbook(tmp_path)
            full_text = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                full_text += f"\n--- Sheet: {sheet} ---\n"
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell else "" for cell in row])
                    if row_text.strip():
                        full_text += row_text + "\n"
            metadata = extract_metadata(full_text, "spreadsheet")
            chunks, chunker_type = ingest_document_with_semantic_chunks(full_text, file.filename, "spreadsheet", file_path, metadata)
            chunker_emoji = "ðŸ§ " if chunker_type == "semantic" else "ðŸ“"
            os.unlink(tmp_path)
            return IngestResponse(success=True, message=f"Ã¢Å“â€¦ Ingested {file.filename} with {chunks} chunks {chunker_emoji}", chunks=chunks, chunker_type=chunker_type)
        
        # ==================== PPTX ====================
        elif suffix == ".pptx":
            doc_id = f"pptx_{uuid.uuid4().hex[:8]}"
            images = extract_images_from_pptx(tmp_path, doc_id)
            
            if UNSTRUCTURED_AVAILABLE:
                print(f"ðŸ”Â§ Using Unstructured.io for PPTX: {file.filename}")
                try:
                    chunks_with_metadata, chunker_type = parse_with_unstructured(tmp_path, "pptx")
                    if chunks_with_metadata:
                        full_text = "\n\n".join([c["text"] for c in chunks_with_metadata])
                        metadata = extract_metadata(full_text[:3000], "presentation")
                        num_chunks, chunker_type, element_types = ingest_with_unstructured(
                            chunks_with_metadata, full_text, file.filename, "presentation", 
                            file_path, metadata, chunker_type=chunker_type, images=images
                        )
                        element_summary = ", ".join([f"{v} {k}" for k, v in element_types.items()])
                        os.unlink(tmp_path)
                        return IngestResponse(success=True, message=f"Ã¢Å“â€¦ Ingested {file.filename} with {num_chunks} chunks ðŸ”Â§ ({element_summary}), {len(images)} images", chunks=num_chunks, chunker_type=chunker_type, element_types=element_types)
                except Exception as e:
                    print(f"âš ï¸ Unstructured PPTX failed: {e}, using fallback")
            
            from pptx import Presentation
            prs = Presentation(tmp_path)
            full_text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                full_text += f"\n--- Slide {slide_num} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        full_text += shape.text + "\n"
            metadata = extract_metadata(full_text, "presentation")
            chunks, chunker_type = ingest_document_with_semantic_chunks(full_text, file.filename, "presentation", file_path, metadata, images=images)
            chunker_emoji = "ðŸ§ " if chunker_type == "semantic" else "ðŸ“"
            os.unlink(tmp_path)
            return IngestResponse(success=True, message=f"Ã¢Å“â€¦ Ingested {file.filename} with {chunks} chunks {chunker_emoji}, {len(images)} images", chunks=chunks, chunker_type=chunker_type)
        
        # ==================== Audio/Video ====================
        elif suffix in [".mp4", ".mp3", ".wav", ".m4a", ".webm", ".avi", ".mov"]:
            file_size_mb = os.path.getsize(tmp_path) / (1024 * 1024)
            if file_size_mb > 25:
                os.unlink(tmp_path)
                os.remove(file_path)
                return IngestResponse(success=False, message=f"Ã¢ÂÅ’ File too large ({file_size_mb:.1f}MB). Maximum size is 25MB for audio/video transcription.")
            
            from openai import OpenAI
            client = OpenAI()
            with open(tmp_path, "rb") as audio_file:
                transcription = client.audio.transcriptions.create(model="whisper-1", file=audio_file)
            
            full_text = transcription.text
            metadata = extract_metadata(full_text, "video")
            chunks, chunker_type = ingest_document_with_semantic_chunks(full_text, file.filename, "video", file_path, metadata)
            chunker_emoji = "ðŸ§ " if chunker_type == "semantic" else "ðŸ“"
            os.unlink(tmp_path)
            return IngestResponse(success=True, message=f"Ã¢Å“â€¦ Transcribed {file.filename} with {chunks} chunks {chunker_emoji}", chunks=chunks, chunker_type=chunker_type)
        
        # ==================== Images (Enhanced Prompt) ====================
        elif suffix in [".jpg", ".jpeg", ".png", ".gif", ".webp"]:
            import base64
            from langchain_core.messages import HumanMessage
            
            with open(tmp_path, "rb") as f:
                image_data = base64.b64encode(f.read()).decode("utf-8")
            
            media_types = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png", ".gif": "image/gif", ".webp": "image/webp"}
            media_type = media_types.get(suffix, "image/jpeg")
            
            # Enhanced prompt for better table and form extraction
            extraction_prompt = """Extract ALL content from this image carefully.

IMPORTANT RULES:
1. If there are TABLES in the image:
   - Preserve table structure using | (pipe) separators
   - Format each row as: Column1 | Column2 | Column3
   - Put each row on a new line
   - Include headers if visible

2. If there are FORMS or INVOICES:
   - Extract field labels and their values
   - Format as: Label: Value
   - Preserve any numbers, dates, amounts exactly

3. For Arabic or mixed language text:
   - Extract text in its original language
   - Maintain the reading order

4. For handwritten text:
   - Do your best to transcribe accurately
   - Mark unclear parts with [unclear]

Example table format:
Product | Quantity | Price
Item A | 10 | $50.00
Item B | 5 | $25.00

Extract all text and data from the image:"""
            
            message = HumanMessage(
                content=[
                    {"type": "text", "text": extraction_prompt},
                    {"type": "image_url", "image_url": {"url": f"data:{media_type};base64,{image_data}"}},
                ],
            )
            
            response = llm.invoke([message])
            full_text = response.content
            metadata = extract_metadata(full_text, "image")
            chunks, chunker_type = ingest_document_with_semantic_chunks(full_text, file.filename, "image", file_path, metadata)
            chunker_emoji = "ðŸ§ " if chunker_type == "semantic" else "ðŸ“"
            os.unlink(tmp_path)
            return IngestResponse(success=True, message=f"Ã¢Å“â€¦ Processed {file.filename} with {chunks} chunks {chunker_emoji}", chunks=chunks, chunker_type=chunker_type)
        
        else:
            os.unlink(tmp_path)
            os.remove(file_path)
            return IngestResponse(success=False, message=f"Unsupported file type: {suffix}")
    
    except Exception as e:
        return IngestResponse(success=False, message=str(e))

@app.post("/scrape", response_model=IngestResponse)
def scrape_website(request: UrlRequest):
    try:
        url = request.url
        headers = {'User-Agent': 'Mozilla/5.0'}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        html_content = response.text
        
        # Extract images BEFORE parsing
        images = extract_images_from_page(url, html_content)
        print(f"ðŸ–¼ï¸ Extracted {len(images)} images from {url}")
        
        from urllib.parse import urlparse
        domain = urlparse(url).netloc
        
        # Use Unstructured.io for smart parsing
        if UNSTRUCTURED_AVAILABLE:
            print(f"ðŸ”Â§ Using Unstructured.io for {url}")
            full_text, chunks_with_metadata, chunker_type = parse_html_with_unstructured(html_content, url)
            
            if not full_text or not chunks_with_metadata:
                return IngestResponse(success=False, message="No content found")
            
            metadata = extract_metadata(full_text[:3000], "website")
            
            num_chunks, chunker_type, element_types = ingest_with_unstructured(
                chunks_with_metadata, full_text, url, "website", "", metadata, 
                chunker_type=chunker_type, images=images
            )
            
            element_summary = ", ".join([f"{v} {k}" for k, v in element_types.items()])
            
            return IngestResponse(
                success=True, 
                message=f"Ã¢Å“â€¦ Scraped {domain} with {num_chunks} chunks ðŸ”Â§ ({element_summary}), {len(images)} images", 
                chunks=num_chunks, 
                chunker_type=chunker_type,
                element_types=element_types
            )
        else:
            # Fallback to BeautifulSoup
            print(f"âš ï¸ Unstructured not available, using BeautifulSoup for {url}")
            soup = BeautifulSoup(html_content, 'html.parser')
            for element in soup(['script', 'style', 'nav', 'footer', 'header']):
                element.decompose()
            
            text = soup.get_text(separator='\n', strip=True)
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            full_text = '\n'.join(lines)
            
            if not full_text:
                return IngestResponse(success=False, message="No content found")
            
            metadata = extract_metadata(full_text, "website")
            chunks, chunker_type = ingest_document_with_semantic_chunks(full_text, url, "website", "", metadata, images=images)
            chunker_emoji = "ðŸ§ " if chunker_type == "semantic" else "ðŸ“"
            
            return IngestResponse(success=True, message=f"Ã¢Å“â€¦ Scraped {domain} with {chunks} chunks {chunker_emoji}, {len(images)} images", chunks=chunks, chunker_type=chunker_type)
    
    except Exception as e:
        return IngestResponse(success=False, message=str(e))

# ==================== Recursive Web Scraper ====================

class RecursiveScrapeRequest(BaseModel):
    url: str
    max_depth: int = 2
    max_pages: int = 20
    same_path_only: bool = True

class RecursiveScrapeResponse(BaseModel):
    success: bool
    message: str
    total_pages: int = 0
    total_chunks: int = 0
    pages_scraped: List[str] = []

def extract_links_from_page(url: str, html_content: str, same_path_only: bool = True) -> List[str]:
    """Extract valid links from a page that belong to the same context"""
    from urllib.parse import urlparse, urljoin, unquote
    
    soup = BeautifulSoup(html_content, 'html.parser')
    parsed_base = urlparse(url)
    base_domain = parsed_base.netloc
    
    # Get path segments (decode URL encoding)
    base_path_decoded = unquote(parsed_base.path)
    base_segments = [s for s in base_path_decoded.split('/') if s]
    
    # For same_path_only, match first N segments (more flexible)
    # e.g., /docs/Atomsphere/Integration/ -> match first 3 segments
    if len(base_segments) >= 3:
        match_segments = base_segments[:3]  # Match first 3 path segments
    elif len(base_segments) >= 2:
        match_segments = base_segments[:2]
    else:
        match_segments = base_segments[:1] if base_segments else []
    
    match_path = '/' + '/'.join(match_segments) if match_segments else ''
    
    links = set()
    
    for a_tag in soup.find_all('a', href=True):
        href = a_tag['href']
        
        # Skip anchors, javascript, mailto, etc.
        if href.startswith('#') or href.startswith('javascript:') or href.startswith('mailto:') or href.startswith('tel:'):
            continue
        
        # Convert relative URLs to absolute
        full_url = urljoin(url, href)
        parsed_url = urlparse(full_url)
        
        # Must be same domain
        if parsed_url.netloc != base_domain:
            continue
        
        # Skip non-http links
        if parsed_url.scheme not in ['http', 'https']:
            continue
        
        # Skip file downloads
        if any(parsed_url.path.lower().endswith(ext) for ext in ['.pdf', '.zip', '.doc', '.docx', '.xls', '.xlsx', '.png', '.jpg', '.jpeg', '.gif']):
            continue
        
        # If same_path_only, check if the link shares the same base path segments
        if same_path_only and match_path:
            link_path_decoded = unquote(parsed_url.path)
            if not link_path_decoded.startswith(match_path):
                continue
        
        # Clean URL (remove fragments and trailing slashes)
        clean_url = f"{parsed_url.scheme}://{parsed_url.netloc}{parsed_url.path}"
        clean_url = clean_url.rstrip('/')
        
        links.add(clean_url)
    
    return list(links)

def extract_images_from_page(url: str, html_content: str, min_size: int = 100) -> List[dict]:
    """Extract meaningful images from a page (skip icons, logos, etc.)"""
    from urllib.parse import urlparse, urljoin
    
    soup = BeautifulSoup(html_content, 'html.parser')
    images = []
    
    for img in soup.find_all('img'):
        src = img.get('src', '')
        if not src:
            continue
        
        # Convert relative URLs to absolute
        img_url = urljoin(url, src)
        
        # Skip data URIs and tiny images (likely icons)
        if img_url.startswith('data:'):
            continue
        
        # Skip common icon/logo patterns
        skip_patterns = ['icon', 'logo', 'avatar', 'emoji', 'badge', 'button', 'sprite', 'tracking', 'pixel', 'spacer']
        if any(pattern in img_url.lower() for pattern in skip_patterns):
            continue
        
        # Get alt text and surrounding context
        alt_text = img.get('alt', '')
        title = img.get('title', '')
        
        # Try to get context from parent elements
        context = ""
        parent = img.parent
        for _ in range(3):  # Go up 3 levels max
            if parent:
                # Look for nearby text
                text = parent.get_text(strip=True)[:200] if parent.get_text(strip=True) else ""
                if text and len(text) > len(context):
                    context = text
                parent = parent.parent
        
        # Check for figure caption
        figure = img.find_parent('figure')
        if figure:
            figcaption = figure.find('figcaption')
            if figcaption:
                context = figcaption.get_text(strip=True)
        
        images.append({
            'url': img_url,
            'alt': alt_text,
            'title': title,
            'context': context[:300] if context else alt_text,
            'source_page': url
        })
    
    return images

def scrape_single_page(url: str) -> tuple:
    """Scrape a single page and return (text_content, html_content, images, success)"""
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        response = requests.get(url, headers=headers, timeout=15)
        response.raise_for_status()
        
        html_content = response.text
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Extract images BEFORE removing elements
        images = extract_images_from_page(url, html_content)
        
        # Remove unwanted elements
        for element in soup(['script', 'style', 'nav', 'footer', 'header', 'aside', 'iframe', 'noscript']):
            element.decompose()
        
        # Extract text
        text = soup.get_text(separator='\n', strip=True)
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        full_text = '\n'.join(lines)
        
        return full_text, html_content, images, True
    except Exception as e:
        print(f"Error scraping {url}: {e}")
        return "", "", [], False

@app.post("/scrape-recursive", response_model=RecursiveScrapeResponse)
def scrape_website_recursive(request: RecursiveScrapeRequest):
    """Recursively scrape a website and all its subpages within the same context"""
    try:
        from urllib.parse import urlparse
        
        start_url = request.url.rstrip('/')
        max_depth = min(request.max_depth, 3)  # Cap at 3 to prevent infinite loops
        max_pages = min(request.max_pages, 50)  # Cap at 50 pages
        same_path_only = request.same_path_only
        
        parsed_start = urlparse(start_url)
        domain = parsed_start.netloc
        
        # Track visited URLs and pages to scrape
        visited = set()
        to_visit = [(start_url, 0)]  # (url, depth)
        pages_scraped = []
        total_chunks = 0
        
        while to_visit and len(pages_scraped) < max_pages:
            current_url, depth = to_visit.pop(0)
            
            # Skip if already visited
            if current_url in visited:
                continue
            
            visited.add(current_url)
            
            # Scrape the page (now returns 4 values including images)
            text_content, html_content, images, success = scrape_single_page(current_url)
            
            if not success or not text_content or len(text_content) < 100:
                continue
            
            # Ingest the content with images
            try:
                metadata = extract_metadata(text_content, "website")
                chunks, chunker_type = ingest_document_with_semantic_chunks(text_content, current_url, "website", "", metadata, images=images)
                total_chunks += chunks
                pages_scraped.append(current_url)
                chunker_emoji = "ðŸ§ " if chunker_type == "semantic" else "ðŸ“"
                print(f"Ã¢Å“â€¦ Scraped: {current_url} ({chunks} chunks {chunker_emoji}, {len(images)} images)")
            except Exception as e:
                print(f"Error ingesting {current_url}: {e}")
                continue
            
            # Find more links if we haven't reached max depth
            if depth < max_depth and len(pages_scraped) < max_pages:
                new_links = extract_links_from_page(current_url, html_content, same_path_only)
                
                for link in new_links:
                    if link not in visited and (link, depth + 1) not in to_visit:
                        to_visit.append((link, depth + 1))
        
        if not pages_scraped:
            return RecursiveScrapeResponse(
                success=False,
                message="No pages could be scraped",
                total_pages=0,
                total_chunks=0,
                pages_scraped=[]
            )
        
        return RecursiveScrapeResponse(
            success=True,
            message=f"Ã¢Å“â€¦ Scraped {len(pages_scraped)} pages from {domain} with {total_chunks} total chunks",
            total_pages=len(pages_scraped),
            total_chunks=total_chunks,
            pages_scraped=pages_scraped
        )
    
    except Exception as e:
        return RecursiveScrapeResponse(
            success=False,
            message=str(e),
            total_pages=0,
            total_chunks=0,
            pages_scraped=[]
        )

@app.post("/youtube", response_model=IngestResponse)
def process_youtube(request: UrlRequest):
    try:
        url = request.url
        patterns = [r'(?:v=|\/)([0-9A-Za-z_-]{11}).*', r'(?:youtu\.be\/)([0-9A-Za-z_-]{11})']
        video_id = None
        for pattern in patterns:
            match = re.search(pattern, url)
            if match:
                video_id = match.group(1)
                break
        
        if not video_id:
            return IngestResponse(success=False, message="Invalid YouTube URL")
        
        ytt_api = YouTubeTranscriptApi()
        transcript_data = ytt_api.fetch(video_id)
        full_text = ' '.join([entry.text for entry in transcript_data])
        
        if not full_text:
            return IngestResponse(success=False, message="No transcript available")
        
        metadata = extract_metadata(full_text, "youtube video")
        chunks, chunker_type = ingest_document_with_semantic_chunks(full_text, url, "youtube", "", metadata)
        chunker_emoji = "ðŸ§ " if chunker_type == "semantic" else "ðŸ“"
        
        return IngestResponse(success=True, message=f"Ã¢Å“â€¦ Processed YouTube with {chunks} chunks {chunker_emoji}", chunks=chunks, chunker_type=chunker_type)
    
    except Exception as e:
        return IngestResponse(success=False, message=str(e))

# ==================== Debug Endpoint ====================

@app.get("/debug/search")
def debug_search(query: str, top_k: int = 10):
    """
    Debug endpoint to see exactly what chunks are retrieved for a query.
    Use this to diagnose retrieval issues.
    """
    try:
        # Get query variations from LLM
        query_variations = rewrite_query_with_llm(query)
        
        all_results = []
        seen_content = set()
        
        for variation in query_variations:
            # Embed and search
            query_vector = embeddings.embed_query(variation)
            results = index.query(vector=query_vector, top_k=top_k, include_metadata=True)
            
            for match in results.matches:
                content = (match.metadata.get("content", "") or match.metadata.get("text", ""))[:200]
                if content not in seen_content:
                    seen_content.add(content)
                    all_results.append({
                        "variation_used": variation,
                        "score": round(float(match.score), 4),
                        "source": match.metadata.get("source", "")[:100],
                        "title": match.metadata.get("title", "")[:100],
                        "type": match.metadata.get("type", ""),
                        "content_preview": content
                    })
        
        # Sort by score
        all_results.sort(key=lambda x: x["score"], reverse=True)
        
        return {
            "original_query": query,
            "query_variations": query_variations,
            "total_results": len(all_results),
            "top_results": all_results[:top_k]
        }
        
    except Exception as e:
        return {"error": str(e)}

@app.get("/debug/chunks")
def debug_chunks_by_source(source_contains: str = "", limit: int = 20):
    """
    Debug endpoint to see what chunks exist for a specific source.
    Example: /debug/chunks?source_contains=azure
    """
    try:
        # This requires listing vectors which isn't directly supported
        # So we'll search with a broad query
        query_vector = embeddings.embed_query(source_contains or "pricing")
        results = index.query(vector=query_vector, top_k=100, include_metadata=True)
        
        filtered = []
        for match in results.matches:
            source = match.metadata.get("source", "").lower()
            title = match.metadata.get("title", "").lower()
            if source_contains.lower() in source or source_contains.lower() in title:
                content = match.metadata.get("content", "") or match.metadata.get("text", "")
                filtered.append({
                    "score": round(float(match.score), 4),
                    "source": match.metadata.get("source", ""),
                    "title": match.metadata.get("title", ""),
                    "type": match.metadata.get("type", ""),
                    "content_preview": content[:300]
                })
        
        return {
            "filter": source_contains,
            "found": len(filtered),
            "chunks": filtered[:limit]
        }
        
    except Exception as e:
        return {"error": str(e)}

@app.get("/stats")
def get_stats():
    try:
        # Get Pinecone stats
        try:
            stats = index.describe_index_stats()
            total_vectors = stats.total_vector_count
        except Exception as e:
            print(f"Pinecone stats error: {e}")
            total_vectors = 0
        
        # Count local documents
        try:
            os.makedirs(DOCUMENTS_FOLDER, exist_ok=True)
            doc_count = len([f for f in os.listdir(DOCUMENTS_FOLDER) if f.endswith('.json')])
        except Exception as e:
            print(f"Document count error: {e}")
            doc_count = 0
        
        # Count BM25 index
        try:
            bm25_data = load_bm25_index()
            bm25_count = len(bm25_data.get("documents", []))
        except Exception as e:
            print(f"BM25 index error: {e}")
            bm25_count = 0
        
        return {
            "total_vectors": total_vectors,
            "total_documents": doc_count,
            "bm25_index_size": bm25_count,
            "index_name": os.getenv("PINECONE_INDEX_NAME"),
            "chunking_type": "Semantic",
            "search_type": "Hybrid (All Sources + Memory)"
        }
    except Exception as e:
        print(f"Stats error: {e}")
        return {
            "total_vectors": 0,
            "total_documents": 0,
            "bm25_index_size": 0,
            "index_name": os.getenv("PINECONE_INDEX_NAME"),
            "chunking_type": "Semantic",
            "search_type": "Hybrid (All Sources + Memory)",
            "error": str(e)
        }

# ==================== Pipeline Monitor Endpoints ====================

@app.get("/pipeline/logs")
def get_pipeline_logs(limit: int = 20, pipeline_type: str = None):
    """Get recent pipeline execution logs"""
    try:
        logs = load_pipeline_logs()
        
        # Filter by type if specified
        if pipeline_type:
            logs = [l for l in logs if l.get("pipeline_type") == pipeline_type]
        
        # Return most recent first
        logs = sorted(logs, key=lambda x: x.get("start_time", ""), reverse=True)
        
        return {
            "logs": logs[:limit],
            "total": len(logs)
        }
    except Exception as e:
        return {"logs": [], "total": 0, "error": str(e)}

@app.get("/pipeline/logs/{log_id}")
def get_pipeline_log(log_id: str):
    """Get a specific pipeline log by ID"""
    try:
        logs = load_pipeline_logs()
        for log in logs:
            if log.get("id") == log_id:
                return log
        return {"error": "Pipeline log not found"}
    except Exception as e:
        return {"error": str(e)}

@app.get("/pipeline/active")
def get_active_pipelines():
    """Get currently running pipelines"""
    return {
        "active": [p.to_dict() for p in active_pipelines.values()]
    }

@app.get("/pipeline/stats")
def get_pipeline_stats():
    """Get aggregated pipeline statistics"""
    try:
        logs = load_pipeline_logs()
        
        if not logs:
            return {
                "total_runs": 0,
                "ingestion": {"count": 0, "avg_duration_ms": 0},
                "retrieval": {"count": 0, "avg_duration_ms": 0}
            }
        
        ingestion_logs = [l for l in logs if l.get("pipeline_type") == "ingestion"]
        retrieval_logs = [l for l in logs if l.get("pipeline_type") == "retrieval"]
        
        def calc_avg(logs_list):
            if not logs_list:
                return 0
            durations = [l.get("total_duration_ms", 0) for l in logs_list]
            return round(sum(durations) / len(durations), 2)
        
        def calc_step_avg(logs_list, step_name):
            durations = []
            for log in logs_list:
                steps = log.get("steps", {})
                if step_name in steps:
                    durations.append(steps[step_name].get("duration_ms", 0))
            return round(sum(durations) / len(durations), 2) if durations else 0
        
        # Calculate average duration for each step
        ingestion_steps = ["extract", "chunk", "embed", "store_vectors", "store_bm25"]
        retrieval_steps = ["embed_query", "pinecone_search", "bm25_search", "rrf_merge", "rerank", "generate"]
        
        ingestion_step_avgs = {step: calc_step_avg(ingestion_logs, step) for step in ingestion_steps}
        retrieval_step_avgs = {step: calc_step_avg(retrieval_logs, step) for step in retrieval_steps}
        
        return {
            "total_runs": len(logs),
            "ingestion": {
                "count": len(ingestion_logs),
                "avg_duration_ms": calc_avg(ingestion_logs),
                "step_averages": ingestion_step_avgs
            },
            "retrieval": {
                "count": len(retrieval_logs),
                "avg_duration_ms": calc_avg(retrieval_logs),
                "step_averages": retrieval_step_avgs
            },
            "last_24h": {
                "ingestion": len([l for l in ingestion_logs if l.get("start_time", "") > (datetime.now() - timedelta(hours=24)).isoformat()]),
                "retrieval": len([l for l in retrieval_logs if l.get("start_time", "") > (datetime.now() - timedelta(hours=24)).isoformat()])
            }
        }
    except Exception as e:
        return {"error": str(e)}

@app.delete("/pipeline/logs")
def clear_pipeline_logs():
    """Clear all pipeline logs"""
    try:
        if os.path.exists(PIPELINE_LOG_FILE):
            os.remove(PIPELINE_LOG_FILE)
        return {"success": True, "message": "Pipeline logs cleared"}
    except Exception as e:
        return {"success": False, "error": str(e)}

# ==================== Cost Monitor Endpoints ====================

@app.get("/cost/summary")
def get_cost_summary():
    """Get overall cost summary"""
    try:
        summary = load_cost_summary()
        return summary
    except Exception as e:
        return {"error": str(e)}

@app.get("/cost/logs")
def get_cost_logs(limit: int = 50, service: str = None, operation: str = None):
    """Get cost log entries"""
    try:
        logs = load_cost_logs()
        
        # Filter by service
        if service:
            logs = [l for l in logs if l.get("service") == service]
        
        # Filter by operation
        if operation:
            logs = [l for l in logs if l.get("operation") == operation]
        
        # Sort by timestamp (newest first)
        logs = sorted(logs, key=lambda x: x.get("timestamp", ""), reverse=True)
        
        return {
            "logs": logs[:limit],
            "total": len(logs)
        }
    except Exception as e:
        return {"logs": [], "total": 0, "error": str(e)}

@app.get("/cost/by-document")
def get_cost_by_document():
    """Get costs grouped by document/source"""
    try:
        logs = load_cost_logs()
        
        by_source = {}
        for log in logs:
            source = log.get("source", "Unknown")
            if source not in by_source:
                by_source[source] = {
                    "total_cost": 0,
                    "operations": {},
                    "count": 0
                }
            by_source[source]["total_cost"] += log.get("total_cost", 0)
            by_source[source]["count"] += 1
            
            operation = log.get("operation", "unknown")
            if operation not in by_source[source]["operations"]:
                by_source[source]["operations"][operation] = 0
            by_source[source]["operations"][operation] += log.get("total_cost", 0)
        
        # Round all costs
        for source in by_source:
            by_source[source]["total_cost"] = round(by_source[source]["total_cost"], 6)
            for op in by_source[source]["operations"]:
                by_source[source]["operations"][op] = round(by_source[source]["operations"][op], 6)
        
        return {"by_document": by_source}
    except Exception as e:
        return {"by_document": {}, "error": str(e)}

@app.get("/cost/pricing")
def get_pricing():
    """Get current API pricing"""
    return {"pricing": API_PRICING}

@app.get("/cost/calculate")
def calculate_cost(text: str = "", tokens: int = 0, operation: str = "embedding"):
    """Calculate estimated cost for a given operation"""
    try:
        if operation == "embedding":
            if text:
                result = CostTracker.calculate_embedding_cost(text)
            elif tokens:
                cost = (tokens / 1_000_000) * API_PRICING["openai"]["text-embedding-3-large"]["price_per_million_tokens"]
                result = {"tokens": tokens, "cost": round(cost, 8)}
            else:
                result = {"error": "Provide text or tokens"}
        elif operation == "rerank":
            result = CostTracker.calculate_rerank_cost(1)
        else:
            result = {"error": f"Unknown operation: {operation}"}
        
        return result
    except Exception as e:
        return {"error": str(e)}

@app.delete("/cost/logs")
def clear_cost_logs():
    """Clear all cost logs"""
    try:
        if os.path.exists(COST_LOG_FILE):
            os.remove(COST_LOG_FILE)
        if os.path.exists(COST_SUMMARY_FILE):
            os.remove(COST_SUMMARY_FILE)
        return {"success": True, "message": "Cost logs cleared"}
    except Exception as e:
        return {"success": False, "error": str(e)}

@app.get("/cost/logs/{log_id}")
def get_cost_log_detail(log_id: str):
    """Get detailed info for a specific cost log"""
    try:
        logs = load_cost_logs()
        for log in logs:
            if log.get("id") == log_id:
                return {"log": log}
        return {"error": "Cost log not found"}
    except Exception as e:
        return {"error": str(e)}

@app.get("/cost/analyze")
def analyze_query_cost(query: str):
    """Pre-analyze a query to estimate cost and provide optimization suggestions"""
    try:
        query_tokens = CostTracker.estimate_tokens(query)
        
        # Estimate typical context size based on average
        estimated_context_tokens = 3000  # Average context size
        estimated_output_tokens = 500    # Average response size
        
        total_input = query_tokens + estimated_context_tokens + 200  # +200 for system prompt
        
        # Calculate estimated cost
        pricing = API_PRICING["openai"]["gpt-4o"]
        input_cost = (total_input / 1_000_000) * pricing["input_price_per_million_tokens"]
        output_cost = (estimated_output_tokens / 1_000_000) * pricing["output_price_per_million_tokens"]
        rerank_cost = API_PRICING["cohere"]["rerank-v3.5"]["price_per_search"]
        embed_cost = (query_tokens / 1_000_000) * API_PRICING["openai"]["text-embedding-3-large"]["price_per_million_tokens"]
        
        total_estimated_cost = input_cost + output_cost + rerank_cost + embed_cost
        
        # Generate suggestions
        suggestions = CostTracker.generate_optimization_suggestions(
            question=query,
            context_tokens=estimated_context_tokens,
            total_tokens=total_input,
            operation="retrieval"
        )
        
        return {
            "query_tokens": query_tokens,
            "estimated_total_tokens": total_input + estimated_output_tokens,
            "estimated_cost": round(total_estimated_cost, 6),
            "breakdown": {
                "query_embedding": round(embed_cost, 8),
                "reranking": rerank_cost,
                "llm_input": round(input_cost, 6),
                "llm_output": round(output_cost, 6)
            },
            "suggestions": suggestions
        }
    except Exception as e:
        return {"error": str(e)}

@app.get("/monitor")
def serve_monitor():
    """Serve the Pipeline Monitor dashboard"""
    return FileResponse("frontend/monitor.html")

@app.get("/files")
def list_files():
    try:
        os.makedirs(UPLOADS_FOLDER, exist_ok=True)
        files = []
        for filename in os.listdir(UPLOADS_FOLDER):
            file_path = os.path.join(UPLOADS_FOLDER, filename)
            if os.path.isfile(file_path):
                files.append({
                    "filename": filename,
                    "size": os.path.getsize(file_path),
                    "download_url": f"/download/{filename}"
                })
        return {"files": files}
    except Exception as e:
        print(f"List files error: {e}")
        return {"files": [], "error": str(e)}

# ==================== Admin Dashboard Endpoints ====================

@app.get("/admin/documents")
def get_all_documents():
    try:
        os.makedirs(DOCUMENTS_FOLDER, exist_ok=True)
        documents = []
        
        for filename in os.listdir(DOCUMENTS_FOLDER):
            if filename.endswith('.json'):
                try:
                    doc_path = os.path.join(DOCUMENTS_FOLDER, filename)
                    with open(doc_path, 'r', encoding='utf-8') as f:
                        doc = json.load(f)
                        documents.append({
                            "id": doc['id'],
                            "title": doc['metadata'].get('title', 'Unknown'),
                            "source": doc['metadata'].get('source', 'Unknown'),
                            "type": doc['metadata'].get('type', 'document'),
                            "file_path": doc['metadata'].get('file_path', ''),
                            "summary": doc['metadata'].get('summary', ''),
                            "created_at": doc.get('created_at', ''),
                            "content_preview": doc['content'][:500] + "..." if len(doc['content']) > 500 else doc['content']
                        })
                except Exception as e:
                    print(f"Error reading document {filename}: {e}")
                    continue
        
        documents.sort(key=lambda x: x.get('created_at', ''), reverse=True)
        return {"documents": documents, "total": len(documents)}
    except Exception as e:
        print(f"Admin documents error: {e}")
        return {"documents": [], "total": 0, "error": str(e)}

@app.get("/admin/documents/{doc_id}")
def get_document_detail(doc_id: str):
    doc = load_full_document(doc_id)
    
    if not doc:
        return {"error": "Document not found"}
    
    chunks = []
    try:
        dummy_vector = embeddings.embed_query("test")
        results = index.query(
            vector=dummy_vector,
            top_k=100,
            include_metadata=True,
            filter={"parent_id": {"$eq": doc_id}}
        )
        
        for match in results['matches']:
            chunks.append({
                "id": match['id'],
                "text": match['metadata'].get('text', ''),
                "chunk_index": match['metadata'].get('chunk_index', 0),
                "chunk_type": match['metadata'].get('chunk_type', 'unknown')
            })
        
        chunks.sort(key=lambda x: x.get('chunk_index', 0))
    except Exception as e:
        print(f"Error fetching chunks: {e}")
    
    return {
        "document": doc,
        "chunks": chunks,
        "chunk_count": len(chunks)
    }

@app.delete("/admin/documents/{doc_id}")
def delete_document(doc_id: str):
    doc_path = os.path.join(DOCUMENTS_FOLDER, f"{doc_id}.json")
    if os.path.exists(doc_path):
        with open(doc_path, 'r') as f:
            doc = json.load(f)
        
        file_path = doc['metadata'].get('file_path', '')
        if file_path and os.path.exists(file_path):
            os.remove(file_path)
        
        os.remove(doc_path)
    
    try:
        dummy_vector = embeddings.embed_query("test")
        results = index.query(
            vector=dummy_vector,
            top_k=100,
            include_metadata=True,
            filter={"parent_id": {"$eq": doc_id}}
        )
        
        chunk_ids = [match['id'] for match in results['matches']]
        if chunk_ids:
            index.delete(ids=chunk_ids)
    except Exception as e:
        print(f"Error deleting from Pinecone: {e}")
    
    try:
        bm25_data = load_bm25_index()
        new_docs = []
        new_meta = []
        
        for doc_text, meta in zip(bm25_data['documents'], bm25_data['metadata']):
            if meta.get('parent_id') != doc_id:
                new_docs.append(doc_text)
                new_meta.append(meta)
        
        bm25_data['documents'] = new_docs
        bm25_data['metadata'] = new_meta
        save_bm25_index(bm25_data)
    except Exception as e:
        print(f"Error updating BM25: {e}")
    
    return {"success": True, "message": f"Deleted document {doc_id}"}

@app.get("/admin/memories")
def get_all_memories():
    memories = []
    
    try:
        dummy_vector = embeddings.embed_query("memory")
        results = index.query(
            vector=dummy_vector,
            top_k=100,
            include_metadata=True,
            filter={"type": {"$eq": "memory"}}
        )
        
        for match in results['matches']:
            memories.append({
                "id": match['id'],
                "question": match['metadata'].get('question', ''),
                "answer": match['metadata'].get('answer', '')[:200] + "...",
                "user_id": match['metadata'].get('user_id', ''),
                "timestamp": match['metadata'].get('timestamp', '')
            })
        
        memories.sort(key=lambda x: x.get('timestamp', ''), reverse=True)
    except Exception as e:
        print(f"Error fetching memories: {e}")
    
    return {"memories": memories, "total": len(memories)}

@app.delete("/admin/memories/{memory_id}")
def delete_memory(memory_id: str):
    try:
        index.delete(ids=[memory_id])
        return {"success": True, "message": f"Deleted memory {memory_id}"}
    except Exception as e:
        return {"success": False, "message": str(e)}

@app.delete("/admin/memories")
def clear_all_memories():
    try:
        dummy_vector = embeddings.embed_query("memory")
        results = index.query(
            vector=dummy_vector,
            top_k=1000,
            include_metadata=True,
            filter={"type": {"$eq": "memory"}}
        )
        
        memory_ids = [match['id'] for match in results['matches']]
        if memory_ids:
            index.delete(ids=memory_ids)
        
        return {"success": True, "message": f"Deleted {len(memory_ids)} memories"}
    except Exception as e:
        return {"success": False, "message": str(e)}

@app.delete("/admin/clear-all")
def clear_all_knowledge_base(include_memories: bool = False, include_conversations: bool = False):
    """
    Clear entire knowledge base (documents, websites, YouTube).
    Optionally include memories and conversations.
    """
    results = {
        "vectors_deleted": 0,
        "documents_deleted": 0,
        "bm25_cleared": False,
        "memories_deleted": 0,
        "conversations_deleted": 0
    }
    
    try:
        # 1. Delete all vectors from Pinecone (except memories unless specified)
        try:
            dummy_vector = embeddings.embed_query("document")
            
            if include_memories:
                # Delete ALL vectors
                all_results = index.query(
                    vector=dummy_vector,
                    top_k=10000,
                    include_metadata=True
                )
                all_ids = [match['id'] for match in all_results['matches']]
            else:
                # Delete only non-memory vectors
                all_results = index.query(
                    vector=dummy_vector,
                    top_k=10000,
                    include_metadata=True
                )
                all_ids = [
                    match['id'] for match in all_results['matches'] 
                    if match['metadata'].get('type') != 'memory'
                ]
            
            if all_ids:
                # Delete in batches of 1000
                for i in range(0, len(all_ids), 1000):
                    batch = all_ids[i:i+1000]
                    index.delete(ids=batch)
                results["vectors_deleted"] = len(all_ids)
                
        except Exception as e:
            print(f"Error deleting vectors: {e}")
        
        # 2. Delete local document files
        try:
            os.makedirs(DOCUMENTS_FOLDER, exist_ok=True)
            doc_files = [f for f in os.listdir(DOCUMENTS_FOLDER) if f.endswith('.json')]
            for f in doc_files:
                os.remove(os.path.join(DOCUMENTS_FOLDER, f))
            results["documents_deleted"] = len(doc_files)
        except Exception as e:
            print(f"Error deleting documents: {e}")
        
        # 3. Clear BM25 index
        try:
            bm25_path = BM25_INDEX_FILE
            if os.path.exists(bm25_path):
                os.remove(bm25_path)
            results["bm25_cleared"] = True
        except Exception as e:
            print(f"Error clearing BM25: {e}")
        
        # 4. Delete uploaded files
        try:
            os.makedirs(UPLOADS_FOLDER, exist_ok=True)
            upload_files = os.listdir(UPLOADS_FOLDER)
            for f in upload_files:
                file_path = os.path.join(UPLOADS_FOLDER, f)
                if os.path.isfile(file_path):
                    os.remove(file_path)
        except Exception as e:
            print(f"Error deleting uploads: {e}")
        
        # 5. Optionally delete conversations
        if include_conversations:
            try:
                os.makedirs(CONVERSATIONS_FOLDER, exist_ok=True)
                conv_files = [f for f in os.listdir(CONVERSATIONS_FOLDER) if f.endswith('.json')]
                for f in conv_files:
                    os.remove(os.path.join(CONVERSATIONS_FOLDER, f))
                results["conversations_deleted"] = len(conv_files)
            except Exception as e:
                print(f"Error deleting conversations: {e}")
        
        return {
            "success": True, 
            "message": "Knowledge base cleared successfully!",
            "details": results
        }
        
    except Exception as e:
        return {"success": False, "message": str(e), "details": results}

@app.get("/admin/stats")
def get_admin_stats():
    try:
        # Get Pinecone stats
        try:
            stats = index.describe_index_stats()
            total_vectors = stats.total_vector_count
        except Exception as e:
            print(f"Pinecone stats error: {e}")
            total_vectors = 0
        
        # Count local documents
        try:
            os.makedirs(DOCUMENTS_FOLDER, exist_ok=True)
            doc_count = len([f for f in os.listdir(DOCUMENTS_FOLDER) if f.endswith('.json')])
        except Exception as e:
            print(f"Document count error: {e}")
            doc_count = 0
        
        # Load BM25 data
        try:
            bm25_data = load_bm25_index()
        except Exception as e:
            print(f"BM25 load error: {e}")
            bm25_data = {"documents": [], "metadata": []}
        
        type_counts = {"document": 0, "website": 0, "youtube": 0, "image": 0, "video": 0}
        for meta in bm25_data.get('metadata', []):
            doc_type = meta.get('type', 'document')
            if doc_type in type_counts:
                type_counts[doc_type] += 1
        
        return {
            "total_vectors": total_vectors,
            "total_documents": doc_count,
            "bm25_chunks": len(bm25_data.get('documents', [])),
            "type_breakdown": type_counts,
            "index_name": os.getenv("PINECONE_INDEX_NAME")
        }
    except Exception as e:
        print(f"Admin stats error: {e}")
        return {
            "total_vectors": 0,
            "total_documents": 0,
            "bm25_chunks": 0,
            "type_breakdown": {"document": 0, "website": 0, "youtube": 0, "image": 0, "video": 0},
            "index_name": os.getenv("PINECONE_INDEX_NAME"),
            "error": str(e)
        }

# ==================== Streaming Chat Endpoint ====================

async def generate_stream(question: str, user_id: str, conversation_id: str):
    """Generate streaming response with proper async"""
    
    # Create or load conversation
    if not conversation_id:
        conversation = create_conversation(question)
        conversation_id = conversation['id']
    
    # Add user message
    add_message_to_conversation(conversation_id, "user", question)
    
    # Get conversation history (LIMITED to prevent token overflow)
    conv = load_conversation(conversation_id)
    conversation_history = ""
    MAX_HISTORY_TOKENS = 2000  # Limit history to ~2000 tokens
    
    if conv and len(conv['messages']) > 1:
        recent_messages = conv['messages'][-5:-1]  # Reduced from 7 to 5
        if recent_messages:
            temp_history = "Recent Conversation:\n"
            for msg in recent_messages:
                role = "User" if msg['role'] == 'user' else "Assistant"
                content = msg['content'][:800] + "..." if len(msg['content']) > 800 else msg['content']  # Reduced from 1500
                temp_history += f"{role}: {content}\n\n"
            
            # Check token limit
            history_tokens = len(temp_history) // 4
            if history_tokens > MAX_HISTORY_TOKENS:
                # Truncate to fit
                max_chars = MAX_HISTORY_TOKENS * 4
                conversation_history = temp_history[:max_chars] + "\n... [history truncated]\n"
            else:
                conversation_history = temp_history
    
    # Build search query
    search_query = question
    if conv and len(conv['messages']) > 2:
        user_messages = [m['content'] for m in conv['messages'] if m['role'] == 'user'][-3:]
        search_query = " ".join(user_messages)
    
    # Get document context WITH tracking enabled for retrieval
    doc_context, sources, images, tracker = get_document_context_with_sources(search_query, track_pipeline=True)
    
    # Send sources first (no need to send images separately now)
    sources_json = json.dumps({"type": "sources", "data": sources, "conversation_id": conversation_id})
    yield f"data: {sources_json}\n\n"
    
    # Add image info to prompt with full URLs for inline embedding
    image_info = ""
    if images:
        image_info = "\n\n=== SCREENSHOTS AVAILABLE (YOU MUST EMBED THESE) ===\n"
        for i, img in enumerate(images[:10], 1):
            url = img.get('url', '')
            image_info += f"SCREENSHOT_{i}: ![Screenshot {i}]({url})\n"
        image_info += "=== END SCREENSHOTS ===\n"
        image_info += "\nIMPORTANT: Copy and paste the SCREENSHOT lines above into your response where relevant!\n"
    
    # Build prompt
    prompt = f"""You are a friendly and helpful AI assistant with access to a knowledge base.

{conversation_history}

Available Information (Documents + Memory):
{doc_context if doc_context.strip() else "No specific information found."}
{image_info}

User Message: {question}

Instructions:
- Be friendly, helpful, and conversational
- If greeting, respond with a friendly greeting
- Use conversation history to understand context

**IMPORTANT - Be Smart About Matching:**
- Use the information above to answer, but be INTELLIGENT about matching
- Understand that users may use different names for the same product/service
- Example: "Database cloud service" likely refers to "Base Database Service"
- Example: "compute E4" means "Compute - Standard - E4"
- Use your knowledge to connect user's question to the available information
- If information seems related but named differently, use it and answer helpfully

**When to say "I don't have information":**
- ONLY say this if there's truly NO relevant information in the context
- Do NOT say this just because the exact wording doesn't match
- Try to find connections and answer helpfully first

- Format using markdown when helpful (headers, bullet points, bold)
- NEVER mention "documents", "context", "memory" in your response
- NEVER list sources in text - shown separately
- Be thorough - include ALL relevant details
- When asked to list ALL items, include EVERY item found

**SCREENSHOT EMBEDDING RULE**:
If SCREENSHOTS are listed above, you MUST embed them in your response.
Simply COPY the exact line like: ![Screenshot 1](/doc-images/xxx.jpeg)
Place screenshots after the relevant step they illustrate.

Response:"""
    
    # CRITICAL: Final token check before sending to OpenAI
    MAX_TOTAL_TOKENS = 25000  # Leave headroom below 30000 limit
    prompt_tokens = len(prompt) // 4
    
    if prompt_tokens > MAX_TOTAL_TOKENS:
        print(f"âš ï¸ Prompt too large ({prompt_tokens} tokens), truncating context...")
        # Calculate how much to reduce
        excess_tokens = prompt_tokens - MAX_TOTAL_TOKENS + 1000  # Extra buffer
        excess_chars = excess_tokens * 4
        
        # Truncate doc_context
        if len(doc_context) > excess_chars:
            doc_context = doc_context[:len(doc_context) - excess_chars] + "\n... [context truncated to fit token limit]"
            
            # Rebuild prompt with truncated context
            prompt = f"""You are a friendly and helpful AI assistant with access to a knowledge base.

{conversation_history}

Available Information (Documents + Memory):
{doc_context if doc_context.strip() else "No specific information found."}
{image_info}

User Message: {question}

Instructions:
- Be friendly, helpful, and conversational
- If greeting, respond with a friendly greeting
- Use conversation history to understand context

**IMPORTANT - Be Smart About Matching:**
- Use the information above to answer, but be INTELLIGENT about matching
- Understand that users may use different names for the same product/service
- If information seems related but named differently, use it and answer helpfully
- ONLY say "I don't have information" if there's truly NO relevant information

- Format using markdown when helpful (headers, bullet points, bold)
- NEVER mention "documents", "context", "memory" in your response
- NEVER list sources in text - shown separately
- Be thorough - include ALL relevant details

**SCREENSHOT EMBEDDING RULE**:
If SCREENSHOTS are listed above, you MUST embed them in your response.
Simply COPY the exact line like: ![Screenshot 1](/doc-images/xxx.jpeg)
Place screenshots after the relevant step they illustrate.

Response:"""
            print(f"Ã¢Å“â€¦ Prompt reduced to ~{len(prompt) // 4} tokens")
    
    # Start generate step tracking
    if tracker:
        tracker.start_step("generate")
    
    # Calculate prompt breakdown for token transparency
    system_prompt = """You are a friendly and helpful AI assistant with access to a knowledge base.
Instructions:
- Be friendly, helpful, and conversational
- Use conversation history to understand context
- Be INTELLIGENT about matching - users may use different names for same products
- Example: "Database cloud service" = "Base Database Service"
- If information seems related but named differently, use it and answer helpfully
- ONLY say "I don't have information" if there's truly NO relevant data
- Format using markdown when helpful
- NEVER mention "documents", "context", "memory" in your response
- Be thorough - include ALL relevant details"""
    
    prompt_breakdown = CostTracker.calculate_prompt_breakdown(
        system_prompt=system_prompt,
        conversation_history=conversation_history,
        document_context=doc_context,
        user_question=question
    )
    
    # Generate optimization suggestions
    context_tokens = prompt_breakdown["components"]["document_context"]["tokens"]
    total_input_tokens = prompt_breakdown["total_input_tokens"]
    suggestions = CostTracker.generate_optimization_suggestions(
        question=question,
        context_tokens=context_tokens,
        total_tokens=total_input_tokens,
        operation="retrieval"
    )
    
    # Stream the response
    full_response = ""
    
    try:
        # Use streaming LLM with async - Ã˜Â§Ã˜Â³Ã˜ÂªÃ˜Â®Ã˜Â¯Ã™â€¦ astream Ã˜Â¨Ã˜Â¯Ã™â€ž stream
        streaming_llm = ChatOpenAI(model="gpt-4o", temperature=0.7, streaming=True)
        
        async for chunk in streaming_llm.astream(prompt):
            if chunk.content:
                full_response += chunk.content
                chunk_json = json.dumps({"type": "chunk", "data": chunk.content})
                yield f"data: {chunk_json}\n\n"
        
        # Calculate and log LLM cost with breakdown
        llm_calc = CostTracker.calculate_llm_cost(prompt, full_response, "gpt-4o", breakdown=prompt_breakdown)
        CostTracker.log_cost(
            service="openai",
            model="gpt-4o",
            operation="generation",
            quantity=llm_calc["total_tokens"],
            unit="tokens",
            unit_price=0,
            total_cost=llm_calc["total_cost"],
            source=question[:50],
            pipeline_id=tracker.id if tracker else "",
            input_tokens=llm_calc["input_tokens"],
            output_tokens=llm_calc["output_tokens"],
            breakdown=prompt_breakdown
        )
        
        # Complete pipeline tracking
        if tracker:
            tracker.complete_step("generate", 
                model="gpt-4o",
                prompt_length=len(prompt),
                response_length=len(full_response),
                cost=llm_calc["total_cost"]
            )
            tracker.finish(
                sources_count=len(sources),
                has_images=len(images) > 0,
                total_cost=llm_calc["total_cost"]
            )
            if tracker.id in active_pipelines:
                del active_pipelines[tracker.id]
        
        # Send token info to frontend
        token_info = {
            "type": "token_info",
            "data": {
                "input_tokens": llm_calc["input_tokens"],
                "output_tokens": llm_calc["output_tokens"],
                "total_tokens": llm_calc["total_tokens"],
                "cost": llm_calc["total_cost"],
                "breakdown": {
                    "system_prompt": prompt_breakdown["components"]["system_prompt"]["tokens"],
                    "conversation_history": prompt_breakdown["components"]["conversation_history"]["tokens"],
                    "document_context": prompt_breakdown["components"]["document_context"]["tokens"],
                    "user_question": prompt_breakdown["components"]["user_question"]["tokens"]
                },
                "suggestions": suggestions,
                "formula": llm_calc["formula"]
            }
        }
        yield f"data: {json.dumps(token_info)}\n\n"
        
        # Send done signal
        yield f"data: {json.dumps({'type': 'done'})}\n\n"
        
        # Save to conversation and memory
        add_message_to_conversation(conversation_id, "assistant", full_response, sources)
        
        if doc_context.strip():
            save_to_memory(question, full_response, user_id)
            
    except Exception as e:
        if tracker:
            tracker.error_step("generate", str(e))
            tracker.finish(error=str(e))
            if tracker.id in active_pipelines:
                del active_pipelines[tracker.id]
        error_json = json.dumps({"type": "error", "data": str(e)})
        yield f"data: {error_json}\n\n"

@app.get("/chat/stream")
async def chat_stream(message: str, user_id: str = "default_user", conversation_id: str = None):
    """Streaming chat endpoint using Server-Sent Events"""
    return StreamingResponse(
        generate_stream(message, user_id, conversation_id),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no"
        }
    )

# ==================== Export Endpoints ====================

class ExportRequest(BaseModel):
    content: str
    title: str = "Export"
    images: List[dict] = []

@app.post("/export/word")
async def export_to_word(request: ExportRequest):
    """Export content with inline images to Word document"""
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    import io
    import re
    
    doc = Document()
    
    # Add title
    title = doc.add_heading(request.title, 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Parse markdown-like content and add to document
    content = request.content
    lines = content.split('\n')
    
    # Regex to find markdown images: ![alt](url)
    img_pattern = re.compile(r'!\[([^\]]*)\]\(([^)]+)\)')
    
    def get_image_from_url(img_url):
        """Get image bytes from URL or local path, convert to PNG for compatibility"""
        from PIL import Image as PILImage
        
        try:
            print(f"ðŸ–¼ï¸ Trying to load image: {img_url}")
            print(f"ðŸ“Â DOC_IMAGES_FOLDER: {DOC_IMAGES_FOLDER}")
            
            image_data = None
            
            # Check if it's a local doc-image
            if img_url.startswith('/doc-images/'):
                filename = img_url.replace('/doc-images/', '')
                local_path = os.path.join(DOC_IMAGES_FOLDER, filename)
                print(f"ðŸ“Â Local path: {local_path}")
                print(f"Ã¢Å“â€¦ File exists: {os.path.exists(local_path)}")
                
                if os.path.exists(local_path):
                    with open(local_path, 'rb') as f:
                        image_data = f.read()
                    print(f"ðŸ“¦ Loaded {len(image_data)} bytes")
                else:
                    print(f"Ã¢ÂÅ’ File not found at {local_path}")
                    return None
            else:
                # Try to fetch from URL
                print(f"ðŸŒÂ Trying HTTP fetch for: {img_url}")
                headers = {'User-Agent': 'Mozilla/5.0'}
                response = requests.get(img_url, headers=headers, timeout=10)
                if response.status_code == 200:
                    image_data = response.content
                else:
                    return None
            
            if image_data:
                # Convert to PNG using PIL for compatibility with python-docx
                print(f"ðŸ”â€ž Converting image to PNG...")
                pil_img = PILImage.open(io.BytesIO(image_data))
                
                # Convert to RGB if necessary
                if pil_img.mode in ('RGBA', 'LA', 'P'):
                    background = PILImage.new('RGB', pil_img.size, (255, 255, 255))
                    if pil_img.mode == 'P':
                        pil_img = pil_img.convert('RGBA')
                    if pil_img.mode == 'RGBA':
                        background.paste(pil_img, mask=pil_img.split()[-1])
                    else:
                        background.paste(pil_img)
                    pil_img = background
                elif pil_img.mode != 'RGB':
                    pil_img = pil_img.convert('RGB')
                
                # Save as PNG to BytesIO
                png_buffer = io.BytesIO()
                pil_img.save(png_buffer, format='PNG')
                png_buffer.seek(0)
                print(f"Ã¢Å“â€¦ Converted to PNG: {png_buffer.getbuffer().nbytes} bytes")
                return png_buffer
                
        except Exception as e:
            print(f"Ã¢ÂÅ’ Error loading/converting image {img_url}: {e}")
            import traceback
            traceback.print_exc()
        return None
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # Check for inline images in the line
        img_matches = list(img_pattern.finditer(line))
        
        if img_matches:
            # Process line with images
            last_end = 0
            for match in img_matches:
                # Add text before image
                text_before = line[last_end:match.start()].strip()
                if text_before:
                    p = doc.add_paragraph()
                    parts = re.split(r'\*\*(.*?)\*\*', text_before)
                    for i, part in enumerate(parts):
                        if i % 2 == 1:
                            p.add_run(part).bold = True
                        else:
                            p.add_run(part)
                
                # Add image
                alt_text = match.group(1)
                img_url = match.group(2)
                
                image_stream = get_image_from_url(img_url)
                if image_stream:
                    try:
                        print(f"ðŸ“Å½ Adding picture to Word doc: {img_url}")
                        doc.add_picture(image_stream, width=Inches(5))
                        print(f"Ã¢Å“â€¦ Picture added successfully!")
                        
                        # Add caption
                        if alt_text:
                            cap_para = doc.add_paragraph(alt_text[:200])
                            cap_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            cap_para.runs[0].font.size = Pt(10)
                            cap_para.runs[0].font.italic = True
                    except Exception as e:
                        print(f"Ã¢ÂÅ’ Error adding picture: {e}")
                        import traceback
                        traceback.print_exc()
                        doc.add_paragraph(f"[Image: {alt_text or 'Could not load'}]")
                else:
                    print(f"Ã¢ÂÅ’ No image stream for: {img_url}")
                    doc.add_paragraph(f"[Image: {alt_text or 'Could not load'}]")
                
                last_end = match.end()
            
            # Add remaining text after last image
            text_after = line[last_end:].strip()
            if text_after:
                p = doc.add_paragraph()
                parts = re.split(r'\*\*(.*?)\*\*', text_after)
                for i, part in enumerate(parts):
                    if i % 2 == 1:
                        p.add_run(part).bold = True
                    else:
                        p.add_run(part)
        else:
            # No images - handle normally
            # Handle headers
            if line.startswith('### '):
                doc.add_heading(line[4:], level=3)
            elif line.startswith('## '):
                doc.add_heading(line[3:], level=2)
            elif line.startswith('# '):
                doc.add_heading(line[2:], level=1)
            # Handle bullet points
            elif line.startswith('- ') or line.startswith('* '):
                p = doc.add_paragraph(style='List Bullet')
                text = line[2:]
                parts = re.split(r'\*\*(.*?)\*\*', text)
                for i, part in enumerate(parts):
                    if i % 2 == 1:
                        p.add_run(part).bold = True
                    else:
                        p.add_run(part)
            # Handle numbered lists
            elif re.match(r'^\d+\.', line):
                p = doc.add_paragraph(style='List Number')
                text = re.sub(r'^\d+\.\s*', '', line)
                parts = re.split(r'\*\*(.*?)\*\*', text)
                for i, part in enumerate(parts):
                    if i % 2 == 1:
                        p.add_run(part).bold = True
                    else:
                        p.add_run(part)
            # Regular paragraph
            else:
                p = doc.add_paragraph()
                parts = re.split(r'\*\*(.*?)\*\*', line)
                for i, part in enumerate(parts):
                    if i % 2 == 1:
                        p.add_run(part).bold = True
                    else:
                        p.add_run(part)
    
    # Save to bytes
    doc_bytes = io.BytesIO()
    doc.save(doc_bytes)
    doc_bytes.seek(0)
    
    # Return as downloadable file
    filename = f"{request.title.replace(' ', '_')}.docx"
    return StreamingResponse(
        doc_bytes,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

@app.post("/export/pdf")
async def export_to_pdf(request: ExportRequest):
    """Export content with inline images to PDF"""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, ListFlowable, ListItem
        from reportlab.lib.enums import TA_CENTER, TA_LEFT
        import io
        import re
        from PIL import Image as PILImage
        
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72)
        
        styles = getSampleStyleSheet()
        
        # Custom styles
        styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=styles['Title'],
            fontSize=24,
            spaceAfter=30,
            alignment=TA_CENTER
        ))
        styles.add(ParagraphStyle(
            name='CustomHeading',
            parent=styles['Heading1'],
            fontSize=16,
            spaceAfter=12,
            spaceBefore=20
        ))
        styles.add(ParagraphStyle(
            name='CustomBody',
            parent=styles['Normal'],
            fontSize=11,
            spaceAfter=8,
            leading=14
        ))
        styles.add(ParagraphStyle(
            name='Caption',
            parent=styles['Normal'],
            fontSize=9,
            alignment=TA_CENTER,
            italic=True,
            spaceAfter=15
        ))
        
        story = []
        
        # Add title
        story.append(Paragraph(request.title, styles['CustomTitle']))
        story.append(Spacer(1, 20))
        
        # Regex to find markdown images
        img_pattern = re.compile(r'!\[([^\]]*)\]\(([^)]+)\)')
        
        def load_and_convert_image(img_url):
            """Load image and convert to PNG for ReportLab"""
            try:
                img_data = None
                
                # Check if it's a local doc-image
                if img_url.startswith('/doc-images/'):
                    filename = img_url.replace('/doc-images/', '')
                    local_path = os.path.join(DOC_IMAGES_FOLDER, filename)
                    if os.path.exists(local_path):
                        with open(local_path, 'rb') as f:
                            img_data = f.read()
                else:
                    # Try to fetch from URL
                    headers = {'User-Agent': 'Mozilla/5.0'}
                    response = requests.get(img_url, headers=headers, timeout=15)
                    if response.status_code == 200:
                        img_data = response.content
                
                if img_data:
                    # Convert to PNG using PIL
                    pil_img = PILImage.open(io.BytesIO(img_data))
                    
                    # Convert to RGB if necessary
                    if pil_img.mode in ('RGBA', 'LA', 'P'):
                        background = PILImage.new('RGB', pil_img.size, (255, 255, 255))
                        if pil_img.mode == 'P':
                            pil_img = pil_img.convert('RGBA')
                        if pil_img.mode == 'RGBA':
                            background.paste(pil_img, mask=pil_img.split()[-1])
                        else:
                            background.paste(pil_img)
                        pil_img = background
                    elif pil_img.mode != 'RGB':
                        pil_img = pil_img.convert('RGB')
                    
                    img_buffer = io.BytesIO()
                    pil_img.save(img_buffer, format='PNG')
                    img_buffer.seek(0)
                    return img_buffer
                    
            except Exception as e:
                print(f"Error loading image {img_url}: {e}")
            return None
        
        # Process content line by line, handling images inline
        content = request.content
        lines = content.split('\n')
        current_list = []
        
        for line in lines:
            line = line.strip()
            if not line:
                if current_list:
                    story.append(ListFlowable(current_list, bulletType='bullet'))
                    current_list = []
                continue
            
            # Check for image in line
            img_match = img_pattern.search(line)
            if img_match:
                # Flush any pending list
                if current_list:
                    story.append(ListFlowable(current_list, bulletType='bullet'))
                    current_list = []
                
                # Get text before and after image
                text_before = line[:img_match.start()].strip()
                text_after = line[img_match.end():].strip()
                
                # Add text before image
                if text_before:
                    text_before = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', text_before)
                    story.append(Paragraph(text_before, styles['CustomBody']))
                
                # Add the image
                alt_text = img_match.group(1)
                img_url = img_match.group(2)
                
                img_buffer = load_and_convert_image(img_url)
                if img_buffer:
                    try:
                        img_obj = RLImage(img_buffer, width=5*inch, height=4*inch, kind='proportional')
                        story.append(Spacer(1, 10))
                        story.append(img_obj)
                        if alt_text:
                            story.append(Paragraph(alt_text, styles['Caption']))
                        story.append(Spacer(1, 10))
                    except Exception as e:
                        print(f"Error adding image to PDF: {e}")
                        story.append(Paragraph(f"[Image: {alt_text}]", styles['Caption']))
                else:
                    story.append(Paragraph(f"[Image: {alt_text}]", styles['Caption']))
                
                # Add text after image
                if text_after:
                    text_after = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', text_after)
                    story.append(Paragraph(text_after, styles['CustomBody']))
                
                continue
            
            # Convert markdown bold to HTML
            line = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', line)
            
            # Handle headers
            if line.startswith('### '):
                if current_list:
                    story.append(ListFlowable(current_list, bulletType='bullet'))
                    current_list = []
                story.append(Paragraph(line[4:], styles['Heading3']))
            elif line.startswith('## '):
                if current_list:
                    story.append(ListFlowable(current_list, bulletType='bullet'))
                    current_list = []
                story.append(Paragraph(line[3:], styles['Heading2']))
            elif line.startswith('# '):
                if current_list:
                    story.append(ListFlowable(current_list, bulletType='bullet'))
                    current_list = []
                story.append(Paragraph(line[2:], styles['CustomHeading']))
            elif line.startswith('- ') or line.startswith('* '):
                current_list.append(ListItem(Paragraph(line[2:], styles['CustomBody'])))
            elif re.match(r'^\d+\.', line):
                if current_list:
                    story.append(ListFlowable(current_list, bulletType='bullet'))
                    current_list = []
                text = re.sub(r'^\d+\.\s*', '', line)
                story.append(Paragraph(text, styles['CustomBody']))
            else:
                if current_list:
                    story.append(ListFlowable(current_list, bulletType='bullet'))
                    current_list = []
                story.append(Paragraph(line, styles['CustomBody']))
        
        # Flush any remaining list
        if current_list:
            story.append(ListFlowable(current_list, bulletType='bullet'))
        
        # Build PDF
        doc.build(story)
        buffer.seek(0)
        
        filename = f"{request.title.replace(' ', '_')}.pdf"
        return StreamingResponse(
            buffer,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    
    except Exception as e:
        print(f"PDF Export Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return {"error": str(e)}

# ==================== Debug Endpoint ====================

@app.get("/debug/images")
def debug_images(url: str):
    """Debug endpoint to test image extraction from a URL"""
    try:
        headers = {'User-Agent': 'Mozilla/5.0'}
        response = requests.get(url, headers=headers, timeout=15)
        response.raise_for_status()
        
        images = extract_images_from_page(url, response.text)
        
        return {
            "url": url,
            "images_found": len(images),
            "images": images
        }
    except Exception as e:
        return {"error": str(e)}

@app.get("/debug/search")
def debug_search(query: str):
    """Debug search endpoint to see what sources are found"""
    search_results = hybrid_search(query, top_k=10, source_filter=None)
    reranked = rerank_results(query, search_results, top_n=5)
    
    return {
        "query": query,
        "search_results_count": len(search_results),
        "search_results": [
            {
                "type": r.get('metadata', {}).get('type'),
                "title": r.get('metadata', {}).get('title'),
                "parent_id": r.get('metadata', {}).get('parent_id'),
                "score": r.get('score'),
                "is_memory": r.get('is_memory', False)
            }
            for r in search_results[:10]
        ],
        "reranked": [
            {
                "type": r.get('metadata', {}).get('type'),
                "title": r.get('metadata', {}).get('title'),
                "score": r.get('score')
            }
            for r in reranked
        ]
    }

@app.get("/debug/doc-images")
def debug_doc_images():
    """Debug endpoint to list all extracted document images"""
    images = []
    if os.path.exists(DOC_IMAGES_FOLDER):
        for filename in os.listdir(DOC_IMAGES_FOLDER):
            filepath = os.path.join(DOC_IMAGES_FOLDER, filename)
            images.append({
                "filename": filename,
                "url": f"/doc-images/{filename}",
                "size_kb": round(os.path.getsize(filepath) / 1024, 2),
                "full_path": filepath,
                "exists": os.path.exists(filepath)
            })
    
    return {
        "folder": DOC_IMAGES_FOLDER,
        "folder_exists": os.path.exists(DOC_IMAGES_FOLDER),
        "total_images": len(images),
        "images": images
    }

@app.get("/debug/test-image-load")
def debug_test_image_load(img_url: str):
    """Test if an image can be loaded for export"""
    import io
    
    result = {
        "url": img_url,
        "is_local": img_url.startswith('/doc-images/'),
        "doc_images_folder": DOC_IMAGES_FOLDER
    }
    
    try:
        if img_url.startswith('/doc-images/'):
            filename = img_url.replace('/doc-images/', '')
            local_path = os.path.join(DOC_IMAGES_FOLDER, filename)
            result["local_path"] = local_path
            result["file_exists"] = os.path.exists(local_path)
            
            if os.path.exists(local_path):
                with open(local_path, 'rb') as f:
                    data = f.read()
                result["file_size_kb"] = round(len(data) / 1024, 2)
                result["success"] = True
            else:
                result["success"] = False
                result["error"] = "File not found"
        else:
            result["success"] = False
            result["error"] = "Not a local doc-image URL"
    except Exception as e:
        result["success"] = False
        result["error"] = str(e)
    
    return result

@app.post("/debug/export-content")
async def debug_export_content(request: ExportRequest):
    """Debug endpoint to see what content is being sent to export"""
    import re
    
    content = request.content
    img_pattern = re.compile(r'!\[([^\]]*)\]\(([^)]+)\)')
    
    # Find all images in content
    images_found = []
    for match in img_pattern.finditer(content):
        images_found.append({
            "alt": match.group(1),
            "url": match.group(2)
        })
    
    return {
        "title": request.title,
        "content_length": len(content),
        "content_preview": content[:500],
        "content_full": content,
        "images_in_request": request.images,
        "images_found_in_content": images_found,
        "total_images_found": len(images_found)
    }

@app.get("/debug/document-images/{doc_title}")
def debug_document_images(doc_title: str):
    """Debug endpoint to see images for a specific document"""
    # Search in individual document files
    if os.path.exists(DOCUMENTS_FOLDER):
        for filename in os.listdir(DOCUMENTS_FOLDER):
            if filename.endswith('.json') and filename.startswith('parent_'):
                filepath = os.path.join(DOCUMENTS_FOLDER, filename)
                try:
                    with open(filepath, "r", encoding='utf-8') as f:
                        doc = json.load(f)
                    
                    title = doc.get('metadata', {}).get('title', '')
                    if doc_title.lower() in title.lower():
                        return {
                            "doc_id": doc.get('id'),
                            "title": title,
                            "images_count": len(doc.get('images', [])),
                            "images": doc.get('images', [])
                        }
                except Exception as e:
                    continue
    
    return {"error": f"Document '{doc_title}' not found"}

@app.get("/debug/all-documents")
def debug_all_documents():
    """Debug endpoint to list all documents with their image counts"""
    documents = []
    if os.path.exists(DOCUMENTS_FOLDER):
        for filename in os.listdir(DOCUMENTS_FOLDER):
            if filename.endswith('.json') and filename.startswith('parent_'):
                filepath = os.path.join(DOCUMENTS_FOLDER, filename)
                try:
                    with open(filepath, "r", encoding='utf-8') as f:
                        doc = json.load(f)
                    
                    documents.append({
                        "id": doc.get('id'),
                        "title": doc.get('metadata', {}).get('title', 'Unknown'),
                        "type": doc.get('metadata', {}).get('type', 'Unknown'),
                        "images_count": len(doc.get('images', []))
                    })
                except Exception as e:
                    continue
    
    return {
        "total_documents": len(documents),
        "documents": documents
    }

@app.get("/debug/prompt")
def debug_prompt(query: str):
    """Debug endpoint to see what images are being passed to the LLM"""
    doc_context, sources, images, _ = get_document_context_with_sources(query, track_pipeline=False)
    
    return {
        "query": query,
        "sources_count": len(sources),
        "images_count": len(images),
        "images": images,
        "context_preview": doc_context[:500] if doc_context else "No context"
    }



# ==================== Configuration Endpoints ====================

@app.get("/config/presets")
def get_config_presets():
    """Get available configuration presets for the frontend"""
    return {
        "presets": get_preset_info(),
        "defaults": AgentConfig().to_dict()
    }

@app.get("/config/current")
def get_current_config():
    """Get current default configuration"""
    return AgentConfig().to_dict()

@app.post("/config/validate")
def validate_config(settings: dict):
    """Validate user configuration settings"""
    try:
        config = AgentConfig.from_user_settings(settings)
        return {
            "valid": True,
            "config": config.to_dict()
        }
    except Exception as e:
        return {
            "valid": False,
            "error": str(e)
        }

# ==================== Run Server ====================

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
